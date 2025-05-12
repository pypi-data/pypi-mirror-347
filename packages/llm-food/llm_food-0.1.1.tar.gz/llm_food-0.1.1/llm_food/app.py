import asyncio
import base64

from datetime import datetime
import hashlib
import json
from io import BytesIO
import os
from typing import List, Optional
import uuid
from fastapi import (
    FastAPI,
    UploadFile,
    File,
    BackgroundTasks,
    HTTPException,
    Query,
    Depends,
    Form,
)
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials

import httpx
import mammoth
import duckdb

# Imports for GCS
from google.cloud import storage
from google.oauth2 import service_account  # For local testing with service account
from markdownify import markdownify

from striprtf.striprtf import rtf_to_text
from pptx import Presentation
import trafilatura

from .config import (
    get_pdf_backend,
    get_gemini_prompt,
    get_api_auth_token,
    get_gcs_project_id,
    get_max_file_size_bytes,
    DUCKDB_FILE,
    GCS_BATCH_BUCKET,
    GEMINI_MODEL_FOR_VISION,
    SUPPORTED_EXTENSIONS,
)

from .models import (
    BatchJobFailedFileOutput,
    BatchJobOutputResponse,
    BatchOutputItem,
    ConversionResponse,
    BatchJobStatusResponse,
)

# --- Conditional imports based on the PDF backend ---
match get_pdf_backend():
    case "pymupdf4llm":
        from pymupdf4llm import to_markdown
        import pymupdf
    case "pypdf2":
        from pypdf import PdfReader
    case "gemini":
        from pdf2image import convert_from_bytes
        from google import genai
        from google.genai.types import CreateBatchJobConfig, JobState

        OCR_PROMPT = get_gemini_prompt()
    case invalid_backend:
        raise ValueError(f"Invalid PDF backend: {invalid_backend}")

# --- Main application setup ---
app = FastAPI(
    title="LLM Food API",
    description="API for converting various document formats to Markdown or text, with batch processing capabilities.",
    version="0.2.0",
)

# --- Security ---
bearer_scheme = HTTPBearer(
    auto_error=False
)  # auto_error=False to handle optional token & custom errors


async def authenticate_request(
    credentials: Optional[HTTPAuthorizationCredentials] = Depends(bearer_scheme),
) -> None:
    configured_token = get_api_auth_token()
    if configured_token:  # Only enforce auth if a token is configured server-side
        if credentials is None:
            raise HTTPException(
                status_code=401,
                detail="Not authenticated. Authorization header is missing.",
                headers={"WWW-Authenticate": "Bearer"},
            )
        if credentials.scheme != "Bearer":
            raise HTTPException(
                status_code=401,
                detail="Invalid authentication scheme. Only Bearer is supported.",
                headers={"WWW-Authenticate": "Bearer"},
            )
        if credentials.token != configured_token:
            raise HTTPException(
                status_code=403,
                detail="Invalid token.",
            )
    # If no token is configured server-side, or if authentication passes, do nothing.
    return


# For local GCS testing with a service account JSON file
def get_gcs_credentials():
    credentials_path = os.getenv("GOOGLE_APPLICATION_CREDENTIALS")
    if credentials_path:
        return service_account.Credentials.from_service_account_file(credentials_path)
    return None  # Fallback to default environment auth if not set


def get_gemini_client():
    project = get_gcs_project_id()
    location = os.getenv("GOOGLE_CLOUD_LOCATION", "us-central1")
    client = genai.Client(vertexai=True, location=location, project=project)
    return client


TASKS = {}


def get_db_connection():
    return duckdb.connect(DUCKDB_FILE)


def initialize_db_schema():
    con = get_db_connection()
    try:
        # Main batch jobs table
        con.execute("""
            CREATE TABLE IF NOT EXISTS batch_jobs (
                job_id VARCHAR PRIMARY KEY,
                output_gcs_path VARCHAR NOT NULL,
                status VARCHAR NOT NULL,
                submitted_at TIMESTAMP NOT NULL,
                total_input_files INTEGER NOT NULL,
                overall_processed_count INTEGER DEFAULT 0,
                overall_failed_count INTEGER DEFAULT 0,
                last_updated_at TIMESTAMP
            )
        """)
        # Gemini PDF sub-jobs (one per Gemini Batch API call)
        con.execute("""
            CREATE TABLE IF NOT EXISTS gemini_pdf_batch_sub_jobs (
                gemini_batch_sub_job_id VARCHAR PRIMARY KEY,
                batch_job_id VARCHAR NOT NULL REFERENCES batch_jobs(job_id),
                gemini_api_job_name VARCHAR,
                status VARCHAR NOT NULL,
                payload_gcs_uri VARCHAR,
                gemini_output_gcs_uri_prefix VARCHAR,
                total_pdf_pages_for_batch INTEGER DEFAULT 0,
                processed_pdf_pages_count INTEGER DEFAULT 0,
                failed_pdf_pages_count INTEGER DEFAULT 0,
                error_message VARCHAR,
                created_at TIMESTAMP NOT NULL,
                updated_at TIMESTAMP
            )
        """)
        # Individual file tasks (for non-PDFs, or individual pages of PDFs before aggregation)
        con.execute("""
            CREATE TABLE IF NOT EXISTS file_tasks (
                file_task_id VARCHAR PRIMARY KEY,
                batch_job_id VARCHAR NOT NULL REFERENCES batch_jobs(job_id),
                gemini_batch_sub_job_id VARCHAR REFERENCES gemini_pdf_batch_sub_jobs(gemini_batch_sub_job_id), -- Link to a Gemini batch if it's a PDF page
                original_filename VARCHAR NOT NULL,
                file_type VARCHAR NOT NULL, -- e.g., 'pdf_page', 'docx'
                status VARCHAR NOT NULL, -- pending, processing, image_uploaded_to_gcs, completed, failed
                gcs_input_image_uri VARCHAR, -- For PDF pages, GCS URI of the image sent to Gemini
                gcs_output_markdown_uri VARCHAR, -- GCS URI of the final .md (for non-PDFs or aggregated PDFs)
                page_number INTEGER, -- For PDF pages
                gemini_request_id VARCHAR, -- The 'id' used in payload.jsonl for this PDF page
                error_message VARCHAR,
                created_at TIMESTAMP NOT NULL,
                updated_at TIMESTAMP
            )
        """)
    finally:
        con.close()


# Call initialization at startup
initialize_db_schema()


def _process_docx_sync(content_bytes: bytes) -> List[str]:
    try:
        doc = BytesIO(content_bytes)
        doc_html = mammoth.convert_to_html(doc).value
        doc_md = markdownify(doc_html).strip()
        return [doc_md]
    except Exception as e:
        return [f"Error processing DOCX: {str(e)}"]


def _process_rtf_sync(content_bytes: bytes) -> List[str]:
    try:
        return [rtf_to_text(content_bytes.decode("utf-8", errors="ignore"))]
    except Exception as e:
        return [f"Error processing RTF: {str(e)}"]


def _process_pptx_sync(content_bytes: bytes) -> List[str]:
    try:
        prs = Presentation(BytesIO(content_bytes))
        # Corrected list comprehension for PPTX to build a single string per slide, then list of slide texts
        slide_texts = []
        for slide in prs.slides:
            text_on_slide = "\n".join(
                shape.text
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text
            )
            if text_on_slide:  # Only add if there's text
                slide_texts.append(text_on_slide)
        return (
            slide_texts if slide_texts else [""]
        )  # Return list of slide texts, or list with empty string if no text
    except Exception as e:
        return [f"Error processing PPTX: {str(e)}"]


def _process_html_sync(content_bytes: bytes) -> List[str]:
    try:
        extracted_text = trafilatura.extract(
            content_bytes.decode("utf-8", errors="ignore"), output_format="markdown"
        )
        return [extracted_text if extracted_text is not None else ""]
    except Exception as e:
        return [f"Error processing HTML: {str(e)}"]


def _process_pdf_pymupdf4llm_sync(content_bytes: bytes) -> List[str]:
    try:
        pymupdf_doc = pymupdf.Document(stream=content_bytes, filetype="pdf")
        page_data_list = to_markdown(pymupdf_doc, page_chunks=True)
        return [page_dict.get("text", "") for page_dict in page_data_list]
    except Exception as e:
        return [f"Error processing PDF with pymupdf4llm: {str(e)}"]


def _process_pdf_pypdf2_sync(content_bytes: bytes) -> List[str]:
    try:
        reader = PdfReader(BytesIO(content_bytes))
        return [p.extract_text() or "" for p in reader.pages]
    except Exception as e:
        return [f"Error processing PDF with pypdf: {str(e)}"]


async def _process_file_content(
    ext: str, content: bytes, pdf_backend_choice: str
) -> List[str]:
    texts_list: List[str] = []
    if ext == ".pdf":
        if pdf_backend_choice == "pymupdf4llm":
            texts_list = await asyncio.to_thread(_process_pdf_pymupdf4llm_sync, content)
        elif pdf_backend_choice == "pypdf2":
            texts_list = await asyncio.to_thread(_process_pdf_pypdf2_sync, content)
        elif pdf_backend_choice == "gemini":
            pages = convert_from_bytes(content)
            images_b64 = []
            for page in pages:
                buffer = BytesIO()
                page.save(buffer, format="PNG")
                image_data = buffer.getvalue()
                b64_str = base64.b64encode(image_data).decode("utf-8")
                images_b64.append(b64_str)
            client = get_gemini_client()
            payloads = [
                [
                    {"inline_data": {"data": b64_str, "mime_type": "image/png"}},
                    {"text": OCR_PROMPT},
                ]
                for b64_str in images_b64
            ]
            results = await asyncio.gather(
                *[
                    client.aio.models.generate_content(
                        model=GEMINI_MODEL_FOR_VISION, contents=payload
                    )
                    for payload in payloads
                ]
            )
            texts_list = [result.text for result in results]
        else:
            texts_list = ["Invalid PDF backend specified."]
    elif ext in [".docx"]:
        texts_list = await asyncio.to_thread(_process_docx_sync, content)
    elif ext in [".rtf"]:
        texts_list = await asyncio.to_thread(_process_rtf_sync, content)
    elif ext in [".pptx"]:
        texts_list = await asyncio.to_thread(_process_pptx_sync, content)
    elif ext in [".html", ".htm"]:
        texts_list = await asyncio.to_thread(_process_html_sync, content)
    else:
        texts_list = ["Unsupported file type encountered in _process_file_content."]
    return texts_list


@app.post(
    "/convert",
    response_model=ConversionResponse,
    dependencies=[Depends(authenticate_request)],
)
async def convert_file_upload(file: UploadFile = File(...)):
    ext = os.path.splitext(file.filename)[1].lower()
    content = await file.read()

    max_size = get_max_file_size_bytes()
    if max_size is not None and len(content) > max_size:
        raise HTTPException(
            status_code=413,
            detail=f"File size {len(content) / (1024 * 1024):.2f}MB exceeds maximum allowed size of {max_size / (1024 * 1024):.2f}MB.",
        )

    content_hash = hashlib.sha256(content).hexdigest()
    pdf_backend_choice = get_pdf_backend()

    texts_list = await _process_file_content(ext, content, pdf_backend_choice)

    if texts_list and (
        texts_list[0].startswith("Error processing")
        or texts_list[0].startswith("Invalid PDF backend")
        or texts_list[0].startswith("Unsupported file type")
    ):
        raise HTTPException(status_code=400, detail=texts_list[0])

    return ConversionResponse(
        filename=file.filename, content_hash=content_hash, texts=texts_list
    )


@app.get(
    "/convert",
    response_model=ConversionResponse,
    dependencies=[Depends(authenticate_request)],
)
async def convert_url(
    url: str = Query(..., description="URL of the webpage to convert to Markdown"),
):
    try:
        async with httpx.AsyncClient() as client:
            response = await client.get(url)
            response.raise_for_status()
            html_content = response.text
            content_bytes = html_content.encode("utf-8")
    except httpx.RequestError as e:
        raise HTTPException(status_code=400, detail=f"Error fetching URL: {str(e)}")
    except httpx.HTTPStatusError as e:
        raise HTTPException(
            status_code=e.response.status_code,
            detail=f"Error fetching URL: {e.response.reason_phrase}",
        )

    if not html_content:
        raise HTTPException(status_code=400, detail="Fetched content is empty.")

    content_hash = hashlib.sha256(content_bytes).hexdigest()

    extracted_text = trafilatura.extract(html_content, output_format="markdown")
    texts_list = [extracted_text if extracted_text is not None else ""]

    filename = os.path.basename(url) or url

    return ConversionResponse(
        filename=filename, content_hash=content_hash, texts=texts_list
    )


@app.get(
    "/status/{task_id}",
    response_model=BatchJobStatusResponse,
    dependencies=[Depends(authenticate_request)],
)
def status(task_id: str):
    con = get_db_connection()
    try:
        job_status_row = con.execute(
            "SELECT * FROM batch_jobs WHERE job_id = ?", (task_id,)
        ).fetchone()

        if not job_status_row:
            raise HTTPException(
                status_code=404, detail=f"Batch job with ID {task_id} not found."
            )

        job_dict = dict(zip([desc[0] for desc in con.description], job_status_row))

        gemini_sub_jobs_rows = con.execute(
            "SELECT * FROM gemini_pdf_batch_sub_jobs WHERE batch_job_id = ?",
            (task_id,),
        ).fetchall()
        job_dict["gemini_pdf_processing_details"] = [
            dict(zip([desc[0] for desc in con.description], sub_job_row))
            for sub_job_row in gemini_sub_jobs_rows
        ]

        file_tasks_rows = con.execute(
            "SELECT original_filename, file_type, status, gcs_output_markdown_uri, error_message, page_number FROM file_tasks WHERE batch_job_id = ?",
            (task_id,),
        ).fetchall()
        job_dict["file_processing_details"] = [
            dict(zip([desc[0] for desc in con.description], task_row))
            for task_row in file_tasks_rows
        ]

        # Pydantic will validate the structure of job_dict against BatchJobStatusResponse
        return job_dict
    finally:
        con.close()


@app.post("/batch", dependencies=[Depends(authenticate_request)])
async def batch_files_upload(
    background_tasks: BackgroundTasks,
    files: List[UploadFile] = File(...),
    output_gcs_path: str = Form(...),
):
    main_batch_job_id = str(uuid.uuid4())
    current_time = datetime.utcnow()

    pdf_files_data_for_batch: List[
        tuple[str, bytes]
    ] = []  # Changed: Store (filename, content_bytes)
    non_pdf_files_for_individual_processing: List[
        UploadFile
    ] = []  # Remains UploadFile for now as content is read just before adding task

    if not GCS_BATCH_BUCKET:
        raise HTTPException(
            status_code=500,
            detail="GCS_BATCH_TEMP_BUCKET is not configured on the server.",
        )
    if not output_gcs_path.startswith("gs://"):
        raise HTTPException(
            status_code=400, detail="Output GCS path must start with gs://"
        )

    print("reading files...")
    for f in files:
        ext = os.path.splitext(f.filename)[1].lower()
        if ext == ".pdf":
            # Read PDF content here and store bytes
            content_bytes = await f.read()
            pdf_files_data_for_batch.append((f.filename, content_bytes))
        elif ext in SUPPORTED_EXTENSIONS:  # Excludes .pdf as it's handled above
            non_pdf_files_for_individual_processing.append(f)
        else:
            # Optionally log or report unsupported files
            print(f"Skipping unsupported file: {f.filename}")

    con = get_db_connection()
    try:
        con.execute(
            "INSERT INTO batch_jobs (job_id, output_gcs_path, status, submitted_at, total_input_files, last_updated_at) VALUES (?, ?, ?, ?, ?, ?)",
            (
                main_batch_job_id,
                output_gcs_path,
                "pending",
                current_time,
                len(pdf_files_data_for_batch)
                + len(non_pdf_files_for_individual_processing),
                current_time,
            ),
        )
        con.commit()

        # Process non-PDF files
        for upload_file in non_pdf_files_for_individual_processing:
            file_task_id = str(uuid.uuid4())
            file_ext = os.path.splitext(upload_file.filename)[1].lower()
            con.execute(
                "INSERT INTO file_tasks (file_task_id, batch_job_id, original_filename, file_type, status, created_at, updated_at) VALUES (?, ?, ?, ?, ?, ?, ?)",
                (
                    file_task_id,
                    main_batch_job_id,
                    upload_file.filename,
                    file_ext,
                    "pending",
                    current_time,
                    current_time,
                ),
            )
            # Read content before passing to background task
            content_bytes = await upload_file.read()
            await upload_file.seek(
                0
            )  # Reset cursor for safety, though might not be strictly needed if not read again

            background_tasks.add_task(
                _process_single_non_pdf_file_and_upload,
                content_bytes,  # Pass content bytes
                file_ext,
                upload_file.filename,
                output_gcs_path,
                main_batch_job_id,
                file_task_id,
            )
        con.commit()

        # Process PDF files via Gemini Batch
        if pdf_files_data_for_batch:
            gemini_batch_sub_job_id = str(uuid.uuid4())
            con.execute(
                "INSERT INTO gemini_pdf_batch_sub_jobs (gemini_batch_sub_job_id, batch_job_id, status, created_at, updated_at) VALUES (?, ?, ?, ?, ?)",
                (
                    gemini_batch_sub_job_id,
                    main_batch_job_id,
                    "pending_preparation",
                    current_time,
                    current_time,
                ),
            )
            # Pass the list of (filename, content_bytes)
            background_tasks.add_task(
                _run_gemini_pdf_batch_conversion,
                pdf_files_data_for_batch,  # Pass the list of (filename, content_bytes)
                output_gcs_path,
                main_batch_job_id,
                gemini_batch_sub_job_id,
            )
            print("Added background tasks for batch prediction PDF")
            con.commit()  # Commit after starting PDF batch task prep

        # Update batch job status to processing if there are tasks
        if pdf_files_data_for_batch or non_pdf_files_for_individual_processing:
            con.execute(
                "UPDATE batch_jobs SET status = ?, last_updated_at = ? WHERE job_id = ?",
                ("processing", datetime.utcnow(), main_batch_job_id),
            )
        else:  # No files to process
            con.execute(
                "UPDATE batch_jobs SET status = ?, last_updated_at = ? WHERE job_id = ?",
                ("completed_no_files", datetime.utcnow(), main_batch_job_id),
            )
        con.commit()

    finally:
        await _check_and_finalize_batch_job_status(
            main_batch_job_id, con
        )  # Call before closing
        con.commit()  # Ensure final status update is committed
        con.close()
    await asyncio.sleep(0.01)  # tick so the event loop gets forward
    return {"task_id": main_batch_job_id}


@app.get(
    "/batch/{task_id}",
    response_model=BatchJobOutputResponse,
    dependencies=[Depends(authenticate_request)],
)
async def get_batch_output(task_id: str):
    con = get_db_connection()
    try:
        job_details_tuple = con.execute(
            "SELECT job_id, status, output_gcs_path FROM batch_jobs WHERE job_id = ?",
            (task_id,),
        ).fetchone()
        if not job_details_tuple:
            raise HTTPException(
                status_code=404, detail=f"Batch job {task_id} not found."
            )

        job_id, job_status, output_gcs_path = job_details_tuple

        outputs_list = []
        errors_list = []
        message = None

        completed_statuses_for_output = ["completed", "completed_with_errors"]

        if job_status in completed_statuses_for_output:
            # Fetch distinct successful outputs
            # A file is considered successfully processed if at least one of its file_tasks (e.g. a page for a PDF, or the file itself for docx)
            # has a gcs_output_markdown_uri and status completed.
            # We need the original_filename and the gcs_output_markdown_uri of the *final aggregated file*.
            # The current DB schema stores the aggregated URI in each page's file_task if that page was part of a successful aggregate.
            successful_files_query = """
                SELECT DISTINCT original_filename, gcs_output_markdown_uri 
                FROM file_tasks 
                WHERE batch_job_id = ? AND status = 'completed' AND gcs_output_markdown_uri IS NOT NULL
            """
            successful_file_uris_tuples = con.execute(
                successful_files_query, (task_id,)
            ).fetchall()

            if successful_file_uris_tuples:
                storage_client = storage.Client(
                    project=get_gcs_project_id(), credentials=get_gcs_credentials()
                )
                for original_fn, gcs_uri in successful_file_uris_tuples:
                    try:
                        bucket_name, blob_name = gcs_uri.replace("gs://", "").split(
                            "/", 1
                        )
                        bucket = storage_client.bucket(bucket_name)
                        blob = bucket.blob(blob_name)
                        markdown_content = await asyncio.to_thread(
                            blob.download_as_text
                        )
                        outputs_list.append(
                            BatchOutputItem(
                                original_filename=original_fn,
                                markdown_content=markdown_content,
                                gcs_output_uri=gcs_uri,
                            )
                        )
                    except Exception as e:
                        print(f"Error downloading GCS content for {gcs_uri}: {e}")
                        # If we can't download a supposedly successful file, list it as an error for this retrieval attempt
                        errors_list.append(
                            BatchJobFailedFileOutput(
                                original_filename=original_fn,
                                file_type="unknown_at_retrieval",  # We don't easily have file_type here from this query
                                error_message=f"Failed to download content from GCS: {str(e)}",
                                status="retrieval_error",
                            )
                        )

            if job_status == "completed_with_errors":
                message = "Job completed with some errors. See errors list."
                failed_tasks_query = """
                    SELECT original_filename, file_type, page_number, error_message, status 
                    FROM file_tasks 
                    WHERE batch_job_id = ? AND status = 'failed'
                """
                failed_tasks_tuples = con.execute(
                    failed_tasks_query, (task_id,)
                ).fetchall()
                for (
                    ft_orig_fn,
                    ft_type,
                    ft_page_num,
                    ft_err_msg,
                    ft_status,
                ) in failed_tasks_tuples:
                    errors_list.append(
                        BatchJobFailedFileOutput(
                            original_filename=ft_orig_fn,
                            file_type=ft_type,
                            page_number=ft_page_num,
                            error_message=ft_err_msg,
                            status=ft_status,
                        )
                    )
            elif (
                not outputs_list and not errors_list
            ):  # Status was 'completed' but no files/uris found or downloaded
                message = "Job reported as completed, but no output files were found or could be retrieved."

        elif job_status == "completed_no_files":
            message = "Job completed, but no files were processed (e.g., no supported files in input)."
        else:  # Pending, processing, or failed states where individual outputs are not expected
            message = f"Job is not yet fully completed or has failed. Current status: {job_status}"

        return BatchJobOutputResponse(
            job_id=job_id,
            status=job_status,
            outputs=outputs_list,
            errors=errors_list,
            message=message,
        )
    finally:
        con.close()


async def _process_single_non_pdf_file_and_upload(
    content_bytes: bytes,
    file_ext: str,
    original_filename: str,
    output_gcs_path_str: str,
    main_batch_job_id: str,
    file_task_id: str,
):
    current_time = datetime.utcnow()
    con = get_db_connection()
    try:
        con.execute(
            "UPDATE file_tasks SET status = ?, updated_at = ? WHERE file_task_id = ?",
            ("processing", current_time, file_task_id),
        )
        con.execute(
            "UPDATE batch_jobs SET last_updated_at = ? WHERE job_id = ?",
            (current_time, main_batch_job_id),
        )
        con.commit()

        # Re-use the existing _process_file_content logic
        # Ensure pdf_backend_choice is not relevant or handled if _process_file_content expects it
        # For non-PDFs, pdf_backend_choice is not used by _process_file_content.
        markdown_texts = await _process_file_content(
            file_ext, content_bytes, get_pdf_backend()
        )  # get_pdf_backend() is for PDF choice, not directly used for non-PDFs here

        if (
            not markdown_texts
            or markdown_texts[0].startswith("Error processing")
            or markdown_texts[0].startswith("Unsupported file type")
        ):
            error_message = (
                markdown_texts[0] if markdown_texts else "Unknown processing error"
            )
            con.execute(
                "UPDATE file_tasks SET status = ?, error_message = ?, updated_at = ? WHERE file_task_id = ?",
                ("failed", error_message, datetime.utcnow(), file_task_id),
            )
            con.execute(
                "UPDATE batch_jobs SET overall_failed_count = overall_failed_count + 1, last_updated_at = ? WHERE job_id = ?",
                (datetime.utcnow(), main_batch_job_id),
            )
            con.commit()
            print(
                f"Failed to process non-PDF file {original_filename}: {error_message}"
            )
            return

        full_markdown_output = "\n\n---\n\n".join(markdown_texts)

        # Upload to GCS
        storage_client = storage.Client(
            project=get_gcs_project_id(), credentials=get_gcs_credentials()
        )
        output_bucket_name, output_prefix = output_gcs_path_str.replace(
            "gs://", ""
        ).split("/", 1)
        output_bucket = storage_client.bucket(output_bucket_name)

        output_blob_name = (
            output_prefix.rstrip("/")
            + "/"
            + os.path.splitext(original_filename)[0]
            + ".md"
        )
        output_blob_obj = output_bucket.blob(output_blob_name)

        output_blob_obj.upload_from_string(
            full_markdown_output, content_type="text/markdown"
        )
        gcs_output_url = f"gs://{output_bucket_name}/{output_blob_name}"

        con.execute(
            "UPDATE file_tasks SET status = ?, gcs_output_markdown_uri = ?, updated_at = ? WHERE file_task_id = ?",
            ("completed", gcs_output_url, datetime.utcnow(), file_task_id),
        )
        con.execute(
            "UPDATE batch_jobs SET overall_processed_count = overall_processed_count + 1, last_updated_at = ? WHERE job_id = ?",
            (datetime.utcnow(), main_batch_job_id),
        )
        con.commit()
        print(
            f"Successfully processed and uploaded non-PDF file {original_filename} to {gcs_output_url}"
        )

    except Exception as e:
        error_str = f"Error in _process_single_non_pdf_file_and_upload for {original_filename}: {str(e)}"
        print(error_str)
        try:
            # Attempt to mark as failed in DB
            con.execute(
                "UPDATE file_tasks SET status = ?, error_message = ?, updated_at = ? WHERE file_task_id = ?",
                ("failed", error_str, datetime.utcnow(), file_task_id),
            )
            con.execute(
                "UPDATE batch_jobs SET overall_failed_count = overall_failed_count + 1, last_updated_at = ? WHERE job_id = ?",
                (datetime.utcnow(), main_batch_job_id),
            )
            con.commit()
        except Exception as db_err:
            print(
                f"Additionally, failed to update DB for task {file_task_id} failure: {db_err}"
            )
    finally:
        await _check_and_finalize_batch_job_status(
            main_batch_job_id, con
        )  # Call before closing
        con.commit()  # Ensure final status update is committed
        con.close()


async def _run_gemini_pdf_batch_conversion(
    pdf_inputs_list: List[
        tuple[str, bytes]
    ],  # Expect list of (filename, content_bytes)
    output_gcs_path_str: str,
    main_batch_job_id: str,
    gemini_batch_sub_job_id: str,
):
    await asyncio.sleep(0.1)  # tick event loop

    current_time = datetime.utcnow()
    con = get_db_connection()
    storage_client = storage.Client(
        project=get_gcs_project_id(), credentials=get_gcs_credentials()
    )
    gemini_client = get_gemini_client()  # Ensure this client is suitable for batch

    temp_gcs_input_prefix = (
        f"gemini_batch_jobs/{main_batch_job_id}/{gemini_batch_sub_job_id}/inputs"
    )
    temp_gcs_images_prefix = f"{temp_gcs_input_prefix}/images"
    temp_gcs_output_prefix = (
        f"gemini_batch_jobs/{main_batch_job_id}/{gemini_batch_sub_job_id}/outputs"
    )

    payload_items_for_jsonl = []
    total_pages_for_this_gemini_batch = 0

    try:
        con.execute(
            "UPDATE gemini_pdf_batch_sub_jobs SET status = ?, updated_at = ? WHERE gemini_batch_sub_job_id = ?",
            ("preparing_images_and_payload", current_time, gemini_batch_sub_job_id),
        )
        con.execute(
            "UPDATE batch_jobs SET last_updated_at = ? WHERE job_id = ?",
            (current_time, main_batch_job_id),
        )
        con.commit()

        temp_bucket = storage_client.bucket(GCS_BATCH_BUCKET)

        for (
            original_pdf_filename,
            pdf_content_bytes,
        ) in pdf_inputs_list:  # Changed: Iterate through (filename, content_bytes)
            # original_pdf_filename = pdf_upload_file.filename # Removed
            # pdf_content_bytes = await pdf_upload_file.read() # Removed
            # await pdf_upload_file.close() # Removed

            try:
                page_images = convert_from_bytes(pdf_content_bytes, fmt="png")
            except Exception as e:
                print(f"Failed to convert PDF {original_pdf_filename} to images: {e}")
                # Mark all potential pages for this PDF as failed in file_tasks if desired, or just log
                # For now, we skip this PDF and it won't contribute to total_pages_for_this_gemini_batch
                # and won't have file_tasks entries created here.
                # A more robust approach might create failed file_tasks entries for it.
                con.execute(
                    "UPDATE gemini_pdf_batch_sub_jobs SET error_message = COALESCE(error_message || CHR(10), '') || ? WHERE gemini_batch_sub_job_id = ?",
                    (
                        f"Failed to convert PDF {original_pdf_filename} to images: {str(e)}",
                        gemini_batch_sub_job_id,
                    ),
                )
                con.commit()
                continue  # Skip to the next PDF file

            total_pages_for_this_gemini_batch += len(page_images)

            for i, page_image in enumerate(page_images):
                page_num = i + 1
                gemini_request_id = f"{os.path.splitext(original_pdf_filename)[0]}_p{page_num}_{uuid.uuid4().hex[:8]}"
                file_task_id = str(uuid.uuid4())

                con.execute(
                    "INSERT INTO file_tasks (file_task_id, batch_job_id, gemini_batch_sub_job_id, original_filename, file_type, status, page_number, gemini_request_id, created_at, updated_at) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)",
                    (
                        file_task_id,
                        main_batch_job_id,
                        gemini_batch_sub_job_id,
                        original_pdf_filename,
                        "pdf_page",
                        "pending_image_upload",
                        page_num,
                        gemini_request_id,
                        datetime.utcnow(),
                        datetime.utcnow(),
                    ),
                )
                con.commit()

                try:
                    img_byte_arr = BytesIO()
                    page_image.save(img_byte_arr, format="PNG")
                    img_byte_arr = img_byte_arr.getvalue()

                    image_blob_name = (
                        f"{temp_gcs_images_prefix}/{gemini_request_id}.png"
                    )
                    image_blob = temp_bucket.blob(image_blob_name)
                    image_blob.upload_from_string(
                        img_byte_arr, content_type="image/png"
                    )
                    image_gcs_uri = f"gs://{GCS_BATCH_BUCKET}/{image_blob_name}"

                    con.execute(
                        "UPDATE file_tasks SET gcs_input_image_uri = ?, status = ?, updated_at = ? WHERE file_task_id = ?",
                        (
                            image_gcs_uri,
                            "image_uploaded_to_gcs",
                            datetime.utcnow(),
                            file_task_id,
                        ),
                    )
                    con.commit()
                    payload_items_for_jsonl.append(
                        {
                            "id": gemini_request_id,
                            "request": {
                                "contents": [
                                    {
                                        "role": "user",
                                        "parts": [
                                            {
                                                "file_data": {
                                                    "file_uri": image_gcs_uri,
                                                    "mime_type": "image/png",
                                                }
                                            },
                                            {"text": OCR_PROMPT},
                                        ],
                                    }
                                ]
                            },
                        }
                    )
                except Exception as e:
                    print(
                        f"Failed to upload image for {original_pdf_filename} page {page_num}: {e}"
                    )
                    con.execute(
                        "UPDATE file_tasks SET status = ?, error_message = ?, updated_at = ? WHERE file_task_id = ?",
                        (
                            "failed",
                            f"Image upload error: {str(e)}",
                            datetime.utcnow(),
                            file_task_id,
                        ),
                    )
                    # also update gemini_pdf_batch_sub_jobs about this failure
                    con.execute(
                        "UPDATE gemini_pdf_batch_sub_jobs SET failed_pdf_pages_count = failed_pdf_pages_count + 1, error_message = COALESCE(error_message || CHR(10), '') || ? WHERE gemini_batch_sub_job_id = ?",
                        (
                            f"Failed image upload for {original_pdf_filename} page {page_num}: {str(e)}",
                            gemini_batch_sub_job_id,
                        ),
                    )
                    con.commit()
                    # Continue to next page, this page will be marked as failed.

        con.execute(
            "UPDATE gemini_pdf_batch_sub_jobs SET total_pdf_pages_for_batch = ? WHERE gemini_batch_sub_job_id = ?",
            (total_pages_for_this_gemini_batch, gemini_batch_sub_job_id),
        )
        con.commit()

        if not payload_items_for_jsonl:
            print(
                f"No pages successfully prepared for Gemini batch sub job {gemini_batch_sub_job_id}. Aborting Gemini submission."
            )
            con.execute(
                "UPDATE gemini_pdf_batch_sub_jobs SET status = ?, error_message = COALESCE(error_message || CHR(10), '') || 'No pages to process', updated_at = ? WHERE gemini_batch_sub_job_id = ?",
                ("failed_no_payload", datetime.utcnow(), gemini_batch_sub_job_id),
            )
            # Update main batch job if all PDFs failed here and there were no non-PDFs
            # This logic might need to be more sophisticated at the batch_jobs level based on counts
            con.execute(
                "UPDATE batch_jobs SET overall_failed_count = overall_failed_count + ? WHERE job_id = ?",
                (len(pdf_inputs_list), main_batch_job_id),
            )
            con.commit()
            return

        # Create and upload payload.jsonl
        payload_jsonl_content = "\n".join(
            [json.dumps(item) for item in payload_items_for_jsonl]
        )
        payload_blob_name = f"{temp_gcs_input_prefix}/payload.jsonl"
        payload_blob = temp_bucket.blob(payload_blob_name)
        payload_blob.upload_from_string(
            payload_jsonl_content, content_type="application/jsonl"
        )
        payload_gcs_uri = f"gs://{GCS_BATCH_BUCKET}/{payload_blob_name}"

        con.execute(
            "UPDATE gemini_pdf_batch_sub_jobs SET payload_gcs_uri = ?, status = ?, updated_at = ? WHERE gemini_batch_sub_job_id = ?",
            (
                payload_gcs_uri,
                "submitting_to_gemini",
                datetime.utcnow(),
                gemini_batch_sub_job_id,
            ),
        )
        con.commit()

        # Submit to Gemini Batch API
        gemini_output_uri_for_job = f"gs://{GCS_BATCH_BUCKET}/{temp_gcs_output_prefix}"
        batch_job_config = CreateBatchJobConfig(dest=gemini_output_uri_for_job)

        print(
            f"Submitting batch job to Gemini for {gemini_batch_sub_job_id} with model {GEMINI_MODEL_FOR_VISION}"
        )
        gemini_job = gemini_client.batches.create(
            model=GEMINI_MODEL_FOR_VISION,
            src=payload_gcs_uri,
            config=batch_job_config,
        )
        # The create call above blocks until job completion or failure for google-genai,
        # or at least starts it and returns an object that can be polled. Let's assume it returns quickly with a job object.
        # The example in the article suggests it returns a job object immediately, and then polling is done.
        # client.batches.create() is a synchronous call that waits for the batch job to complete.
        # This might make the background task run for a long time.

        con.execute(
            "UPDATE gemini_pdf_batch_sub_jobs SET gemini_api_job_name = ?, status = ?, gemini_output_gcs_uri_prefix = ?, updated_at = ? WHERE gemini_batch_sub_job_id = ?",
            (
                gemini_job.name,
                str(gemini_job.state),
                gemini_output_uri_for_job,
                datetime.utcnow(),
                gemini_batch_sub_job_id,
            ),
        )
        con.commit()
        print(f"Gemini job {gemini_job.name} submitted. State: {gemini_job.state}")

        # Polling loop (client.batches.create might already do this, check SDK behavior)
        # If client.batches.create() is indeed blocking until completion, this loop is redundant for it,
        # but useful if it returns immediately.
        # For google-genai, `create` is blocking and polls until completion. `get` is then used to refresh state if needed.
        # Let's assume `gemini_job` is the final state object after `create()` completes.

        # Corrected Polling Logic based on user feedback
        polling_interval_seconds = 30
        completed_job_states = {
            JobState.JOB_STATE_SUCCEEDED,
            JobState.JOB_STATE_FAILED,
            JobState.JOB_STATE_CANCELLED,
            JobState.JOB_STATE_PAUSED,  # Considering PAUSED as a state to stop active polling, might need review
        }

        while gemini_job.state not in completed_job_states:
            await asyncio.sleep(polling_interval_seconds)
            try:
                refreshed_job = gemini_client.batches.get(name=gemini_job.name)
                if refreshed_job.state != gemini_job.state:
                    gemini_job = refreshed_job  # Update job object only if state changed or to get latest info
                    print(
                        f"Gemini job {gemini_job.name} polling. Current state: {gemini_job.state}"
                    )
                    con.execute(
                        "UPDATE gemini_pdf_batch_sub_jobs SET status = ?, updated_at = ? WHERE gemini_batch_sub_job_id = ?",
                        (
                            str(gemini_job.state),
                            datetime.utcnow(),
                            gemini_batch_sub_job_id,
                        ),
                    )
                    con.commit()
                else:
                    # Optional: print a less verbose polling message if state hasn't changed
                    # print(f"Gemini job {gemini_job.name} still in state: {gemini_job.state}")
                    pass  # State unchanged, continue polling
            except Exception as e:
                print(
                    f"Error during Gemini job polling for {gemini_job.name}: {e}. Will retry polling."
                )
                # Decide if certain errors should break the loop or be retried.
                # For now, simple retry after sleep.

        # Final state update after loop
        con.execute(
            "UPDATE gemini_pdf_batch_sub_jobs SET status = ?, updated_at = ? WHERE gemini_batch_sub_job_id = ?",
            (str(gemini_job.state), datetime.utcnow(), gemini_batch_sub_job_id),
        )
        con.commit()
        print(
            f"Gemini job {gemini_job.name} polling finished. Final state: {gemini_job.state}"
        )

        if gemini_job.state == JobState.JOB_STATE_SUCCEEDED:
            con.execute(
                "UPDATE gemini_pdf_batch_sub_jobs SET status = ?, updated_at = ? WHERE gemini_batch_sub_job_id = ?",
                (
                    "processing_gemini_results",
                    datetime.utcnow(),
                    gemini_batch_sub_job_id,
                ),
            )
            con.commit()

            # Find the predictions.jsonl file. It's in a timestamped or job-id named subfolder.
            # The article uses a utility: get_latest_folder. Simpler: list blobs with prefix.
            # Output is typically: gs://{dest_bucket}/{dest_prefix}/gemini_batch_JOB_ID_TIMESTAMP/predictions.jsonl
            # The gemini_job.output_files might contain the direct path.
            # Let's assume the Batch API produces output like: {gemini_output_uri_for_job}/{gemini_job.name}/predictions.jsonl (this structure might vary)
            # Or more likely: {gemini_output_uri_for_job}/predictions_JOB_ID_PART_OF_URI.jsonl or similar fixed name patterns.

            # The result files are directly in the output_uri specified. e.g. {gemini_output_uri_for_job}/predictions.jsonl-00000-of-00001
            # Or, it might be a folder. The SDK or docs should clarify this structure precisely.
            # The article implies: `output_uri/TIMESTAMPED_FOLDER/predictions.jsonl`
            # Let's try to list blobs to find it.
            output_blobs = list(
                storage_client.list_blobs(
                    GCS_BATCH_BUCKET, prefix=f"{temp_gcs_output_prefix}/"
                )
            )
            predictions_file_blob = None
            for blob_item in output_blobs:
                if blob_item.name.endswith("predictions.jsonl") or (
                    "predictions.jsonl-" in blob_item.name
                ):  # Batch output can be sharded
                    predictions_file_blob = blob_item
                    break  # Take the first one found, or handle multiple shards if necessary

            if not predictions_file_blob:
                raise Exception(
                    f"predictions.jsonl not found in Gemini output: gs://{GCS_BATCH_BUCKET}/{temp_gcs_output_prefix}/"
                )

            print(f"Downloading predictions from: {predictions_file_blob.name}")
            predictions_content = predictions_file_blob.download_as_text()

            page_markdown_results = {}
            processed_gemini_pages = 0
            failed_gemini_pages = 0

            for line in predictions_content.splitlines():
                if not line.strip():
                    continue
                try:
                    prediction = json.loads(line)
                    request_id = prediction["id"]  # This is our gemini_request_id
                    # Assuming direct text output as we didn't ask for structured JSON response for the text itself
                    markdown_text = prediction["response"]["candidates"][0]["content"][
                        "parts"
                    ][0]["text"]
                    page_markdown_results[request_id] = markdown_text
                    processed_gemini_pages += 1
                except Exception as e:
                    print(f"Error parsing prediction line: {line}, Error: {e}")
                    # How to map this failure back to a specific page if ID is unparsable?
                    # For now, just count as a general parsing failure for the batch.
                    failed_gemini_pages += 1  # This is a prediction parsing failure, not a page processing failure yet.

            # Aggregate results by original PDF and upload
            # Group file_tasks by original_filename for those in this gemini_batch_sub_job_id
            pdf_page_tasks_cursor = con.execute(
                "SELECT file_task_id, original_filename, page_number, gemini_request_id FROM file_tasks WHERE gemini_batch_sub_job_id = ? ORDER BY original_filename, page_number ASC",
                (gemini_batch_sub_job_id,),
            ).fetchall()

            aggregated_pdfs = {}
            for task_row_tuple in pdf_page_tasks_cursor:
                task_id, pdf_name, page_no, req_id = task_row_tuple
                if pdf_name not in aggregated_pdfs:
                    aggregated_pdfs[pdf_name] = []

                markdown = page_markdown_results.get(req_id)
                if markdown:
                    aggregated_pdfs[pdf_name].append(
                        (page_no, markdown, task_id, "completed")
                    )
                else:
                    # This page's markdown wasn't found in predictions, or failed parsing earlier
                    aggregated_pdfs[pdf_name].append(
                        (
                            page_no,
                            "Error: OCR output not found for this page.",
                            task_id,
                            "failed",
                        )
                    )
                    con.execute(
                        "UPDATE gemini_pdf_batch_sub_jobs SET failed_pdf_pages_count = failed_pdf_pages_count + 1 WHERE gemini_batch_sub_job_id = ?",
                        (gemini_batch_sub_job_id,),
                    )
                    con.commit()

            final_pdfs_processed_count = 0
            final_pdfs_failed_count = 0

            for pdf_name, pages_data in aggregated_pdfs.items():
                pages_data.sort(key=lambda x: x[0])  # Sort by page_number
                full_pdf_markdown = "\n\n---\n\n".join(
                    [p_data[1] for p_data in pages_data]
                )

                output_bucket_name_final, output_prefix_final = (
                    output_gcs_path_str.replace("gs://", "").split("/", 1)
                )
                output_bucket_final = storage_client.bucket(output_bucket_name_final)
                final_blob_name = (
                    output_prefix_final.rstrip("/")
                    + "/"
                    + os.path.splitext(pdf_name)[0]
                    + ".md"
                )
                final_blob_obj = output_bucket_final.blob(final_blob_name)

                try:
                    final_blob_obj.upload_from_string(
                        full_pdf_markdown, content_type="text/markdown"
                    )
                    gcs_final_url = f"gs://{output_bucket_name_final}/{final_blob_name}"
                    # Mark all page tasks for this PDF based on their individual outcomes
                    all_pages_succeeded = True
                    for p_data_item in pages_data:
                        p_task_id, p_status = p_data_item[2], p_data_item[3]
                        p_error_msg = p_data_item[1] if p_status == "failed" else None
                        con.execute(
                            "UPDATE file_tasks SET status = ?, gcs_output_markdown_uri = ?, error_message = ?, updated_at = ? WHERE file_task_id = ?",
                            (
                                p_status,
                                gcs_final_url if p_status == "completed" else None,
                                p_error_msg,
                                datetime.utcnow(),
                                p_task_id,
                            ),
                        )
                        if p_status == "failed":
                            all_pages_succeeded = False

                    if all_pages_succeeded:
                        final_pdfs_processed_count += 1
                        con.execute(
                            "UPDATE gemini_pdf_batch_sub_jobs SET processed_pdf_pages_count = processed_pdf_pages_count + ? WHERE gemini_batch_sub_job_id = ?",
                            (len(pages_data), gemini_batch_sub_job_id),
                        )
                    else:
                        final_pdfs_failed_count += (
                            1  # Count PDF as failed if any page failed
                        )
                    con.commit()
                except Exception as upload_exc:
                    print(
                        f"Failed to upload final aggregated PDF {pdf_name}: {upload_exc}"
                    )
                    final_pdfs_failed_count += 1
                    # Mark all page tasks for this PDF as failed due to final upload error
                    for p_data_item in pages_data:
                        p_task_id = p_data_item[2]
                        con.execute(
                            "UPDATE file_tasks SET status = ?, error_message = ?, updated_at = ? WHERE file_task_id = ?",
                            (
                                "failed",
                                f"Final GCS upload error: {str(upload_exc)}",
                                datetime.utcnow(),
                                p_task_id,
                            ),
                        )
                    con.execute(
                        "UPDATE gemini_pdf_batch_sub_jobs SET failed_pdf_pages_count = failed_pdf_pages_count + ? WHERE gemini_batch_sub_job_id = ?",
                        (len(pages_data), gemini_batch_sub_job_id),
                    )
                    con.commit()

            con.execute(
                "UPDATE batch_jobs SET overall_processed_count = overall_processed_count + ?, overall_failed_count = overall_failed_count + ?, last_updated_at = ? WHERE job_id = ?",
                (
                    final_pdfs_processed_count,
                    final_pdfs_failed_count,
                    datetime.utcnow(),
                    main_batch_job_id,
                ),
            )
            con.execute(
                "UPDATE gemini_pdf_batch_sub_jobs SET status = ?, updated_at = ? WHERE gemini_batch_sub_job_id = ?",
                ("completed", datetime.utcnow(), gemini_batch_sub_job_id),
            )
            con.commit()
            print(
                f"Gemini PDF batch sub job {gemini_batch_sub_job_id} completed successfully."
            )

        else:  # Gemini job did not succeed
            error_msg_from_job = (
                str(gemini_job.error)
                if hasattr(gemini_job, "error") and gemini_job.error
                else "Unknown Gemini job failure"
            )
            print(
                f"Gemini job {gemini_job.name} failed. State: {gemini_job.state}, Error: {error_msg_from_job}"
            )
            con.execute(
                "UPDATE gemini_pdf_batch_sub_jobs SET status = ?, error_message = COALESCE(error_message || CHR(10), '') || ?, updated_at = ? WHERE gemini_batch_sub_job_id = ?",
                (
                    f"failed_gemini_job_{str(gemini_job.state)}",
                    error_msg_from_job,
                    datetime.utcnow(),
                    gemini_batch_sub_job_id,
                ),
            )
            # Mark all associated file_tasks as failed
            con.execute(
                "UPDATE file_tasks SET status='failed', error_message=?, updated_at=? WHERE gemini_batch_sub_job_id = ?",
                (
                    f"Gemini job failed: {error_msg_from_job}",
                    datetime.utcnow(),
                    gemini_batch_sub_job_id,
                ),
            )
            con.execute(
                "UPDATE batch_jobs SET overall_failed_count = overall_failed_count + ? WHERE job_id = ?",
                (len(pdf_inputs_list), main_batch_job_id),
            )  # Approx count
            con.commit()

    except Exception as e:
        error_str = f"Error in _run_gemini_pdf_batch_conversion for sub-job {gemini_batch_sub_job_id}: {str(e)}"
        print(error_str)
        # Attempt to mark sub-job as failed in DB
        try:
            con.execute(
                "UPDATE gemini_pdf_batch_sub_jobs SET status = ?, error_message = COALESCE(error_message || CHR(10), '') || ?, updated_at = ? WHERE gemini_batch_sub_job_id = ?",
                (
                    "failed_internal_error",
                    error_str,
                    datetime.utcnow(),
                    gemini_batch_sub_job_id,
                ),
            )
            # Also mark relevant file_tasks as failed if they weren't already
            con.execute(
                "UPDATE file_tasks SET status='failed', error_message=? WHERE gemini_batch_sub_job_id = ? AND status NOT IN ('completed', 'failed')",
                (
                    f"Internal error in batch: {error_str[:200]}",
                    gemini_batch_sub_job_id,
                ),
            )
            # Update overall batch job with failures - this is an approximation
            con.execute(
                "UPDATE batch_jobs SET overall_failed_count = overall_failed_count + (SELECT COUNT(*) FROM file_tasks WHERE gemini_batch_sub_job_id = ? AND status = 'failed'), last_updated_at = ? WHERE job_id = ?",
                (gemini_batch_sub_job_id, datetime.utcnow(), main_batch_job_id),
            )
            con.commit()
        except Exception as db_err:
            print(
                f"Additionally, failed to update DB for PDF batch sub-job {gemini_batch_sub_job_id} failure: {db_err}"
            )
    finally:
        # Clean up temporary image files from GCS_BATCH_TEMP_BUCKET/images to save costs (optional)
        # This requires listing and deleting, can be a separate cleanup task or done here.
        # For simplicity, not implemented in this iteration.
        # print(f"Cleaning up GCS temp files for {gemini_batch_sub_job_id} - not implemented")
        await _check_and_finalize_batch_job_status(
            main_batch_job_id, con
        )  # Call before closing
        con.commit()  # Ensure final status update is committed
        con.close()


async def _check_and_finalize_batch_job_status(
    main_batch_job_id: str, con: duckdb.DuckDBPyConnection
):
    """Checks if all tasks for a batch job are completed and updates the main job status.
    This function assumes the connection `con` is open and does not close it.
    """
    try:
        job_info = con.execute(
            "SELECT total_input_files, overall_processed_count, overall_failed_count, status FROM batch_jobs WHERE job_id = ?",
            (main_batch_job_id,),
        ).fetchone()

        if not job_info:
            print(
                f"_check_and_finalize_batch_job_status: Batch job {main_batch_job_id} not found."
            )
            return

        total_files, processed_count, failed_count, current_status = job_info

        # If already in a final state, no need to update further by this check.
        if current_status in [
            "completed",
            "completed_with_errors",
            "failed_catastrophic",
            "completed_no_files",
        ]:
            return

        if (processed_count + failed_count) >= total_files:
            new_status = ""
            if failed_count > 0:
                new_status = "completed_with_errors"
            else:
                new_status = "completed"

            if new_status:
                print(
                    f"Finalizing batch job {main_batch_job_id} to status: {new_status}"
                )
                con.execute(
                    "UPDATE batch_jobs SET status = ?, last_updated_at = ? WHERE job_id = ?",
                    (new_status, datetime.utcnow(), main_batch_job_id),
                )
                # No commit here, assuming caller will commit or it's part of a larger transaction within the caller.
                # However, for safety if called at the very end, a commit might be desired directly here.
                # For now, let the caller manage commit after this check.
    except Exception as e:
        print(
            f"Error in _check_and_finalize_batch_job_status for {main_batch_job_id}: {e}"
        )
        # Optionally, update main batch job to a specific error state if this check itself fails critically


# --- Main function for Uvicorn ---
def main():
    import uvicorn

    # Read host and port from environment variables or use defaults
    host = os.getenv("LLM_FOOD_HOST", "0.0.0.0")
    port = int(os.getenv("LLM_FOOD_PORT", "8000"))
    reload = (
        os.getenv("LLM_FOOD_RELOAD", "false").lower() == "true"
    )  # Added reload option

    print(
        f"Starting server on {host}:{port} with reload={'enabled' if reload else 'disabled'}"
    )
    uvicorn.run(
        "llm_food.app:app", host=host, port=port, reload=reload
    )  # Corrected to pass app string for reload


if __name__ == "__main__":
    # This allows running the FastAPI app directly using `python -m llm_food.app`
    main()
