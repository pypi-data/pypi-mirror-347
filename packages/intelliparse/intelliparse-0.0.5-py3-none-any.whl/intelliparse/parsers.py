"""
Module: parsers

This module provides a collection of classes designed for parsing various file types and extracting their content in a structured format.
It offers a flexible and extensible architecture for handling different document formats,
ranging from plain text and office documents to images, audio, video, and specialized formats like XML and compressed archives.

**Core Functionality:**

The central component of this module is the `ParserProtocol` abstract base class.
Concrete implementations of `ParserProtocol` are responsible for handling specific file types.
The module also includes a facade class, `FileParser`, which intelligently routes incoming `RawFile` objects to the appropriate parser based on the file extension.

**Key Features:**

*   **Extensibility:** Easily add support for new file types by creating new classes that inherit from `ParserProtocol` and implement the `parse_async` method.
*   **Strategy-based Parsing:**  The `ParsingStrategy` enum allows you to control the level of detail and processing applied during parsing. Strategies can range from basic text extraction to advanced content analysis using AI agents.
*   **AI-Powered Content Enrichment:** Integration with AI agents (specifically, visual and audio description agents) to generate rich textual descriptions for images, videos, and audio files when using higher parsing strategies.
*   **Structured Output:** Parsers return a `ParsedFile` object, which contains structured content extracted from the file, including sections, text, Markdown representation, images, and items like tables.
*   **Dependency Management:** Uses decorators like `@ensure_module_installed` to manage optional dependencies required for parsing specific file types, enhancing robustness and user experience.

**Module Structure:**

*   **Abstract Base Class:** `ParserProtocol` - Defines the interface for all file parsers.
*   **Facade Parser:** `FileParser` - Acts as the entry point for parsing files, delegating to specific parsers based on file type.
*   **Concrete Parsers:** Classes like `TxtFileParser`, `OfficeFileParser`, `PDFFileParser`, `ImageFileParser`, `AudioFileParser`, `VideoFileParser`, `XMLFileParser`, `CompressedFileParser`, `DWGFileParser`, `PKTFileParser`, `AlgFileParser`, `MarkitdownFileParser` - Implement parsing logic for specific file formats.
*   **Data Structures:** `ParsedFile`, `SectionContent`, `Image`, `TablePageItem` - Define the structured output format for parsed file content.
*   **Enums and Constants:** `ParsingStrategy` - Defines different levels of parsing detail.
*   **Exceptions:** `InvalidFileExtension` - Custom exception raised for unsupported file types.

**Getting Started:**

To parse a file, you would typically use the `FileParser`. First, create a `RawFile` object representing the file you want to parse. Then, instantiate an `FileParser` and call the `parse` or `parse_async` method.

**Example:**

```python
from architecture.data.files import RawFile
from intelliparse.parsers FileParser, ParsingStrategy
import asyncio

async def main():
    # Assuming you have file contents in bytes and the file name
    file_contents = b"This is a sample text file."
    file_name = "example.txt"

    raw_file = RawFile.from_bytes(
        file_contents,
        extension="txt"
    )

    parser = FileParser(strategy=ParsingStrategy.DEFAULT)
    parsed_file = await parser.parse_async(raw_file)

    print(f"Parsed file name: {parsed_file.name}")
    for section in parsed_file.sections:
        print(f"Section {section.number}:")
        print(f"  Text: {section.text[:50]}...") # Print first 50 characters
        print(f"  Markdown: {section.md[:50]}...") # Print first 50 characters

if __name__ == "__main__":
    asyncio.run(main())
```

This module provides a robust and versatile solution for file content extraction, adaptable to various needs and file formats.
"""

from __future__ import annotations

import abc
import asyncio
import io
import logging
import os
import subprocess
import tempfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Callable, Literal, Never, Optional, Sequence, cast, override
from xml.etree.ElementTree import Element

import msgspec
from architecture import dp, log
from architecture.data.files import RawFile, bytes_to_mime, find_extension
from architecture.utils.decorators import ensure_module_installed
from architecture.utils.functions import run_sync
from intellibricks.agents import Agent
from intellibricks.llms.types import (
    AudioDescription,
    AudioFilePart,
    ChainOfThought,
    ImageFilePart,
    VideoFilePart,
    VisualMediaDescription,
)

from .types import Image, ParsedFile, SectionContent, TablePageItem

debug_logger = log.create_logger(__name__, level=logging.DEBUG)
exception_logger = log.create_logger(__name__, level=logging.ERROR)

_parser_registry: dict[str, type[FileParser]] = {}


def _parses(
    *extensions: str,
) -> Callable[[type[FileParser]], type[FileParser]]:
    """
    Decorator to register a file parser class for specific file extensions.

    This decorator is used to register a file parser class for specific file extensions.
    It adds the parser class to the global parser registry, allowing the `FileParser`
    to automatically select the correct parser based on the file extension. Should be
    used internally only by concrete intellibricks parser classes.

    **Parameters:**

    *   `extensions` (str): One or more file extensions that the parser class supports.

    **Returns:**

    *   `Callable[[type[FileParser]], type[FileParser]]`: A decorator function that registers the parser class.

    **Example:**

    ```python
    from intelliparse.parsers import FileParser, parses

    @_parses("txt")
    class CustomTxtFileParser(FileParser):
        async def parse_async(self, file: RawFile) -> ParsedFile:
            # Add parsing logic here
            pass
    ```
    """

    def decorator(
        parser_cls: type[FileParser],
    ) -> type[FileParser]:
        for extension in extensions:
            _parser_registry[extension] = parser_cls
        return parser_cls

    return decorator


class InvalidFileExtension(Exception):
    """
    Exception raised when a file extension is not supported by the parsers.

    This exception is typically raised by the `FileParser` when it encounters
    a `RawFile` with an extension that no specific parser is registered to handle.

    **Example:**

    ```python
    from architecture.data.files import RawFile
    from intelliparse.parsers import FileParser, InvalidFileExtension, ParsingStrategy
    import asyncio

    async def main():
        file_contents = b"This is an unsupported file."
        file_name = "unknown.extension"
        raw_file = RawFile.from_bytes(
            data=file_contents,
            name=file_name,
            extension="intellibricks_is_the_best"
        )

        parser = FileParser(strategy=ParsingStrategy.DEFAULT)
        try:
            parsed_file = await parser.parse_async(raw_file)
        except InvalidFileExtension as e:
            print(f"Error: {e}") # Output: Error: Unsupported file extension: unknown

    if __name__ == "__main__":
        asyncio.run(main())
    ```
    """


class ParserProtocol(msgspec.Struct, frozen=True, tag_field="type"):
    """
    Abstract base class for file parsers.

    This class defines the interface that all concrete file parser classes in this module must implement.
    It provides a basic structure and common functionality for parsing files and extracting their content.

    **Key Features:**

    *   **Abstract Base Class:** Cannot be instantiated directly. Subclasses must implement the abstract method `parse_async`.
    *   **Strategy-Based Parsing:**  Includes a `strategy` attribute of type `ParsingStrategy` to control the parsing level.
    *   **Synchronous and Asynchronous Interface:** Provides both synchronous (`parse`) and asynchronous (`parse_async`) methods for content extraction. The synchronous method is a wrapper around the asynchronous one, using `run_sync`.

    **Attributes:**

    *   `strategy` (ParsingStrategy):  The parsing strategy to be used. Defaults to `ParsingStrategy.DEFAULT`.

    **Methods:**

    *   `parse(file: RawFile) -> ParsedFile`:
        Synchronously extracts content from a `RawFile`. This method is a convenience wrapper around `parse_async`.

        **Parameters:**

        *   `file` (RawFile): The `RawFile` object representing the file to be parsed.

        **Returns:**

        *   `ParsedFile`: A `ParsedFile` object containing the extracted content.

        **Example:**

        ```python
        from architecture.data.files import RawFile
        from intelliparse.parsers TxtFileParser, ParserProtocol, ParsedFile, ParsingStrategy

        # Assume you have raw file content and name
        file_content = b"This is a text file content."
        file_name = "document.txt"
        raw_file = RawFile.from_bytes(file_content, "txt")

        parser: ParserProtocol = TxtFileParser(strategy=ParsingStrategy.DEFAULT)
        parsed_file: ParsedFile = parser.parse(raw_file)

        print(f"Parsed file name: {parsed_file.name}") # Output: Parsed file name: document.txt
        print(f"Section 1 Text: {parsed_file.sections[0].text}") # Output: Section 1 Text: This is a text file content.
        ```

    *   `parse_async(file: RawFile) -> ParsedFile`:
        **Abstract method** that must be implemented by subclasses. Asynchronously extracts content from a `RawFile`.

        **Parameters:**

        *   `file` (RawFile): The `RawFile` object representing the file to be parsed.

        **Returns:**

        *   `ParsedFile`: A `ParsedFile` object containing the extracted content.

        **Raises:**

        *   `NotImplementedError`: If the method is not implemented in a subclass.

        **Note:**

        Subclasses should override `parse_async` to provide specific parsing logic for their supported file types.
    """

    def parse(
        self,
        file: RawFile,
    ) -> ParsedFile:
        """Extracts content from the file."""
        return run_sync(self.parse_async, file)

    @abc.abstractmethod
    async def parse_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        """Extracts content from the file."""
        raise NotImplementedError("This method should be implemented by subclasses.")


@dp.Facade
class FileParser(ParserProtocol, frozen=True, tag="parser"):
    """
    Facade class that acts as the main entry point for parsing files.

    This class implements the `ParserProtocol` interface and automatically delegates the parsing process
    to the appropriate specialized parser based on the file extension of the input `RawFile`.
    It supports a wide range of file types, including office documents, PDFs, images, audio, video, and more.

    **Key Features:**

    *   **File Type Dispatch:**  Automatically selects the correct parser based on the file extension.
    *   **Facade Pattern:** Simplifies the parsing process by providing a single class to handle various file types.
    *   **AI Agent Integration:** Optionally integrates with visual and audio description agents to enhance content extraction, especially in `HIGH` parsing strategy.
    *   **Configurable Strategy:** Inherits the `strategy` attribute from `ParserProtocol` to control the parsing level.

    **Attributes:**

    *   `strategy` (ParsingStrategy): The parsing strategy to be used. Defaults to `ParsingStrategy.DEFAULT`.
    *   `visual_description_agent` (Optional[Agent[ChainOfThought[VisualMediaDescription]]]):
        Optional agent for generating textual descriptions of images and videos. Used when the parsing strategy is `HIGH` and for file types that support visual content analysis.
    *   `audio_description_agent` (Optional[Agent[ChainOfThought[AudioDescription]]]):
        Optional agent for generating textual descriptions of audio files. Used when the parsing strategy is `HIGH` and for audio file types.

    **Methods:**

    *   `parse_async(file: RawFile) -> ParsedFile`:
        Asynchronously extracts content from the provided `RawFile`. This method determines the file type and delegates the actual parsing to the appropriate specialized parser.

        **Parameters:**

        *   `file` (RawFile): The `RawFile` object representing the file to be parsed.

        **Returns:**

        *   `ParsedFile`: A `ParsedFile` object containing the extracted content.

        **Raises:**

        *   `InvalidFileExtension`: If the file extension is not supported.

        **Example:**

        ```python
        from architecture.data.files import RawFile
        from intelliparse.parsers FileParser, ParsingStrategy, ParsedFile
        import asyncio

        async def main():
            # Example with a DOCX file (assuming you have docx_content in bytes)
            docx_content = b"..." # Your DOCX file content here
            docx_file = RawFile.from_bytes(docx_content, "docx")

            parser = FileParser(strategy=ParsingStrategy.HIGH) # Or DEFAULT, MEDIUM, FAST
            parsed_docx: ParsedFile = await parser.parse_async(docx_file)

            print(f"Parsed DOCX file name: {parsed_docx.name}")
            # Access parsed content from parsed_docx.sections

            # Example with a TXT file
            txt_content = b"This is a plain text file."
            txt_file = RawFile.from_bytes(txt_content, "txt")
            parsed_txt: ParsedFile = await parser.parse_async(txt_file)

            print(f"Parsed TXT file name: {parsed_txt.name}")
            # Access parsed content from parsed_txt.sections

        if __name__ == "__main__":
            asyncio.run(main())
        ```

    **Usage Notes:**

    *   This is the recommended class to use for parsing files in most scenarios.
    *   You can configure the parsing behavior by setting the `strategy` attribute.
    *   To enable AI-powered descriptions for images, videos, and audio, provide instances of `visual_description_agent` and/or `audio_description_agent` during initialization.
    """

    strategy: Literal["low", "medium", "high"] = msgspec.field(
        default_factory=lambda: "low"
    )

    visual_description_agent: Optional[
        Agent[ChainOfThought[VisualMediaDescription]]
    ] = None
    """Agent used for generating textual descriptions of images and videos, if the synapse supports it."""

    audio_description_agent: Optional[Agent[ChainOfThought[AudioDescription]]] = None
    """Agent used for generating textual descriptions of audio files, if the synapse supports it."""

    @override
    async def parse_async(self, file: RawFile) -> ParsedFile:
        parser_cls = _parser_registry.get(file.extension)

        if not parser_cls:
            raise InvalidFileExtension(f"Unsupported extension: {file.extension}")

        return await parser_cls(
            strategy=self.strategy,
            visual_description_agent=self.visual_description_agent,
            audio_description_agent=self.audio_description_agent,
        ).parse_async(file)


@_parses("xml")
class XMLFileParser(FileParser, frozen=True, tag="xml"):
    """
    Parser for XML files (.xml).

    This parser extracts the content from XML files and represents it in two formats:

    *   **Raw XML Text:** The original XML content is preserved in the `text` field of the `SectionContent`.
    *   **Markdown Representation:** The XML structure is converted into a nested Markdown list format, making it more human-readable. This Markdown representation is stored in the `md` field of the `SectionContent`.

    **Key Features:**

    *   **XML to Markdown Conversion:** Transforms XML structure into a clear Markdown representation.
    *   **Error Handling:** Gracefully handles XML parsing errors by falling back to displaying the raw XML within a Markdown code block.
    *   **Single Section Output:**  Outputs the entire XML content as a single section in the `ParsedFile`.

    **Methods:**

    *   `parse_async(file: RawFile) -> ParsedFile`:
        Asynchronously extracts content from an XML `RawFile`.

        **Parameters:**

        *   `file` (RawFile): The `RawFile` object representing the XML file.

        **Returns:**

        *   `ParsedFile`: A `ParsedFile` object containing a single section with the XML content in both raw text and Markdown formats.

        **Example:**

        ```python
        from architecture.data.files import RawFile
        from intelliparse.parsers XMLFileParser, ParsedFile, ParsingStrategy
        import asyncio

        async def main():
            xml_content = b"..."
            xml_file = RawFile.from_bytes(xml_content, "xml")

            parser = XMLFileParser(strategy=ParsingStrategy.DEFAULT)
            parsed_xml: ParsedFile = await parser.parse_async(xml_file)

            section = parsed_xml.sections[0]
            print(f"Section 1 - Raw XML Text:\n{section.text}")
            print(f"Section 1 - Markdown Representation:\n{section.md}")

        if __name__ == "__main__":
            asyncio.run(main())
        ```

    *   `xml_to_md(xml_str: str) -> str`:
        Converts a string containing XML content into a nested Markdown list structure.

        **Parameters:**

        *   `xml_str` (str): The XML content as a string.

        **Returns:**

        *   `str`: The Markdown representation of the XML.

    *   `_convert_element_to_md(element: Element, level: int) -> str`:
        Recursive helper method to convert an XML element and its children to Markdown. (Internal use)
    """

    @override
    async def parse_async(self, file: RawFile) -> ParsedFile:
        raw_xml = file.contents.decode("utf-8", errors="replace")
        md_content = self.xml_to_md(raw_xml)

        section_content = SectionContent(
            number=1,
            text=raw_xml,
            md=md_content,
            images=[],
            items=[],
        )

        return ParsedFile(
            name=file.name,
            sections=[section_content],
        )

    def xml_to_md(self, xml_str: str) -> str:
        """Converts XML content into a nested Markdown list structure."""
        try:
            root: Element = ET.fromstring(xml_str)
            return self._convert_element_to_md(root, level=0)
        except ET.ParseError as e:
            exception_logger.exception("Error parsing XML: %s", e)
            return "```xml\n" + xml_str + "\n```"  # Fallback to raw XML in code block

    def _convert_element_to_md(self, element: Element, level: int) -> str:
        """Recursively converts an XML element and its children to Markdown.

        Args:
            element: The XML element to convert
            level: Current nesting level for indentation
        """
        indent = "  " * level
        lines: list[str] = []

        # Element tag as bold item
        lines.append(f"{indent}- **{element.tag}**")

        # Attributes as sub-items
        if element.attrib:
            lines.append(f"{indent}  - *Attributes*:")
            for key, value in element.attrib.items():
                lines.append(f"{indent}    - `{key}`: `{value}`")

        # Text content
        if element.text and element.text.strip():
            text = element.text.strip().replace("\n", " ")
            lines.append(f"{indent}  - *Text*: {text}")

        # Process child elements recursively
        for child in element:
            lines.append(self._convert_element_to_md(child, level + 1))

        return "\n".join(lines)


@_parses("zip", "rar", "pkz")
class CompressedFileParser(FileParser, frozen=True, tag="compressed"):
    """
    Parser for compressed archive files (ZIP, RAR, PKZ).

    This parser handles compressed files by:

    1.  **Extracting Contents:**  Extracting all files within the archive to a temporary directory.
    2.  **Recursive Parsing:**  For each extracted file, it uses the `FileParser` facade to determine the appropriate parser and extract content.
    3.  **Merging Results:**  Combines the `ParsedFile` objects from each extracted file into a single `ParsedFile` representing the entire archive.

    **Supported Formats:**

    *   ZIP (.zip, .pkz)
    *   RAR (.rar)

    **Key Features:**

    *   **Archive Handling:**  Supports common archive formats to extract and parse contained files.
    *   **Recursive Parsing:** Leverages the `FileParser` to handle diverse file types within archives.
    *   **Merged Output:** Provides a single `ParsedFile` representing the combined content of all files in the archive.

    **Methods:**

    *   `parse_async(file: RawFile) -> ParsedFile`:
        Asynchronously extracts and parses the contents of a compressed `RawFile`.

        **Parameters:**

        *   `file` (RawFile): The `RawFile` object representing the compressed archive file.

        **Returns:**

        *   `ParsedFile`: A `ParsedFile` object containing the merged content of all parsed files from within the archive.

        **Raises:**

        *   `ValueError`: If the file extension is not a supported compressed file type.

        **Example:**

        ```python
        from architecture.data.files import RawFile
        from intelliparse.parsers CompressedFileParser, ParsedFile, ParsingStrategy
        import asyncio
        import zipfile
        import io

        async def main():
            # Create a dummy ZIP file in memory for example
            buffer = io.BytesIO()
            with zipfile.ZipFile(buffer, 'w') as zf:
                zf.writestr("text_file.txt", "This is a text file inside ZIP.")
                zf.writestr("image.png", b"PNG Content Here") # Replace with actual PNG bytes

            zip_content = buffer.getvalue()
            zip_file = RawFile.from_bytes(zip_content, "zip")

            parser = CompressedFileParser(strategy=ParsingStrategy.DEFAULT)
            parsed_zip: ParsedFile = await parser.parse_async(zip_file)

            print(f"Parsed ZIP file name: {parsed_zip.name}")
            for parsed_child_file in parsed_zip.parsed_files:
                print(f"  Child File: {parsed_child_file.name}")
                for section in parsed_child_file.sections:
                    print(f"    Section {section.number}: Text: {section.text[:30]}...") # First 30 chars

        if __name__ == "__main__":
            asyncio.run(main())
        ```

    **Usage Notes:**

    *   Requires the `zipfile` and `rarfile` libraries to be installed (automatically handled if you installed `intellibricks[files]`).
    *   Handles nested archives recursively (parses files within archives within archives).
    """

    @override
    async def parse_async(self, file: RawFile) -> ParsedFile:
        import tempfile
        import zipfile

        import rarfile

        # We'll accumulate ParsedFile objects from each extracted child file
        parsed_files: list[ParsedFile] = []

        # Write the compressed file to a temporary location
        with tempfile.NamedTemporaryFile(delete=True) as tmp:
            tmp.write(file.contents)
            tmp.flush()

            # Decide how to open the archive based on extension
            match file.extension:
                case "zip" | "pkz":
                    # Treat PKZ exactly like ZIP for demo purposes
                    with zipfile.ZipFile(tmp.name, "r") as zip_ref:
                        # Iterate over files inside the archive
                        for info in zip_ref.infolist():
                            # Directories have filename ending with "/"
                            if info.is_dir():
                                continue

                            # Read raw bytes of the child file
                            child_data = zip_ref.read(info)
                            child_name = info.filename
                            child_ext = find_extension(filename=child_name)

                            # Turn that child file into a RawFile
                            child_raw_file = RawFile.from_bytes(
                                data=child_data,
                                name=child_name,
                                extension=child_ext,
                            )

                            # Parse using our FileParser façade
                            # (re-using the same strategy/visual_description_agent)
                            parser = FileParser(
                                strategy=self.strategy,
                                visual_description_agent=self.visual_description_agent,
                            )
                            child_parsed = await parser.parse_async(child_raw_file)
                            parsed_files.append(child_parsed)

                case "rar":
                    with rarfile.RarFile(tmp.name, "r") as rar_ref:
                        for info in rar_ref.infolist():
                            """Type of "isdir" is unknownPylancereportUnknownMemberType"""
                            if info.isdir():  # type: ignore
                                continue

                            child_data = rar_ref.read(info)  # type: ignore

                            child_name = info.filename  # type: ignore

                            child_ext = find_extension(filename=child_name)  # type: ignore

                            child_raw_file = RawFile.from_bytes(
                                data=child_data,  # type: ignore
                                name=child_name,  # type: ignore
                                extension=child_ext,
                            )

                            parser = FileParser(
                                strategy=self.strategy,
                                visual_description_agent=self.visual_description_agent,
                            )
                            child_parsed = await parser.parse_async(child_raw_file)
                            parsed_files.append(child_parsed)

                case _:
                    # Fallback if something else accidentally calls this parser
                    raise ValueError(
                        f"CompressedFileParser does not handle extension: {file.extension}"
                    )

        # Merge all the parsed files into a single ParsedFile
        return ParsedFile.from_parsed_files(parsed_files)


@_parses("dwg")
class DWGFileParser(FileParser, frozen=True, tag="dwg"):
    """
    Parser for DWG (Drawing) files (.dwg).

    This parser handles DWG files by converting them to PDF format first, and then extracting
    images of each page from the PDF. It then uses the `StaticImageFileParser` to process each page image,
    potentially utilizing a visual description agent if configured and the parsing strategy is set to `HIGH`.

    **Key Features:**

    *   **DWG to PDF Conversion:** Converts DWG files to PDF format using the `aspose-cad` library.
    *   **PDF Page to Image Extraction:** Extracts each page of the converted PDF as a PNG image.
    *   **Image Parsing Delegation:** Uses `StaticImageFileParser` to process the extracted images, enabling text extraction and visual description if configured.
    *   **Multi-Page DWG Support:** Handles multi-page DWG files by processing each page individually.

    **Dependencies:**

    *   Requires the `aspose-cad` and `pymupdf` libraries to be installed (automatically handled if you installed `intellibricks[files]`).
    *   Requires a non-ARM64 architecture due to `aspose-cad` limitations.

    **Methods:**

    *   `parse_async(file: RawFile) -> ParsedFile`:
        Asynchronously extracts content from a DWG `RawFile`.

        **Parameters:**

        *   `file` (RawFile): The `RawFile` object representing the DWG file.

        **Returns:**

        *   `ParsedFile`: A `ParsedFile` object containing sections, each representing a page from the DWG, processed as an image.

        **Raises:**

        *   `ValueError`: If running on ARM64 architecture.
        *   `RuntimeError`: If `aspose-cad` module is not installed or if conversion fails.

        **Example:**

        ```python
        from architecture.data.files import RawFile
        from intelliparse.parsers DWGFileParser, ParsedFile, ParsingStrategy
        import asyncio

        async def main():
            # Assume you have dwg_content in bytes
            dwg_content = b"..." # Your DWG file content in bytes
            dwg_file = RawFile.from_bytes(dwg_content, "dwg")

            parser = DWGFileParser(strategy=ParsingStrategy.HIGH) # Or DEFAULT, MEDIUM, FAST
            parsed_dwg: ParsedFile = await parser.parse_async(dwg_file)

            print(f"Parsed DWG file name: {parsed_dwg.name}")
            for section in parsed_dwg.sections:
                print(f"  Section {section.number}: Images: {len(section.images)}") # Number of images per page

        if __name__ == "__main__":
            asyncio.run(main())
        ```

    **Usage Notes:**

    *   Due to the conversion process, the output `ParsedFile` will primarily contain images representing the DWG content, rather than text.
    *   Make sure to install the required dependencies (`aspose-cad`, `pymupdf`) before using this parser.
    *   This parser is resource-intensive due to format conversion and image processing.
    """

    @ensure_module_installed("aspose-cad", "intellibricks[files]")
    @override
    async def parse_async(self, file: RawFile) -> ParsedFile:
        """
        DWG files are kind of tricky. To parse them, Intellibricks converts them to PDF first,
        then takes a "screenshot" of each page of the PDF and uses GenAI to describe the images.
        """
        import platform

        if platform.machine() == "arm64":
            raise ValueError("ARM architecture is not supported by aspose-cad")

        import aspose.cad as cad  # type: ignore

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            file.save_to_file(file_path)

            # Load the DWG file
            image = cad.Image.load(file_path)  # type: ignore

            # Specify PDF Options
            pdfOptions = cad.imageoptions.PdfOptions()  # type: ignore

            output_path = f"{temp_dir}/output.pdf"

            # Save as PDF
            image.save(output_path, pdfOptions)  # type: ignore

            image_bytes_list = self.__pdf_to_images(output_path)

            raw_files = [
                RawFile.from_bytes(
                    data=img, name=f"{file.name}_{i}.png", extension="png"
                )
                for i, img in enumerate(image_bytes_list)
            ]

            parser = StaticImageFileParser(
                strategy=self.strategy,
                visual_description_agent=self.visual_description_agent,
            )

            parsed_files = [await parser.parse_async(f) for f in raw_files]
            sections = [
                section
                for parsed_file in parsed_files
                for section in parsed_file.sections
            ]

            return ParsedFile.from_sections(file.name, sections)

    def __pdf_to_images(self, pdf_path: str) -> Sequence[bytes]:
        """Converts each page of a PDF to image bytes.

        Args:
            pdf_path (str): The path to the PDF file.

        Returns:
            Sequence[bytes]: A list of bytes objects, each containing a PNG image of a PDF page.
        """
        import pymupdf

        image_bytes_list: list[bytes] = []
        doc = pymupdf.open(pdf_path)

        try:
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)  # type: ignore
                pix = page.get_pixmap()  # type: ignore

                # Create a bytes buffer and save the image into it
                buffer = io.BytesIO()
                pix.save(buffer, "png")  # type: ignore
                image_bytes = buffer.getvalue()

                image_bytes_list.append(image_bytes)

        finally:
            doc.close()

        return image_bytes_list


@_parses("pkt")
class PKTFileParser(FileParser, frozen=True, tag="pkt"):
    """
    Parser for Packet Tracer files (.pkt, .pka).

    This parser extracts the XML representation from Packet Tracer files. Packet Tracer files are
    compressed and encrypted, so this parser handles decryption and decompression to obtain the underlying XML content.

    **Key Features:**

    *   **Packet Tracer Format Support:** Specifically designed to parse .pkt and .pka files.
    *   **Decryption and Decompression:** Handles the necessary steps to decrypt and decompress the file content to reveal the XML data.
    *   **XML Output:** Extracts the content as XML text, which is stored in both `text` and `md` fields of the `SectionContent`.
    *   **Single Section Output:**  Outputs the entire XML content as a single section in the `ParsedFile`.

    **Methods:**

    *   `parse_async(file: RawFile) -> ParsedFile`:
        Asynchronously extracts content from a Packet Tracer `RawFile`.

        **Parameters:**

        *   `file` (RawFile): The `RawFile` object representing the Packet Tracer file (.pkt or .pka).

        **Returns:**

        *   `ParsedFile`: A `ParsedFile` object containing a single section with the extracted XML content.

        **Example:**

        ```python
        from architecture.data.files import RawFile
        from intelliparse.parsers PKTFileParser, ParsedFile, ParsingStrategy
        import asyncio

        async def main():
            # Assume you have pkt_content in bytes
            pkt_content = b"..." # Your PKT file content in bytes
            pkt_file = RawFile.from_bytes(pkt_content, "pkt")

            parser = PKTFileParser(strategy=ParsingStrategy.DEFAULT)
            parsed_pkt: ParsedFile = await parser.parse_async(pkt_file)

            section = parsed_pkt.sections[0]
            print(f"Parsed PKT file name: {parsed_pkt.name}")
            print(f"Section 1 - XML Content:\n{section.text[:100]}...") # Print first 100 chars of XML

        if __name__ == "__main__":
            asyncio.run(main())
        ```

    *   `pkt_to_xml_bytes(pkt_file: str) -> bytes`:
        Converts a Packet Tracer file (.pkt/.pka) to its XML representation as bytes.

        **Parameters:**

        *   `pkt_file` (str): Path to the input .pkt or .pka file.

        **Returns:**

        *   `bytes`: The uncompressed XML content as bytes.

    **Usage Notes:**

    *   This parser is specific to Cisco Packet Tracer files.
    *   The extracted content is in XML format, representing the network topology and configuration.
    """

    @override
    async def parse_async(self, file: RawFile) -> ParsedFile:
        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            file.save_to_file(file_path)

            xml_bytes = self.pkt_to_xml_bytes(file_path)

            # For now, we'll just return the XML content as a single page
            xml_text = xml_bytes.decode("utf-8", errors="replace")

            page_content = SectionContent(
                number=1,
                text=xml_text,
                md=xml_text,
                images=[],
                items=[],
            )

            return ParsedFile(
                name=file.name,
                sections=[page_content],
            )

    def pkt_to_xml_bytes(self, pkt_file: str) -> bytes:
        """
        Convert a Packet Tracer file (.pkt/.pka) to its XML representation as bytes.

        :param pkt_file: Path to the input .pkt or .pka file.
        :return: The uncompressed XML content as bytes.
        """
        import zlib

        with open(pkt_file, "rb") as f:
            in_data = bytearray(f.read())

        i_size = len(in_data)
        out = bytearray()

        # Decrypt each byte with decreasing file length
        for byte in in_data:
            out.append(byte ^ (i_size & 0xFF))
            i_size -= 1

        # The first 4 bytes (big-endian) represent the size of the XML when uncompressed
        # (This value is not needed for the actual return, but we parse it for completeness.)
        _uncompressed_size = int.from_bytes(out[:4], byteorder="big")

        # Decompress the data after the first 4 bytes
        xml_data = zlib.decompress(out[4:])

        return xml_data


@_parses("pdf")
class PDFFileParser(FileParser, frozen=True, tag="pdf"):
    """
    Parser for PDF files (.pdf).

    This parser extracts text and images from PDF files. For higher parsing strategies (`ParsingStrategy.HIGH`),
    it can also utilize a visual description agent to generate textual descriptions of images found within the PDF.

    **Key Features:**

    *   **Text Extraction:** Extracts text content from each page of the PDF.
    *   **Image Extraction:** Extracts images embedded within the PDF pages.
    *   **Visual Description (Optional):**  When `ParsingStrategy.HIGH` is used and a `visual_description_agent` is provided, it generates AI-powered descriptions for extracted images.
    *   **Multi-Page Support:** Processes multi-page PDFs, creating a `SectionContent` for each page.

    **Dependencies:**

    *   Requires the `pypdf` library to be installed (automatically handled if you installed `intellibricks[files]`).

    **Methods:**

    *   `parse_async(file: RawFile) -> ParsedFile`:
        Asynchronously extracts content from a PDF `RawFile`.

        **Parameters:**

        *   `file` (RawFile): The `RawFile` object representing the PDF file.

        **Returns:**

        *   `ParsedFile`: A `ParsedFile` object containing sections, each representing a page from the PDF, with extracted text and images.

        **Example:**

        ```python
        from architecture.data.files import RawFile
        from intelliparse.parsers PDFFileParser, ParsedFile, ParsingStrategy
        import asyncio

        async def main():
            # Assume you have pdf_content in bytes
            pdf_content = b"..." # Your PDF file content in bytes
            pdf_file = RawFile.from_bytes(pdf_content, "pdf")

            parser = PDFFileParser(strategy=ParsingStrategy.HIGH) # Or DEFAULT, MEDIUM, FAST
            parsed_pdf: ParsedFile = await parser.parse_async(pdf_file)

            print(f"Parsed PDF file name: {parsed_pdf.name}")
            for section in parsed_pdf.sections:
                print(f"  Section {section.number}: Text: {section.text[:50]}...") # First 50 chars of text
                print(f"  Section {section.number}: Images: {len(section.images)}") # Number of images on page

        if __name__ == "__main__":
            asyncio.run(main())
        ```

    **Usage Notes:**

    *   For `ParsingStrategy.HIGH`, ensure you have configured a `visual_description_agent` in the `FileParser` or `PDFFileParser` instance.
    *   The quality of text extraction can vary depending on the PDF's structure and whether it's text-based or image-based.
    """

    @ensure_module_installed("pypdf", "intellibricks[files]")
    @override
    async def parse_async(self, file: RawFile) -> ParsedFile:
        import hashlib

        from pypdf import PdfReader

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            file.save_to_file(file_path)

            reader = PdfReader(file_path)
            section_contents: list[SectionContent] = []
            image_cache: dict[str, tuple[str, str]] = {}

            for page_num, page in enumerate(reader.pages):
                page_images: list[Image] = []
                image_descriptions: list[str] = []

                if self.visual_description_agent and self.strategy == "high":
                    for image_num, image in enumerate(page.images):
                        image_bytes = image.data
                        image_hash = hashlib.sha256(image_bytes).hexdigest()

                        if image_hash in image_cache:
                            cached_md, cached_ocr = image_cache[image_hash]
                            image_md = cached_md
                            ocr_text = cached_ocr
                        else:
                            agent_input = ImageFilePart(
                                mime_type=bytes_to_mime(image.data), data=image.data
                            )
                            agent_response = (
                                await self.visual_description_agent.run_async(
                                    agent_input
                                )
                            )
                            image_md = agent_response.parsed.final_answer.md
                            ocr_text = agent_response.parsed.final_answer.ocr_text
                            image_cache[image_hash] = (image_md, ocr_text)

                        image_descriptions.append(
                            f"Page Image {image_num + 1}: {image_md}"
                        )
                        page_images.append(
                            Image(
                                contents=image.data,
                                name=image.name,
                                ocr_text=ocr_text,
                            )
                        )

                page_text = [page.extract_text(), "".join(image_descriptions)]
                md = "".join(page_text)
                section_content = SectionContent(
                    number=page_num + 1,
                    text=md,
                    md=md,
                    images=page_images,
                )
                section_contents.append(section_content)

            return ParsedFile(
                name=file.name,
                sections=section_contents,
            )


@_parses("doc", "docx")
class DocxFileParser(FileParser, frozen=True, tag="docx"):
    """
    Parser for DOCX (Microsoft Word) files (.docx, .doc).

    This parser extracts text and images from DOCX files. For higher parsing strategies (`ParsingStrategy.HIGH`),
    it can also utilize a visual description agent to generate textual descriptions of images found within the document.

    **Key Features:**

    *   **Text Extraction:** Extracts text content from paragraphs in the DOCX document.
    *   **Image Extraction:** Extracts images embedded within the DOCX document.
    *   **Visual Description (Optional):** When `ParsingStrategy.HIGH` is used and a `visual_description_agent` is provided, it generates AI-powered descriptions for extracted images.
    *   **Single Section Output:**  Outputs the entire DOCX content as a single section in the `ParsedFile` (DOCX files do not inherently have pages like PDFs).
    *   **.doc Conversion:** Automatically converts older .doc formats to .docx using LibreOffice before parsing.

    **Dependencies:**

    *   Requires the `docx` library (python-docx) to be installed (automatically handled if you installed `intellibricks[files]`).
    *   Requires LibreOffice to be installed and in the system PATH for .doc conversion.

    **Methods:**

    *   `parse_async(file: RawFile) -> ParsedFile`:
        Asynchronously extracts content from a DOCX `RawFile`.

        **Parameters:**

        *   `file` (RawFile): The `RawFile` object representing the DOCX file.

        **Returns:**

        *   `ParsedFile`: A `ParsedFile` object with a single section containing the extracted text and images from the DOCX file.

        **Example:**

        ```python
        from architecture.data.files import RawFile
        from intelliparse.parsers DocxFileParser, ParsedFile, ParsingStrategy
        import asyncio

        async def main():
            # Assume you have docx_content in bytes
            docx_content = b"..." # Your DOCX file content in bytes
            docx_file = RawFile.from_bytes(docx_content, "docx")

            parser = DocxFileParser(strategy=ParsingStrategy.HIGH) # Or DEFAULT, MEDIUM, FAST
            parsed_docx: ParsedFile = await parser.parse_async(docx_file)

            print(f"Parsed DOCX file name: {parsed_docx.name}")
            section = parsed_docx.sections[0]
            print(f"Section 1 - Text: {section.text[:100]}...") # First 100 chars of text
            print(f"Section 1 - Images: {len(section.images)}") # Number of images in document

        if __name__ == "__main__":
            asyncio.run(main())
        ```

    **Usage Notes:**

    *   For `ParsingStrategy.HIGH`, ensure you have configured a `visual_description_agent` in the `FileParser` or `DocxFileParser` instance.
    *   The output `ParsedFile` contains a single section as DOCX documents are typically treated as a continuous flow of content rather than distinct pages.
    """

    @ensure_module_installed("docx", "intellibricks[files]")
    @override
    async def parse_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        import hashlib

        from docx import Document

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            
            # Check if we need to convert from .doc to .docx
            if file.extension == "doc":
                try:
                    file.save_to_file(file_path)
                    converted_docx_file = self._convert_to_docx(file)
                    converted_docx_file.save_to_file(file_path)
                except RuntimeError as e:
                    # If conversion fails, log the error and try with original file
                    exception_logger.exception(f"Doc to docx conversion failed: {str(e)}. Trying with original file.")
                    file.save_to_file(file_path)
            else:
                file.save_to_file(file_path)

            try:
                document = Document(file_path)
            except Exception as e:
                # If document still fails to load, return a simple error message
                exception_logger.exception(f"Failed to load document: {str(e)}")
                return ParsedFile(
                    name=file.name,
                    sections=[
                        SectionContent(
                            number=1,
                            text=f"Error: Could not parse file. The document may be corrupted or in an unsupported format.",
                            md=f"Error: Could not parse file. The document may be corrupted or in an unsupported format.",
                            images=[],
                        )
                    ],
                )
                
            image_cache: dict[str, tuple[str, str]] = {}  # (md, ocr_text)

            paragraph_texts = [p.text for p in document.paragraphs if p.text.strip()]
            doc_text = "\n".join(paragraph_texts)

            doc_images: list[tuple[str, bytes]] = []
            for rel in document.part._rels.values():  # type: ignore[reportPrivateUsage]
                if "image" in rel.reltype:
                    image_part = rel.target_part
                    image_name = image_part.partname.split("/")[-1]
                    image_bytes = image_part.blob
                    doc_images.append((image_name, image_bytes))

            final_images: list[Image] = []
            image_descriptions: list[str] = []
            if self.visual_description_agent and self.strategy == "high":
                for idx, (image_name, image_bytes) in enumerate(doc_images, start=1):
                    image_hash = hashlib.sha256(image_bytes).hexdigest()

                    if image_hash in image_cache:
                        cached_md, cached_ocr = image_cache[image_hash]
                        image_md = cached_md
                        ocr_text = cached_ocr
                    else:
                        agent_input = ImageFilePart(
                            mime_type=bytes_to_mime(image_bytes),
                            data=image_bytes,
                        )
                        agent_response = await self.visual_description_agent.run_async(
                            agent_input
                        )
                        image_md = agent_response.parsed.final_answer.md
                        ocr_text = agent_response.parsed.final_answer.ocr_text
                        image_cache[image_hash] = (image_md, ocr_text)

                    image_descriptions.append(f"Docx Image {idx}: {image_md}")
                    final_images.append(
                        Image(
                            name=image_name,
                            contents=image_bytes,
                            ocr_text=ocr_text,
                        )
                    )

                if image_descriptions:
                    doc_text += "\n\n" + "\n".join(image_descriptions)

            return ParsedFile(
                name=file.name,
                sections=[
                    SectionContent(
                        number=1,
                        text=doc_text,
                        md=doc_text,
                        images=final_images,
                    )
                ],
            )
            
    def _convert_to_docx(self, file: RawFile) -> RawFile:
        """Convert Word files (.doc) to .docx format and return as RawFile.

        Args:
            file: RawFile instance containing the input file data.

        Returns:
            RawFile instance containing converted content.

        Raises:
            RuntimeError: If conversion fails or LibreOffice not installed.
        """

        def _is_libreoffice_installed() -> bool:
            try:
                subprocess.run(
                    ["libreoffice", "--version"],
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                    check=True,
                )
                return True
            except (subprocess.CalledProcessError, FileNotFoundError):
                return False

        if not _is_libreoffice_installed():
            raise RuntimeError("LibreOffice not found in system PATH")

        with tempfile.TemporaryDirectory() as temp_dir:
            # Write input file to temporary directory
            input_path = os.path.join(temp_dir, file.name)
            with open(input_path, "wb") as f:
                f.write(file.contents)

            # Run LibreOffice conversion
            try:
                subprocess.run(
                    [
                        "libreoffice",
                        "--headless",
                        "--convert-to",
                        "docx",
                        "--outdir",
                        temp_dir,
                        input_path,
                    ],
                    check=True,
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.PIPE,
                    timeout=60,
                )
            except subprocess.CalledProcessError as e:
                error_msg = e.stderr.decode().strip() if e.stderr else "Unknown error"
                raise RuntimeError(f"Conversion failed: {error_msg}") from e
            except subprocess.TimeoutExpired:
                raise RuntimeError("Conversion timed out after 60 seconds")

            # Determine output file path
            output_filename = Path(file.name).stem + ".docx"
            output_path = os.path.join(temp_dir, output_filename)

            if not os.path.exists(output_path):
                available_files = os.listdir(temp_dir)
                raise RuntimeError(
                    f"Converted file not found at {output_path}. Found files: {available_files}"
                )

            # Read converted file and return as RawFile
            return RawFile.from_file_path(output_path)


@_parses("ppt", "pptx", "pptm")
class PptxFileParser(FileParser, frozen=True, tag="pptx"):
    """
    Parser for PowerPoint files (.pptx, .ppt, .pptm).

    This parser extracts text from shapes and images from slides in PowerPoint presentations.
    For higher parsing strategies (`ParsingStrategy.HIGH`), it can also utilize a visual description agent
    to generate textual descriptions of images found on the slides.

    **Supported Formats:**

    *   PowerPoint Presentation (.pptx)
    *   PowerPoint Presentation (.ppt) - Converted to .pptx format before parsing.
    *   PowerPoint Macro-Enabled Presentation (.pptm) - Converted to .pptx format before parsing.

    **Key Features:**

    *   **Text Extraction:** Extracts text content from shapes on each slide.
    *   **Image Extraction:** Extracts images embedded within the slides.
    *   **Visual Description (Optional):** When `ParsingStrategy.HIGH` is used and a `visual_description_agent` is provided, it generates AI-powered descriptions for extracted slide images.
    *   **Multi-Slide Support:** Processes multi-slide presentations, creating a `SectionContent` for each slide.
    *   **.ppt and .pptm Conversion:** Automatically converts older .ppt and macro-enabled .pptm formats to .pptx using LibreOffice before parsing.

    **Dependencies:**

    *   Requires the `pptx` library (python-pptx) to be installed (automatically handled if you installed `intellibricks[files]`).
    *   Requires LibreOffice to be installed and in the system PATH for .ppt and .pptm conversion.

    **Methods:**

    *   `parse_async(file: RawFile) -> ParsedFile`:
        Asynchronously extracts content from a PowerPoint `RawFile`.

        **Parameters:**

        *   `file` (RawFile): The `RawFile` object representing the PowerPoint file.

        **Returns:**

        *   `ParsedFile`: A `ParsedFile` object with sections, each representing a slide from the presentation, containing extracted text and images.

        **Raises:**

        *   `RuntimeError`: If LibreOffice is not installed or if .ppt/.pptm conversion fails.

        **Example:**

        ```python
        from architecture.data.files import RawFile
        from intelliparse.parsers PptxFileParser, ParsedFile, ParsingStrategy
        import asyncio

        async def main():
            # Assume you have pptx_content in bytes
            pptx_content = b"..." # Your PPTX file content in bytes
            pptx_file = RawFile.from_bytes(pptx_content, "pptx")

            parser = PptxFileParser(strategy=ParsingStrategy.HIGH) # Or DEFAULT, MEDIUM, FAST
            parsed_pptx: ParsedFile = await parser.parse_async(pptx_file)

            print(f"Parsed PPTX file name: {parsed_pptx.name}")
            for section in parsed_pptx.sections:
                print(f"  Section {section.number} (Slide): Text: {section.text[:100]}...") # First 100 chars of text
                print(f"  Section {section.number} (Slide): Images: {len(section.images)}") # Number of images on slide

        if __name__ == "__main__":
            asyncio.run(main())
        ```

    *   `_convert_to_pptx(file: RawFile) -> RawFile`:
        Converts PowerPoint files (.ppt/.pptm) to .pptx format using LibreOffice. (Internal use)

    **Usage Notes:**

    *   For `ParsingStrategy.HIGH`, ensure you have configured a `visual_description_agent` in the `FileParser` or `PptxFileParser` instance.
    *   For parsing .ppt and .pptm files, LibreOffice must be installed and accessible in the system's PATH.
    *   Each slide in the PowerPoint presentation is represented as a separate section in the output `ParsedFile`.
    """

    @ensure_module_installed("pptx", "intellibricks[files]")
    @override
    async def parse_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        import hashlib

        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.presentation import Presentation as PptxPresentation
        from pptx.shapes.autoshape import Shape
        from pptx.shapes.picture import Picture

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            if file.extension in {"ppt", "pptm"}:
                converted_pptx_file = self._convert_to_pptx(file)
                converted_pptx_file.save_to_file(file_path)
            else:
                file.save_to_file(file_path)

            prs: PptxPresentation = Presentation(file_path)
            sections: list[SectionContent] = []
            processed_images: dict[str, tuple[str, str]] = {}

            for slide_index, slide in enumerate(prs.slides, start=1):
                slide_texts: list[str] = []
                slide_images: list[tuple[str, bytes, str]] = []  # (name, data, hash)

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        shape_with_text = cast(Shape, shape)
                        text_str: str = shape_with_text.text
                        slide_texts.append(text_str)

                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        picture_shape = cast(Picture, shape)
                        image_blob: bytes = picture_shape.image.blob
                        image_hash = hashlib.sha256(image_blob).hexdigest()
                        image_name: str = (
                            shape.name or f"slide_{slide_index}_img_{image_hash[:8]}"
                        )
                        slide_images.append((image_name, image_blob, image_hash))

                combined_text: str = "\n".join(slide_texts)
                final_images: list[Image] = []
                image_descriptions: list[str] = []

                if self.visual_description_agent and self.strategy == "high":
                    for img_idx, (image_name, image_blob, image_hash) in enumerate(
                        slide_images, start=1
                    ):
                        is_cached = image_hash in processed_images
                        if is_cached:
                            cached_md, cached_ocr = processed_images[image_hash]
                            image_descriptions.append(
                                f"Slide {slide_index} - Image {img_idx}: {cached_md}"
                            )
                            final_images.append(
                                Image(
                                    name=image_name,
                                    contents=image_blob,
                                    ocr_text=cached_ocr,
                                )
                            )
                            continue

                        agent_input = ImageFilePart(
                            mime_type=bytes_to_mime(image_blob),
                            data=image_blob,
                        )
                        agent_response = await self.visual_description_agent.run_async(
                            agent_input
                        )
                        image_md: str = agent_response.parsed.final_answer.md
                        image_ocr = agent_response.parsed.final_answer.ocr_text

                        processed_images[image_hash] = (image_md, image_ocr)
                        image_descriptions.append(
                            f"Slide {slide_index} - Image {img_idx}: {image_md}"
                        )
                        final_images.append(
                            Image(
                                name=image_name, contents=image_blob, ocr_text=image_ocr
                            )
                        )

                    if image_descriptions:
                        combined_text += "\n\n" + "\n".join(image_descriptions)

                section_content = SectionContent(
                    number=slide_index,
                    text=combined_text,
                    md=combined_text,
                    images=final_images,
                )
                sections.append(section_content)

            return ParsedFile(
                name=file.name,
                sections=sections,
            )

    def _convert_to_pptx(self, file: RawFile) -> RawFile:
        """Convert PowerPoint files (.ppt/.pptm) to .pptx format and return as RawFile.

        Args:
            file: RawFile instance containing the input file data.

        Returns:
            RawFile instance containing converted content.

        Raises:
            RuntimeError: If conversion fails or LibreOffice not installed.
        """

        def _is_libreoffice_installed() -> bool:
            try:
                subprocess.run(
                    ["libreoffice", "--version"],
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                    check=True,
                )
                return True
            except (subprocess.CalledProcessError, FileNotFoundError):
                return False

        if not _is_libreoffice_installed():
            raise RuntimeError("LibreOffice not found in system PATH")

        with tempfile.TemporaryDirectory() as temp_dir:
            # Write input file to temporary directory
            input_path = os.path.join(temp_dir, file.name)
            with open(input_path, "wb") as f:
                f.write(file.contents)

            # Run LibreOffice conversion
            try:
                subprocess.run(
                    [
                        "libreoffice",
                        "--headless",
                        "--convert-to",
                        "pptx",
                        "--outdir",
                        temp_dir,
                        input_path,
                    ],
                    check=True,
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.PIPE,
                    timeout=60,
                )
            except subprocess.CalledProcessError as e:
                error_msg = e.stderr.decode().strip() if e.stderr else "Unknown error"
                raise RuntimeError(f"Conversion failed: {error_msg}") from e
            except subprocess.TimeoutExpired:
                raise RuntimeError("Conversion timed out after 60 seconds")

            # Determine output file path
            output_filename = Path(file.name).stem + ".pptx"
            output_path = os.path.join(temp_dir, output_filename)

            if not os.path.exists(output_path):
                available_files = os.listdir(temp_dir)
                raise RuntimeError(
                    f"Converted file not found at {output_path}. Found files: {available_files}"
                )

            # Read converted file and return as RawFile
            return RawFile.from_file_path(output_path)


@_parses("xls", "xlsx")
class ExcelFileParser(FileParser, frozen=True, tag="excel"):
    """
    Parser for Excel files (.xlsx, .xls).

    This parser extracts data and images from Excel spreadsheets. It extracts cell values as structured data
    and also as CSV formatted text. For higher parsing strategies (`ParsingStrategy.HIGH`),
    it can also utilize a visual description agent to generate textual descriptions of images found within the worksheets.

    **Supported Formats:**

    *   Excel Workbook (.xlsx)
    *   Excel Workbook (.xls)

    **Key Features:**

    *   **Structured Data Extraction:** Extracts cell values from each worksheet and provides them as lists of rows and a CSV string.
    *   **Image Extraction:** Extracts images embedded within the worksheets.
    *   **Visual Description (Optional):** When `ParsingStrategy.HIGH` is used and a `visual_description_agent` is provided, it generates AI-powered descriptions for extracted worksheet images.
    *   **Multi-Sheet Support:** Processes multi-sheet Excel workbooks, creating a `SectionContent` for each worksheet.
    *   **Table Page Items:** Creates `TablePageItem` objects within each section to represent the structured table data from each worksheet.

    **Dependencies:**

    *   Requires the `openpyxl` library to be installed (automatically handled if you installed `intellibricks[files]`).

    **Methods:**

    *   `parse_async(file: RawFile) -> ParsedFile`:
        Asynchronously extracts content from an Excel `RawFile`.

        **Parameters:**

        *   `file` (RawFile): The `RawFile` object representing the Excel file.

        **Returns:**

        *   `ParsedFile`: A `ParsedFile` object with sections, each representing a worksheet, containing extracted text (CSV), structured table data, and images.

        **Example:**

        ```python
        from architecture.data.files import RawFile
        from intelliparse.parsers ExcelFileParser, ParsedFile, ParsingStrategy
        import asyncio

        async def main():
            # Assume you have xlsx_content in bytes
            xlsx_content = b"..." # Your XLSX file content in bytes
            xlsx_file = RawFile.from_bytes(xlsx_content, "xlsx")

            parser = ExcelFileParser(strategy=ParsingStrategy.HIGH) # Or DEFAULT, MEDIUM, FAST
            parsed_excel: ParsedFile = await parser.parse_async(xlsx_file)

            print(f"Parsed Excel file name: {parsed_excel.name}")
            for section in parsed_excel.sections:
                print(f"  Section {section.number} (Worksheet): Text (CSV):\n{section.items[0].csv[:100]}...") # First 100 chars of CSV
                print(f"  Section {section.number} (Worksheet): Images: {len(section.images)}") # Number of images on worksheet

        if __name__ == "__main__":
            asyncio.run(main())
        ```

    **Usage Notes:**

    *   For `ParsingStrategy.HIGH`, ensure you have configured a `visual_description_agent` in the `FileParser` or `ExcelFileParser` instance.
    *   Each worksheet in the Excel workbook is represented as a separate section in the output `ParsedFile`.
    *   Structured table data (rows, CSV) is available within the `items` attribute of each section, specifically as a `TablePageItem`.
    """

    @ensure_module_installed("openpyxl", "intellibricks[files]")
    @override
    async def parse_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        import csv
        import io

        from openpyxl import Workbook, load_workbook

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            file.save_to_file(file_path)

            wb: Workbook = load_workbook(file_path, data_only=True)
            sections: list[SectionContent] = []

            for sheet_index, sheet in enumerate(wb.worksheets, start=1):
                # Gather structured data
                rows: list[list[str]] = []
                row_texts: list[str] = []
                for row in sheet.iter_rows(values_only=True):
                    # Process cell values
                    cell_values = [
                        str(cell) if cell is not None else "" for cell in row
                    ]
                    rows.append(cell_values)
                    row_texts.append("\t".join(cell_values))

                combined_text = "\n".join(row_texts)

                # Generate CSV content
                csv_buffer = io.StringIO()
                csv_writer = csv.writer(csv_buffer)
                csv_writer.writerows(rows)
                csv_str = csv_buffer.getvalue().strip()

                # Process images
                sheet_images: list[tuple[str, bytes]] = []
                if hasattr(sheet, "_images"):
                    image_list = getattr(sheet, "_images", [])
                    for img_idx, img in enumerate(image_list, start=1):
                        img_data = getattr(img, "_data", None)
                        if img_data is not None:
                            image_name = f"{sheet.title}_img_{img_idx}.png"
                            sheet_images.append((image_name, img_data))

                final_images: list[Image] = []
                # Generate image descriptions if needed
                if self.visual_description_agent and self.strategy == "high":
                    image_descriptions: list[str] = []
                    for img_idx, image_obj in enumerate(sheet_images, start=1):
                        agent_input = ImageFilePart(
                            mime_type=bytes_to_mime(image_obj[1]),
                            data=image_obj[1],
                        )
                        agent_response = await self.visual_description_agent.run_async(
                            agent_input
                        )
                        image_md = agent_response.parsed.final_answer.md
                        image_descriptions.append(
                            f"Worksheet {sheet.title} - Image {img_idx}: {image_md}"
                        )
                        final_images.append(
                            Image(
                                name=image_obj[0],
                                contents=image_obj[1],
                                ocr_text=agent_response.parsed.final_answer.ocr_text,
                            )
                        )

                    if image_descriptions:
                        combined_text += "\n\n" + "\n".join(image_descriptions)

                # Create table page item
                table_item = TablePageItem(
                    md=combined_text, rows=rows, csv=csv_str, is_perfect_table=True
                )

                section_content = SectionContent(
                    number=sheet_index,
                    text=combined_text,
                    md=combined_text,
                    images=final_images,
                    items=[table_item],
                )
                sections.append(section_content)

            return ParsedFile(
                name=file.name,
                sections=sections,
            )


@_parses("txt", "alg")
class TxtFileParser(FileParser, frozen=True, tag="txt"):
    """
    Parser for plain text files (.txt).

    This parser extracts the entire content of a text file as plain text. It creates a single section
    in the `ParsedFile` representing the entire file content.

    **Key Features:**

    *   **Plain Text Extraction:** Extracts the text content directly from .txt files.
    *   **Single Section Output:**  Outputs the entire text content as a single section in the `ParsedFile`.
    *   **Simple and Efficient:**  Provides a straightforward and efficient way to parse plain text files.

    **Methods:**

    *   `parse_async(file: RawFile) -> ParsedFile`:
        Asynchronously extracts content from a TXT `RawFile`.

        **Parameters:**

        *   `file` (RawFile): The `RawFile` object representing the text file.

        **Returns:**

        *   `ParsedFile`: A `ParsedFile` object with a single section containing the plain text content of the file.

        **Example:**

        ```python
        from architecture.data.files import RawFile
        from intelliparse.parsers TxtFileParser, ParsedFile, ParsingStrategy
        import asyncio

        async def main():
            # Assume you have txt_content in bytes
            txt_content = b"This is a sample text file."
            txt_file = RawFile.from_bytes(txt_content, "txt")

            parser = TxtFileParser(strategy=ParsingStrategy.DEFAULT)
            parsed_txt: ParsedFile = await parser.parse_async(txt_file)

            section = parsed_txt.sections[0]
            print(f"Parsed TXT file name: {parsed_txt.name}")
            print(f"Section 1 - Text Content:\n{section.text}")

        if __name__ == "__main__":
            asyncio.run(main())
        ```

    **Usage Notes:**

    *   This parser is ideal for simple text-based files where no complex structure or image extraction is needed.
    *   The output `ParsedFile` will always contain a single section representing the whole file's text content.
    """

    @override
    async def parse_async(self, file: RawFile) -> ParsedFile:
        text_content = file.contents.decode("utf-8", errors="replace")

        page_content = SectionContent(
            number=1,
            text=text_content,
            md=text_content,
        )

        return ParsedFile(
            name=file.name,
            sections=[page_content],
        )


@_parses("png", "jpeg", "tiff", "bmp", "jpg", "jp2")
class StaticImageFileParser(FileParser, frozen=True, tag="static_image"):
    """
    Parser for static image files (PNG, JPEG, TIFF, BMP, JPG).

    This parser extracts content from static image files. For higher parsing strategies (`ParsingStrategy.HIGH`)
    and if a `visual_description_agent` is configured, it can generate an AI-powered textual description of the image.
    For TIFF images, it performs an in-memory conversion to PNG for broader compatibility.

    **Supported Formats:**

    *   PNG (.png)
    *   JPEG (.jpeg, .jpg)
    *   TIFF (.tiff, .tif) - Converted to PNG internally.
    *   BMP (.bmp)

    **Key Features:**

    *   **Image Format Support:** Handles common static image formats.
    *   **TIFF to PNG Conversion:** Automatically converts TIFF images to PNG format in memory for consistent processing.
    *   **Visual Description (Optional):** When `ParsingStrategy.HIGH` is used and a `visual_description_agent` is provided, it generates AI-powered descriptions for the image.
    *   **Single Section Output:**  Outputs the image (and optional description) as a single section in the `ParsedFile`.

    **Dependencies:**

    *   Requires the `Pillow` (PIL) library to be installed (automatically handled if you installed `intellibricks[files]`).

    **Methods:**

    *   `parse_async(file: RawFile) -> ParsedFile`:
        Asynchronously extracts content from a static image `RawFile`.

        **Parameters:**

        *   `file` (RawFile): The `RawFile` object representing the static image file.

        **Returns:**

        *   `ParsedFile`: A `ParsedFile` object with a single section containing the image and optionally a textual description.

        **Example:**

        ```python
        from architecture.data.files import RawFile
        from intelliparse.parsers StaticImageFileParser, ParsedFile, ParsingStrategy
        import asyncio

        async def main():
            # Assume you have png_content in bytes
            png_content = b"..." # Your PNG file content in bytes
            png_file = RawFile.from_bytes(png_content, "png")

            parser = StaticImageFileParser(strategy=ParsingStrategy.HIGH) # Or DEFAULT, MEDIUM, FAST
            parsed_image: ParsedFile = await parser.parse_async(png_file)

            section = parsed_image.sections[0]
            print(f"Parsed Image file name: {parsed_image.name}")
            print(f"Section 1 - Images: {len(section.images)}") # Should be 1
            if section.text:
                print(f"Section 1 - Description:\n{section.text[:100]}...") # First 100 chars of description if strategy is HIGH

        if __name__ == "__main__":
            asyncio.run(main())
        ```

    **Usage Notes:**

    *   For `ParsingStrategy.HIGH`, ensure you have configured a `visual_description_agent` in the `FileParser` or `StaticImageFileParser` instance.
    *   TIFF images are converted to PNG in memory, so the `Image` object in the `ParsedFile` will always contain PNG data for TIFF inputs.
    *   The output `ParsedFile` contains a single section representing the image and its optional description.
    """

    @override
    async def parse_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        from PIL import Image as PILImage

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            file.save_to_file(file_path)

            # Determine the extension
            extension = file.extension

            # Convert to PNG if TIFF
            if extension in {"tiff", "tif"}:
                # Use Pillow to open, then convert to PNG in memory
                with io.BytesIO(file.contents) as input_buffer:
                    with PILImage.open(input_buffer) as pil_img:
                        # Convert to RGBA or RGB if needed
                        if pil_img.mode not in ("RGB", "RGBA"):
                            pil_img = pil_img.convert("RGBA")

                        # Save as PNG into a new buffer
                        output_buffer = io.BytesIO()
                        pil_img.save(output_buffer, format="PNG")
                        converted_bytes = output_buffer.getvalue()

                # Use the converted PNG bytes
                image_bytes = converted_bytes
                current_mime_type = bytes_to_mime(image_bytes)
            else:
                # No conversion needed
                image_bytes = file.contents

                # For demonstration, pick your MIME by extension
                if extension in {"png", "bmp"}:
                    current_mime_type = "image/" + extension
                elif extension in {"jpg", "jpeg"}:
                    current_mime_type = "image/jpeg"
                else:
                    # Fallback to PNG or raise an error if you want
                    current_mime_type = "image/png"

            # Create an Image object

            image_ocr: str | None = None
            # Generate a description if we have an agent + HIGH strategy
            text_content = ""
            if self.visual_description_agent and self.strategy == "high":
                agent_input = ImageFilePart(
                    mime_type=current_mime_type,
                    data=image_bytes,
                )
                agent_response = await self.visual_description_agent.run_async(
                    agent_input
                )
                description_md = agent_response.parsed.final_answer.md
                image_ocr = agent_response.parsed.final_answer.ocr_text
                text_content = description_md

            image_obj = Image(name=file.name, contents=image_bytes, ocr_text=image_ocr)
            # We treat it as a single "page" with one image
            page_content = SectionContent(
                number=1,
                text=text_content,
                md=text_content,
                images=[image_obj],
            )

            return ParsedFile(
                name=file.name,
                sections=[page_content],
            )


@_parses("gif")
class AnimatedImageFileParser(FileParser, frozen=True, tag="animated_image"):
    """
    Parser for animated GIF files (.gif).

    This parser handles animated GIF files by extracting representative frames from the animation.
    It selects up to 3 frames from the GIF, aiming to capture different points in the animation sequence.
    For higher parsing strategies (`ParsingStrategy.HIGH`) and if a `visual_description_agent` is configured,
    it can generate AI-powered textual descriptions for each selected frame.

    **Key Features:**

    *   **Animated GIF Support:** Specifically designed to parse .gif files.
    *   **Frame Selection:** Selects up to 3 representative frames from the animation.
    *   **Visual Description (Optional):** When `ParsingStrategy.HIGH` is used and a `visual_description_agent` is provided, it generates AI-powered descriptions for each selected frame.
    *   **Multi-Section Output:**  Outputs a `ParsedFile` with up to 3 sections, each representing a selected frame from the GIF animation.

    **Dependencies:**

    *   Requires the `Pillow` (PIL) library to be installed (automatically handled if you installed `intellibricks[files]`).

    **Methods:**

    *   `parse_async(file: RawFile) -> ParsedFile`:
        Asynchronously extracts content from an animated GIF `RawFile`.

        **Parameters:**

        *   `file` (RawFile): The `RawFile` object representing the animated GIF file.

        **Returns:**

        *   `ParsedFile`: A `ParsedFile` object with sections, each representing a selected frame from the GIF animation.

        **Example:**

        ```python
        from architecture.data.files import RawFile
        from intelliparse.parsers AnimatedImageFileParser, ParsedFile, ParsingStrategy
        import asyncio

        async def main():
            # Assume you have gif_content in bytes
            gif_content = b"..." # Your GIF file content in bytes
            gif_file = RawFile.from_bytes(gif_content, "gif")

            parser = AnimatedImageFileParser(strategy=ParsingStrategy.HIGH) # Or DEFAULT, MEDIUM, FAST
            parsed_gif: ParsedFile = await parser.parse_async(gif_file)

            print(f"Parsed GIF file name: {parsed_gif.name}")
            for section in parsed_gif.sections:
                print(f"  Section {section.number} (Frame): Images: {len(section.images)}") # Should be 1 per section
                if section.text:
                    print(f"  Section {section.number} (Frame): Description:\n{section.text[:100]}...") # First 100 chars of description if strategy is HIGH

        if __name__ == "__main__":
            asyncio.run(main())
        ```

    **Usage Notes:**

    *   For `ParsingStrategy.HIGH`, ensure you have configured a `visual_description_agent` in the `FileParser` or `AnimatedImageFileParser` instance.
    *   The number of sections in the output `ParsedFile` will be between 0 and 3, depending on the number of frames in the GIF and the frame selection process.
    *   Each section represents a single frame from the animation, converted to PNG format.
    """

    @override
    async def parse_async(
        self,
        file: RawFile,
    ) -> ParsedFile:
        from PIL import Image as PILImage

        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = f"{temp_dir}/{file.name}"
            file.save_to_file(file_path)

            # Safety check: only proceed if it's a .gif
            # or you can attempt detection based on file headers
            extension = file.extension
            if extension not in {"gif"}:
                raise ValueError("AnimatedImageFileParser only supports .gif files.")

            # --- 1. Load all frames from the GIF ---
            frames: list[PILImage.Image] = []
            with PILImage.open(file_path) as gif_img:
                try:
                    while True:
                        frames.append(gif_img.copy())
                        gif_img.seek(gif_img.tell() + 1)
                except EOFError:
                    pass  # we've reached the end of the animation

            num_frames = len(frames)
            if num_frames == 0:
                # No frames => no content
                return ParsedFile(name=file.name, sections=[])

            # --- 2. Pick up to 3 frames, splitting the GIF into 3 segments ---
            # If there are fewer than 3 frames, just use them all.
            # If more than 3, pick three frames spaced across the animation.

            if num_frames <= 3:
                selected_frames = frames
            else:
                # Example approach: pick near 1/3, 2/3, end
                idx1 = max(0, (num_frames // 3) - 1)
                idx2 = max(0, (2 * num_frames // 3) - 1)
                idx3 = num_frames - 1
                # Ensure distinct indexes
                unique_indexes = sorted(set([idx1, idx2, idx3]))
                selected_frames = [frames[i] for i in unique_indexes]

            # --- 3. Convert each selected frame to PNG and (optionally) describe it ---
            pages: list[SectionContent] = []
            for i, frame in enumerate(selected_frames, start=1):
                # Convert frame to PNG in-memory
                png_buffer = io.BytesIO()
                # Convert to RGBA if needed
                if frame.mode not in ("RGB", "RGBA"):
                    frame = frame.convert("RGBA")
                frame.save(png_buffer, format="PNG")
                png_bytes = png_buffer.getvalue()

                frame_image_ocr: str | None = None
                # If strategy is HIGH, pass the frame to the agent
                text_description = ""
                if self.visual_description_agent and self.strategy == "high":
                    agent_input = ImageFilePart(
                        mime_type=bytes_to_mime(png_bytes),
                        data=png_bytes,
                    )
                    agent_response = await self.visual_description_agent.run_async(
                        agent_input
                    )
                    frame_image_ocr = agent_response.parsed.final_answer.ocr_text
                    text_description = agent_response.parsed.final_answer.md

                # Create an Image object
                frame_image = Image(
                    name=f"{file.name}-frame{i}.png",
                    contents=png_bytes,
                    ocr_text=frame_image_ocr,
                )
                # Each frame is its own "page" in the final doc
                page_content = SectionContent(
                    number=i,
                    text=text_description,
                    md=text_description,
                    images=[frame_image],
                )
                pages.append(page_content)

            # --- 4. Return the multi-page ParsedFile ---
            return ParsedFile(
                name=file.name,
                sections=pages,
            )


@_parses("flac", "mp3", "mpeg", "mpga", "m4a", "ogg", "wav", "webm")
class AudioFileParser(FileParser, frozen=True, tag="audio"):
    """
    Parser for audio files (FLAC, MP3, MPEG, MPGA, M4A, OGG, WAV, WEBM).

    This parser handles various audio file formats. It first converts the audio file to MP3 format using FFmpeg
    to ensure compatibility with audio transcription services. Then, it uses an `audio_description_agent` to
    transcribe the audio content into text.

    **Supported Formats:**

    *   FLAC (.flac)
    *   MP3 (.mp3)
    *   MPEG (.mpeg)
    *   MPGA (.mpga)
    *   M4A (.m4a)
    *   OGG (.ogg)
    *   WAV (.wav)
    *   WEBM (.webm)

    **Key Features:**

    *   **Audio Format Conversion:** Converts audio files to MP3 format using FFmpeg for consistent processing.
    *   **Audio Transcription:** Uses an `audio_description_agent` to transcribe the audio content to text.
    *   **Single Section Output:**  Outputs the transcribed text and Markdown representation as a single section in the `ParsedFile`.

    **Dependencies:**

    *   Requires FFmpeg to be installed and in the system PATH.
    *   Requires an `audio_description_agent` to be provided to the parser instance.

    **Methods:**

    *   `parse_async(file: RawFile) -> ParsedFile`:
        Asynchronously extracts content from an audio `RawFile`.

        **Parameters:**

        *   `file` (RawFile): The `RawFile` object representing the audio file.

        **Returns:**

        *   `ParsedFile`: A `ParsedFile` object with a single section containing the transcribed text and Markdown representation of the audio content.

        **Raises:**

        *   `ValueError`: If no `audio_description_agent` is provided.
        *   `RuntimeError`: If FFmpeg is not installed or if audio conversion fails.

        **Example:**

        ```python
        from architecture.data.files import RawFile
        from intelliparse.parsers AudioFileParser, ParsedFile, ParsingStrategy
        from intellibricks.agents import Agent
        from intellibricks.llms.mock_llm import MockLLM  # Example Mock LLM
        from intellibricks.llms.types import ChainOfThought, AudioDescription
        import asyncio

        async def main():
            # Assume you have mp3_content in bytes
            mp3_content = b"..." # Your MP3 file content in bytes
            mp3_file = RawFile.from_bytes(mp3_content, "mp3")

            # Mock audio description agent for example
            mock_llm = MockLLM[ChainOfThought[AudioDescription]](...) # Configure MockLLM as needed
            audio_agent: Agent[ChainOfThought[AudioDescription]] = Agent(llm=mock_llm)

            parser = AudioFileParser(
                strategy=ParsingStrategy.DEFAULT,
                audio_description_agent=audio_agent # Provide the audio agent
            )
            parsed_audio: ParsedFile = await parser.parse_async(mp3_file)

            section = parsed_audio.sections[0]
            print(f"Parsed Audio file name: {parsed_audio.name}")
            print(f"Section 1 - Transcription:\n{section.text[:100]}...") # First 100 chars of transcription

        if __name__ == "__main__":
            asyncio.run(main())
        ```

    *   `_check_ffmpeg_installed() -> None`:
        Checks if FFmpeg is installed and accessible in the system PATH. (Internal use)
    *   `_could_not_transcript() -> Never`:
        Raises a ValueError indicating audio transcription failure. (Internal use)

    **Usage Notes:**

    *   Requires FFmpeg to be installed and configured correctly in your environment.
    *   You must provide an `audio_description_agent` instance when creating an `AudioFileParser`.
    *   The output `ParsedFile` contains a single section with the audio transcription.
    """

    async def parse_async(self, file: RawFile) -> ParsedFile:
        if self.audio_description_agent is None:
            raise ValueError("No audio description agent provided.")

        file_contents: bytes = file.contents
        file_extension = file.extension

        if file_extension in {
            "flac",
            "mpeg",
            "mpga",
            "m4a",
            "ogg",
            "wav",
            "webm",
        }:
            import aiofiles.os as aios
            from aiofiles import open as aio_open

            self._check_ffmpeg_installed()

            # Generate unique temporary filenames
            input_temp = os.path.join(
                tempfile.gettempdir(),
                f"input_{os.urandom(8).hex()}.{file_extension}",
            )
            output_temp = os.path.join(
                tempfile.gettempdir(), f"output_{os.urandom(8).hex()}.mp3"
            )

            # Write input file asynchronously
            async with aio_open(input_temp, "wb") as f:
                await f.write(file_contents)

            # Build FFmpeg command
            command = [
                "ffmpeg",
                "-hide_banner",
                "-loglevel",
                "error",  # Suppress unnecessary logs
                "-y",  # Overwrite output file if exists
                "-i",
                input_temp,
                "-codec:a",
                "libmp3lame",
                "-q:a",
                "2",  # Quality preset (0-9, 0=best)
                output_temp,
            ]

            # Execute FFmpeg
            process = await asyncio.create_subprocess_exec(
                *command, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
            )

            _, stderr = await process.communicate()

            # Handle conversion errors
            if process.returncode != 0:
                await aios.remove(input_temp)
                if await aios.path.exists(output_temp):
                    await aios.remove(output_temp)
                raise RuntimeError(
                    f"Audio conversion failed: {stderr.decode().strip()}"
                )

            # Read converted file
            async with aio_open(output_temp, "rb") as f:
                file_contents = await f.read()

            # Cleanup temporary files
            await aios.remove(input_temp)
            await aios.remove(output_temp)

        transcription = self.audio_description_agent.run(
            AudioFilePart(data=file_contents, mime_type=bytes_to_mime(file_contents))
        )

        return ParsedFile(
            name=file.name,
            sections=[
                SectionContent(
                    number=1,
                    text=transcription.audio_transcription.text
                    if transcription.audio_transcription is not None
                    else self._could_not_transcript(),
                    md=transcription.parsed.final_answer.md,
                    images=[],
                )
            ],
        )

    def _could_not_transcript(self) -> Never:
        raise ValueError("Could not transcribe the audio")

    def _check_ffmpeg_installed(self) -> None:
        import subprocess

        try:
            result = subprocess.run(
                ["ffmpeg", "-version"],
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
            )
            exception_logger.exception("FFmpeg is not installed or not in PATH.")
            if result.returncode != 0:
                raise RuntimeError()
        except FileNotFoundError:
            exception_logger.exception("FFmpeg is not installed or not in PATH.")
            raise RuntimeError()


@_parses("mp4")
class VideoFileParser(FileParser, frozen=True, tag="video"):
    """
    Parser for video files (.mp4).

    This parser handles MP4 video files. It uses a `visual_description_agent` to generate a textual description
    of the video content.

    **Supported Formats:**

    *   MP4 (.mp4)

    **Key Features:**

    *   **Video Description:** Uses a `visual_description_agent` to generate a textual description of the video content.
    *   **Single Section Output:**  Outputs the generated video description as a single section in the `ParsedFile`.

    **Dependencies:**

    *   Requires a `visual_description_agent` to be provided to the parser instance.

    **Methods:**

    *   `parse_async(file: RawFile) -> ParsedFile`:
        Asynchronously extracts content from a video `RawFile`.

        **Parameters:**

        *   `file` (RawFile): The `RawFile` object representing the video file.

        **Returns:**

        *   `ParsedFile`: A `ParsedFile` object with a single section containing the textual description of the video content.

        **Raises:**

        *   `ValueError`: If no `visual_description_agent` is provided or if the file is not an MP4 file.

        **Example:**

        ```python
        from architecture.data.files import RawFile
        from intelliparse.parsers VideoFileParser, ParsedFile, ParsingStrategy
        from intellibricks.agents import Agent
        from intellibricks.llms.mock_llm import MockLLM  # Example Mock LLM
        from intellibricks.llms.types import ChainOfThought, VisualMediaDescription
        import asyncio

        async def main():
            # Assume you have mp4_content in bytes
            mp4_content = b"..." # Your MP4 file content in bytes
            mp4_file = RawFile.from_bytes(mp4_content, "video.mp4", "mp4)

            # Mock visual description agent for example
            mock_llm = MockLLM[ChainOfThought[VisualMediaDescription]](...) # Configure MockLLM as needed
            visual_agent: Agent[ChainOfThought[VisualMediaDescription]] = Agent(llm=mock_llm)

            parser = VideoFileParser(
                strategy=ParsingStrategy.DEFAULT,
                visual_description_agent=visual_agent # Provide the visual agent
            )
            parsed_video: ParsedFile = await parser.parse_async(mp4_file)

            section = parsed_video.sections[0]
            print(f"Parsed Video file name: {parsed_video.name}")
            print(f"Section 1 - Description:\n{section.text[:100]}...") # First 100 chars of video description

        if __name__ == "__main__":
            asyncio.run(main())
        ```

    **Usage Notes:**

    *   You must provide a `visual_description_agent` instance when creating a `VideoFileParser`.
    *   Currently, only MP4 files are supported.
    *   The output `ParsedFile` contains a single section with the AI-generated video description.
    """

    async def parse_async(self, file: RawFile) -> ParsedFile:
        if self.visual_description_agent is None:
            raise ValueError("No visual description agent provided.")

        extension = file.extension
        if extension != "mp4":
            raise ValueError("VideoFileParser only supports .mp4 files.")

        file_contents = file.contents
        visual_media_description = await self.visual_description_agent.run_async(
            VideoFilePart(data=file_contents, mime_type=bytes_to_mime(file_contents))
        )

        return ParsedFile(
            name=file.name,
            sections=[
                SectionContent(
                    number=1,
                    text=visual_media_description.parsed.final_answer.md,
                    md=visual_media_description.parsed.final_answer.md,
                    images=[],
                )
            ],
        )