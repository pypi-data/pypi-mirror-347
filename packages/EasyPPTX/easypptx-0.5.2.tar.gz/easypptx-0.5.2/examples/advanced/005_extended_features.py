"""
Extended features example for EasyPPTX.

This example demonstrates the extended features of EasyPPTX including:
- Percentage-based positioning
- Meiryo as default font
- Custom color definitions
- Auto-alignment of multiple objects
- Reference PowerPoint template support
"""

from pathlib import Path

from pptx.enum.shapes import MSO_SHAPE

from easypptx import Presentation, Text

# Create output directory
output_dir = Path("output")
output_dir.mkdir(exist_ok=True)


# Example 1: Percentage-based positioning
def create_percentage_example():
    """Create a presentation demonstrating percentage-based positioning."""
    print("Creating percentage-based positioning example...")

    # Create a new presentation
    pres = Presentation()

    # Add a title slide
    slide = pres.add_slide()
    text = Text(slide)
    text.add_title("Percentage-Based Positioning", color="blue")

    # Add text at percentage positions
    text.add_paragraph("This text is at 10%, 30%", x="10%", y="30%", width="30%")
    text.add_paragraph("This text is at 50%, 30%", x="50%", y="30%", width="30%")
    text.add_paragraph("This text is at 10%, 50%", x="10%", y="50%", width="30%")
    text.add_paragraph("This text is at 50%, 50%", x="50%", y="50%", width="30%")

    # Add shapes using percentage positioning
    slide.add_shape(shape_type=MSO_SHAPE.RECTANGLE, x="10%", y="70%", width="30%", height="10%", fill_color="red")

    slide.add_shape(shape_type=MSO_SHAPE.OVAL, x="50%", y="70%", width="30%", height="10%", fill_color="green")

    # Save the presentation
    pres.save(output_dir / "percentage_example.pptx")
    print(f"Saved to {output_dir / 'percentage_example.pptx'}")


# Example 2: Default fonts and colors
def create_styling_example():
    """Create a presentation demonstrating default fonts and colors."""
    print("Creating styling example...")

    # Create a new presentation
    pres = Presentation()

    # Add a title slide
    slide = pres.add_slide()
    text = Text(slide)
    text.add_title("Default Fonts and Colors", color="black")

    # Add text with different colors from the color dictionary
    text.add_paragraph("This text uses the default black color", y="30%", color="black", font_size=24)

    text.add_paragraph("This text uses the red color", y="40%", color="red", font_size=24)

    text.add_paragraph("This text uses the green color", y="50%", color="green", font_size=24)

    text.add_paragraph("This text uses the blue color", y="60%", color="blue", font_size=24)

    text.add_paragraph("This text uses the white color on blue background", y="70%", color="white", font_size=24)

    # Add a blue background for the white text
    slide.add_shape(shape_type=MSO_SHAPE.RECTANGLE, x="10%", y="70%", width="80%", height="10%", fill_color="blue")

    # Save the presentation
    pres.save(output_dir / "styling_example.pptx")
    print(f"Saved to {output_dir / 'styling_example.pptx'}")


# Example 3: Auto-alignment of multiple objects
def create_auto_alignment_example():
    """Create a presentation demonstrating auto-alignment of multiple objects."""
    print("Creating auto-alignment example...")

    # Create a new presentation
    pres = Presentation()

    # Add a title slide
    slide = pres.add_slide()
    text = Text(slide)
    text.add_title("Auto-Alignment of Multiple Objects", color="black")

    # Create multiple objects to be automatically aligned in a grid layout
    objects_data = [
        {"type": "text", "text": "Item 1", "font_size": 24, "color": "black"},
        {"type": "text", "text": "Item 2", "font_size": 24, "color": "red"},
        {"type": "text", "text": "Item 3", "font_size": 24, "color": "green"},
        {"type": "text", "text": "Item 4", "font_size": 24, "color": "blue"},
        {"type": "shape", "shape_type": MSO_SHAPE.RECTANGLE, "fill_color": "black"},
        {"type": "shape", "shape_type": MSO_SHAPE.OVAL, "fill_color": "red"},
        {"type": "shape", "shape_type": MSO_SHAPE.ROUNDED_RECTANGLE, "fill_color": "green"},
        {"type": "shape", "shape_type": MSO_SHAPE.ACTION_BUTTON_HOME, "fill_color": "blue"},
    ]

    # Add objects in a grid layout with 5% padding
    slide.add_multiple_objects(
        objects_data=objects_data, layout="grid", padding_percent=5.0, start_y="30%", height="60%"
    )

    # Add a second slide with horizontal layout
    slide = pres.add_slide()
    text = Text(slide)
    text.add_title("Horizontal Layout Example", color="black")

    # Create objects for horizontal layout
    horizontal_objects = [
        {"type": "text", "text": "Left", "color": "red"},
        {"type": "text", "text": "Center", "color": "green"},
        {"type": "text", "text": "Right", "color": "blue"},
    ]

    # Add objects in a horizontal layout
    slide.add_multiple_objects(objects_data=horizontal_objects, layout="horizontal", start_y="50%", height="20%")

    # Add a third slide with vertical layout
    slide = pres.add_slide()
    text = Text(slide)
    text.add_title("Vertical Layout Example", color="black")

    # Create objects for vertical layout
    vertical_objects = [
        {"type": "text", "text": "Top", "color": "red"},
        {"type": "text", "text": "Middle", "color": "green"},
        {"type": "text", "text": "Bottom", "color": "blue"},
    ]

    # Add objects in a vertical layout
    slide.add_multiple_objects(
        objects_data=vertical_objects, layout="vertical", start_x="40%", width="20%", start_y="30%", height="60%"
    )

    # Save the presentation
    pres.save(output_dir / "auto_alignment_example.pptx")
    print(f"Saved to {output_dir / 'auto_alignment_example.pptx'}")


# Execute examples
if __name__ == "__main__":
    print("Creating examples with extended features...\n")

    create_percentage_example()
    print()

    create_styling_example()
    print()

    create_auto_alignment_example()
    print()

    print("All examples created successfully!")
