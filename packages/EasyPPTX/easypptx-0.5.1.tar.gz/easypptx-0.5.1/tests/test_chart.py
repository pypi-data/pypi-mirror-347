"""Tests for the Chart class."""

from unittest.mock import MagicMock, patch

import pandas as pd
import pytest
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

from easypptx import Chart


class TestChart:
    """Test cases for the Chart class."""

    def setup_method(self):
        """Set up test environment before each test method."""
        self.slide = MagicMock()
        self.chart = Chart(self.slide)

        # Mock chart shape and chart objects
        self.mock_chart_shape = MagicMock()
        self.mock_chart = MagicMock()
        self.mock_chart_shape.chart = self.mock_chart

        # Set the return value for shapes.add_chart
        self.slide.pptx_slide.shapes.add_chart.return_value = self.mock_chart_shape

        # Sample data for tests
        self.sample_categories = ["Category A", "Category B", "Category C"]
        self.sample_values = [10, 20, 30]

    def test_chart_init(self):
        """Test Chart initialization."""
        assert self.chart.slide == self.slide

    def test_chart_types(self):
        """Test that chart types dictionary is set correctly."""
        assert "column" in Chart.CHART_TYPES
        assert "bar" in Chart.CHART_TYPES
        assert "line" in Chart.CHART_TYPES
        assert "pie" in Chart.CHART_TYPES
        assert "area" in Chart.CHART_TYPES
        assert "scatter" in Chart.CHART_TYPES

    def test_add_chart_basic(self):
        """Test adding a basic chart."""
        # Call add method with basic parameters
        result = self.chart.add(chart_type="column", categories=self.sample_categories, values=self.sample_values)

        # Verify shapes.add_chart was called with correct parameters
        self.slide.pptx_slide.shapes.add_chart.assert_called_once()
        call_args = self.slide.pptx_slide.shapes.add_chart.call_args[0]

        # Check chart type
        assert call_args[0] == XL_CHART_TYPE.COLUMN_CLUSTERED

        # Since we're mocking the slide, we can't test the exact conversion
        # Instead, just test that the parameters were passed in some form
        assert call_args[1] is not None  # x
        assert call_args[2] is not None  # y
        assert call_args[3] is not None  # width
        assert call_args[4] is not None  # height

        # Verify chart data was created correctly
        chart_data = call_args[5]
        assert isinstance(chart_data, CategoryChartData)

        # Verify the result is the mock chart
        assert result == self.mock_chart

    def test_add_chart_with_position_and_size(self):
        """Test adding a chart with custom position and size."""
        # Call add method with custom position and size
        result = self.chart.add(
            chart_type="bar",
            categories=self.sample_categories,
            values=self.sample_values,
            x=2.0,
            y=3.0,
            width=5.0,
            height=3.5,
        )

        # Verify shapes.add_chart was called with correct parameters
        self.slide.pptx_slide.shapes.add_chart.assert_called_once()
        call_args = self.slide.pptx_slide.shapes.add_chart.call_args[0]

        assert call_args[0] == XL_CHART_TYPE.BAR_CLUSTERED
        # Since we're mocking the slide, we can't test the exact conversion
        # Instead, just test that the parameters were passed in some form
        assert call_args[1] is not None  # x
        assert call_args[2] is not None  # y
        assert call_args[3] is not None  # width
        assert call_args[4] is not None  # height

        # Verify the result is the mock chart
        assert result == self.mock_chart

    def test_add_chart_with_title(self):
        """Test adding a chart with a title."""
        # Call add method with a title
        result = self.chart.add(
            chart_type="column", categories=self.sample_categories, values=self.sample_values, title="Test Chart Title"
        )

        # Verify chart title was set
        assert self.mock_chart.has_title is True
        assert self.mock_chart.chart_title.text_frame.text == "Test Chart Title"

        # Verify the result is the mock chart
        assert result == self.mock_chart

    def test_add_chart_without_title(self):
        """Test adding a chart without a title."""
        # Call add method without a title
        result = self.chart.add(
            chart_type="column", categories=self.sample_categories, values=self.sample_values, title=None
        )

        # Verify chart title was not set
        assert self.mock_chart.has_title is False

        # Verify the result is the mock chart
        assert result == self.mock_chart

    def test_add_chart_with_legend(self):
        """Test adding a chart with a legend."""
        # Call add method with has_legend=True
        result = self.chart.add(
            chart_type="pie", categories=self.sample_categories, values=self.sample_values, has_legend=True
        )

        # Verify legend was enabled
        assert self.mock_chart.has_legend is True

        # Verify the result is the mock chart
        assert result == self.mock_chart

    def test_add_chart_without_legend(self):
        """Test adding a chart without a legend."""
        # Call add method with has_legend=False
        result = self.chart.add(
            chart_type="pie", categories=self.sample_categories, values=self.sample_values, has_legend=False
        )

        # Verify legend was disabled
        assert self.mock_chart.has_legend is False

        # Verify the result is the mock chart
        assert result == self.mock_chart

    def test_add_chart_with_invalid_type(self):
        """Test adding a chart with an invalid chart type raises ValueError."""
        with pytest.raises(ValueError):
            self.chart.add(chart_type="invalid_type", categories=self.sample_categories, values=self.sample_values)

    def test_add_chart_with_mismatched_data(self):
        """Test adding a chart with mismatched categories and values raises ValueError."""
        with pytest.raises(ValueError):
            self.chart.add(
                chart_type="column",
                categories=["A", "B"],  # 2 categories
                values=[1, 2, 3],  # 3 values
            )

    def test_add_bar_chart(self):
        """Test adding a bar chart using the convenience method."""
        # Patch the add method
        with patch.object(self.chart, "add") as mock_add:
            # Call add_bar
            self.chart.add_bar(categories=self.sample_categories, values=self.sample_values, title="Bar Chart")

            # Verify add was called with correct parameters
            mock_add.assert_called_once()
            call_args = mock_add.call_args[1]

            assert call_args["chart_type"] == "bar"
            assert call_args["categories"] == self.sample_categories
            assert call_args["values"] == self.sample_values
            assert call_args["title"] == "Bar Chart"

    def test_add_column_chart(self):
        """Test adding a column chart using the convenience method."""
        # Patch the add method
        with patch.object(self.chart, "add") as mock_add:
            # Call add_column
            self.chart.add_column(categories=self.sample_categories, values=self.sample_values, x=2.0, y=3.0)

            # Verify add was called with correct parameters
            mock_add.assert_called_once()
            call_args = mock_add.call_args[1]

            assert call_args["chart_type"] == "column"
            assert call_args["categories"] == self.sample_categories
            assert call_args["values"] == self.sample_values
            assert call_args["x"] == 2.0
            assert call_args["y"] == 3.0

    def test_add_pie_chart(self):
        """Test adding a pie chart using the convenience method."""
        # Patch the add method
        with patch.object(self.chart, "add") as mock_add:
            # Call add_pie
            self.chart.add_pie(categories=self.sample_categories, values=self.sample_values, width=4.0, height=4.0)

            # Verify add was called with correct parameters
            mock_add.assert_called_once()
            call_args = mock_add.call_args[1]

            assert call_args["chart_type"] == "pie"
            assert call_args["categories"] == self.sample_categories
            assert call_args["values"] == self.sample_values
            assert call_args["width"] == 4.0
            assert call_args["height"] == 4.0

    def test_from_dataframe_basic(self):
        """Test creating a chart from a pandas DataFrame."""
        # Create a DataFrame
        df = pd.DataFrame({
            "Category": ["A", "B", "C"],
            "Value": [10, 20, 30],
        })

        # Patch the add method
        with patch.object(self.chart, "add") as mock_add:
            # Call from_dataframe
            self.chart.from_dataframe(df=df, chart_type="column", category_column="Category", value_column="Value")

            # Verify add was called with correct data
            mock_add.assert_called_once()
            call_args = mock_add.call_args[1]

            assert call_args["chart_type"] == "column"
            assert call_args["categories"] == ["A", "B", "C"]
            assert call_args["values"] == [10, 20, 30]

    def test_from_dataframe_with_parameters(self):
        """Test creating a chart from DataFrame with custom parameters."""
        # Create a DataFrame
        df = pd.DataFrame({
            "Category": ["A", "B", "C"],
            "Value": [10, 20, 30],
        })

        # Patch the add method
        with patch.object(self.chart, "add") as mock_add:
            # Call from_dataframe with custom parameters
            self.chart.from_dataframe(
                df=df,
                chart_type="pie",
                category_column="Category",
                value_column="Value",
                x=2.0,
                y=3.0,
                width=5.0,
                height=5.0,
                title="Pie Chart",
                has_legend=True,
                legend_position=2,
            )

            # Verify add was called with correct parameters
            mock_add.assert_called_once()
            call_args = mock_add.call_args[1]

            assert call_args["chart_type"] == "pie"
            assert call_args["x"] == 2.0
            assert call_args["y"] == 3.0
            assert call_args["width"] == 5.0
            assert call_args["height"] == 5.0
            assert call_args["title"] == "Pie Chart"
            assert call_args["has_legend"] is True
            assert call_args["legend_position"] == 2

    def test_from_dataframe_invalid_category_column(self):
        """Test from_dataframe with an invalid category column raises ValueError."""
        df = pd.DataFrame({
            "Category": ["A", "B", "C"],
            "Value": [10, 20, 30],
        })

        with pytest.raises(ValueError):
            self.chart.from_dataframe(
                df=df, chart_type="column", category_column="NonExistentColumn", value_column="Value"
            )

    def test_from_dataframe_invalid_value_column(self):
        """Test from_dataframe with an invalid value column raises ValueError."""
        df = pd.DataFrame({
            "Category": ["A", "B", "C"],
            "Value": [10, 20, 30],
        })

        with pytest.raises(ValueError):
            self.chart.from_dataframe(
                df=df, chart_type="column", category_column="Category", value_column="NonExistentColumn"
            )
