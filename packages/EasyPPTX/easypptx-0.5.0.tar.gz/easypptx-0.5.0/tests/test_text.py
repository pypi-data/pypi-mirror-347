"""Tests for the Text class."""

from unittest.mock import MagicMock

import pytest
from pptx.dml.color import ColorFormat
from pptx.text.text import Font, TextFrame, _Paragraph
from pptx.util import Pt

from easypptx import Text


@pytest.fixture
def mock_slide():
    """Create a mock slide for testing."""
    slide = MagicMock()
    # Mock the add_text method to return a shape with a text frame
    mock_shape = MagicMock()
    mock_paragraph = MagicMock(spec=_Paragraph)
    mock_font = MagicMock(spec=Font)
    mock_color = MagicMock(spec=ColorFormat)
    mock_font.color = mock_color
    mock_paragraph.font = mock_font

    mock_text_frame = MagicMock(spec=TextFrame)
    mock_text_frame.paragraphs = [mock_paragraph]
    mock_shape.text_frame = mock_text_frame

    slide.add_text.return_value = mock_shape
    return slide


def test_text_init():
    """Test Text initialization."""
    slide = MagicMock()
    text = Text(slide)
    assert text.slide == slide


def test_add_title():
    """Test adding a title to a slide."""
    slide = MagicMock()
    text = Text(slide)

    # Call add_title
    result = text.add_title("Test Title")

    # Verify slide.add_text was called with correct parameters
    slide.add_text.assert_called_once()
    call_args = slide.add_text.call_args[1]
    assert call_args["text"] == "Test Title"
    assert call_args["font_bold"] is True
    assert call_args["font_size"] == 44

    # Verify the result is the return value from slide.add_text
    assert result == slide.add_text.return_value


def test_add_paragraph():
    """Test adding a paragraph to a slide."""
    slide = MagicMock()
    text = Text(slide)

    # Call add_paragraph with default parameters
    result = text.add_paragraph("Test Paragraph")

    # Verify slide.add_text was called with correct parameters
    slide.add_text.assert_called_once()
    call_args = slide.add_text.call_args[1]
    assert call_args["text"] == "Test Paragraph"
    assert call_args["x"] == 1.0
    assert call_args["y"] == 2.0
    assert call_args["width"] == 8.0
    assert call_args["height"] == 1.0
    assert call_args["font_size"] == 18
    assert call_args["font_bold"] is False
    assert call_args["font_italic"] is False

    # Verify the result is the return value from slide.add_text
    assert result == slide.add_text.return_value


def test_add_paragraph_with_formatting():
    """Test adding a paragraph with formatting options."""
    slide = MagicMock()
    text = Text(slide)

    # Set up mock slide

    # Call add_paragraph with formatting
    result = text.add_paragraph(
        "Test Paragraph",
        x=2.0,
        y=3.0,
        width=5.0,
        height=2.0,
        font_size=24,
        font_bold=True,
        font_italic=True,
        color=(255, 0, 0),  # Red
    )

    # Verify slide.add_text was called with correct parameters
    slide.add_text.assert_called_once()
    call_args = slide.add_text.call_args[1]
    assert call_args["text"] == "Test Paragraph"
    assert call_args["x"] == 2.0
    assert call_args["y"] == 3.0
    assert call_args["width"] == 5.0
    assert call_args["height"] == 2.0
    assert call_args["font_size"] == 24
    assert call_args["font_bold"] is True
    assert call_args["font_italic"] is True

    # Verify the result is the return value from slide.add_text
    assert result == slide.add_text.return_value


def test_format_text_frame():
    """Test the static format_text_frame method."""
    mock_text_frame = MagicMock(spec=TextFrame)
    mock_paragraph1 = MagicMock(spec=_Paragraph)
    mock_paragraph2 = MagicMock(spec=_Paragraph)

    mock_font1 = MagicMock(spec=Font)
    mock_font2 = MagicMock(spec=Font)
    mock_color1 = MagicMock(spec=ColorFormat)
    mock_color2 = MagicMock(spec=ColorFormat)

    mock_font1.color = mock_color1
    mock_font2.color = mock_color2
    mock_paragraph1.font = mock_font1
    mock_paragraph2.font = mock_font2

    mock_text_frame.paragraphs = [mock_paragraph1, mock_paragraph2]

    # Call format_text_frame with all parameters
    Text.format_text_frame(
        mock_text_frame,
        font_size=32,
        font_bold=True,
        font_italic=True,
        color=(0, 255, 0),  # Green
    )

    # Check that formatting was applied to both paragraphs
    for mock_paragraph in mock_text_frame.paragraphs:
        assert mock_paragraph.font.size == Pt(32)
        assert mock_paragraph.font.bold is True
        assert mock_paragraph.font.italic is True


def test_format_text_frame_partial_properties():
    """Test format_text_frame with only some properties specified."""
    mock_text_frame = MagicMock(spec=TextFrame)
    mock_paragraph = MagicMock(spec=_Paragraph)
    mock_font = MagicMock(spec=Font)
    mock_paragraph.font = mock_font
    mock_text_frame.paragraphs = [mock_paragraph]

    # Call format_text_frame with only font_size
    Text.format_text_frame(mock_text_frame, font_size=24)

    # Check that only font_size was set
    assert mock_paragraph.font.size == Pt(24)
    assert not mock_paragraph.font.bold.called
    assert not mock_paragraph.font.italic.called
