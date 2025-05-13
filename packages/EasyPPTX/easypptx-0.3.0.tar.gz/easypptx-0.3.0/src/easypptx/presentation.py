"""Core presentation module for EasyPPTX."""

from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING, Any, ClassVar

import pandas as pd
from pptx import Presentation as PPTXPresentation
from pptx.chart.chart import Chart as PPTXChart
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.shapes.autoshape import Shape as PPTXShape
from pptx.util import Inches, Pt

from easypptx.image import Image
from easypptx.pyplot import Pyplot
from easypptx.slide import Slide
from easypptx.table import Table
from easypptx.template import Template, TemplateManager
from easypptx.text import Text

# Use TYPE_CHECKING to avoid circular imports at runtime
if TYPE_CHECKING:
    from easypptx.grid import Grid
else:
    # Import at top level for runtime but avoid circular imports
    from easypptx.grid import Grid


class Presentation:
    """Main presentation class for creating and manipulating PowerPoint presentations.

    This class provides a simplified interface for working with PowerPoint presentations,
    making it easy to create, modify, and save PPTX files.

    For standard aspect ratios (16:9 and 4:3), reference templates are automatically used
    when no specific template is provided, ensuring consistent and attractive presentations.

    Attributes:
        pptx_presentation: The underlying python-pptx Presentation object

    Examples:
        ```python
        # Create a new presentation with default 16:9 aspect ratio
        # (automatically uses the reference_16x9.pptx template)
        presentation = Presentation()

        # Create a presentation with 4:3 aspect ratio
        # (automatically uses the reference_4x3.pptx template)
        presentation = Presentation(aspect_ratio="4:3")

        # Create a presentation with custom dimensions
        # (doesn't use reference templates)
        presentation = Presentation(width_inches=13.33, height_inches=7.5)

        # Create a presentation with a custom template
        # (overrides the reference templates)
        presentation = Presentation(template_path="my_template.pptx")

        # Open an existing presentation
        presentation = Presentation.open("example.pptx")

        # Add a new slide
        slide = presentation.add_slide()

        # Save the presentation
        presentation.save("output.pptx")
        ```
    """

    # Standard aspect ratios in width:height format
    ASPECT_RATIOS: ClassVar[dict[str, tuple[float, float]]] = {
        "16:9": (13.33, 7.5),  # Widescreen (default)
        "4:3": (10, 7.5),  # Standard
        "16:10": (13.33, 8.33),  # Widescreen alternative
        "A4": (11.69, 8.27),  # A4 paper size
        "LETTER": (11, 8.5),  # US Letter paper size
    }

    # Default colors dictionary
    COLORS: ClassVar[dict[str, RGBColor]] = {
        "black": RGBColor(0x10, 0x10, 0x10),
        "darkgray": RGBColor(0x40, 0x40, 0x40),
        "gray": RGBColor(0x80, 0x80, 0x80),
        "lightgray": RGBColor(0xD0, 0xD0, 0xD0),
        "red": RGBColor(0xFF, 0x40, 0x40),
        "green": RGBColor(0x40, 0xFF, 0x40),
        "blue": RGBColor(0x40, 0x40, 0xFF),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "yellow": RGBColor(0xFF, 0xD7, 0x00),
        "cyan": RGBColor(0x00, 0xE5, 0xFF),
        "magenta": RGBColor(0xFF, 0x00, 0xFF),
        "orange": RGBColor(0xFF, 0xA5, 0x00),
    }

    # Text alignment dictionary
    ALIGN: ClassVar[dict[str, PP_ALIGN]] = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

    # Vertical alignment dictionary
    VERTICAL: ClassVar[dict[str, MSO_ANCHOR]] = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }

    # Default font settings
    DEFAULT_FONT = "Meiryo"

    def __init__(
        self,
        aspect_ratio: str | None = "16:9",
        width_inches: float | None = None,
        height_inches: float | None = None,
        template_path: str | None = None,
        reference_pptx: str | None = None,
        blank_layout_index: int | None = None,
        default_bg_color: str | tuple[int, int, int] | None = None,
    ) -> None:
        """Initialize a new empty presentation.

        Args:
            aspect_ratio: Predefined aspect ratio, one of "16:9" (default), "4:3", "16:10", "A4", "LETTER"
            width_inches: Custom width in inches (overrides aspect_ratio if specified)
            height_inches: Custom height in inches (overrides aspect_ratio if specified)
            template_path: Path to a reference PowerPoint template to use for styles (default: None)
            reference_pptx: Path to a custom reference PPTX file to use (default: None)
            blank_layout_index: Index of blank layout in the slide_layouts (default: None, auto-detected)
            default_bg_color: Default background color for slides as string name or RGB tuple (default: None)

        Raises:
            ValueError: If an invalid aspect ratio is specified
            FileNotFoundError: If the template file doesn't exist
        """
        self.default_bg_color = default_bg_color

        # Initialize the Template object and TemplateManager
        self.template = Template()
        self.template_manager = TemplateManager()

        # Track which reference PPTX file we've loaded
        self._loaded_reference = None

        if template_path:
            # Use an existing template
            try:
                self.pptx_presentation = PPTXPresentation(template_path)
                self._loaded_reference = str(template_path)
            except Exception as e:
                raise FileNotFoundError(f"Template file not found or invalid: {e}") from e
        elif reference_pptx:
            # Use a custom reference PPTX file
            reference_path = Path(reference_pptx)
            if not reference_path.exists():
                raise FileNotFoundError(f"Reference PPTX file not found: {reference_pptx}")

            try:
                self.pptx_presentation = PPTXPresentation(str(reference_path))
                self._loaded_reference = str(reference_path)
            except Exception as e:
                raise FileNotFoundError(f"Reference PPTX file not found or invalid: {e}") from e
        else:
            # Check if we should use a reference template based on aspect ratio
            reference_template = None

            # Only use reference templates for the specific aspect ratios we have templates for
            if aspect_ratio == "16:9" and width_inches is None and height_inches is None:
                # Use 16:9 reference template
                reference_template = Path(__file__).parent / "reference_16x9.pptx"
            elif aspect_ratio == "4:3" and width_inches is None and height_inches is None:
                # Use 4:3 reference template
                reference_template = Path(__file__).parent / "reference_4x3.pptx"

            if reference_template and reference_template.exists():
                # Use the appropriate reference template
                self.pptx_presentation = PPTXPresentation(str(reference_template))
            else:
                # Create a new presentation without a template
                self.pptx_presentation = PPTXPresentation()

            # Set slide dimensions based on inputs
            if width_inches is not None and height_inches is not None:
                # Use custom dimensions
                self._set_slide_dimensions(width_inches, height_inches)
            elif aspect_ratio is not None:
                # Use predefined aspect ratio
                if aspect_ratio not in self.ASPECT_RATIOS:
                    valid_ratios = ", ".join(self.ASPECT_RATIOS.keys())
                    raise ValueError(f"Invalid aspect ratio: {aspect_ratio}. Valid options are: {valid_ratios}")

                width, height = self.ASPECT_RATIOS[aspect_ratio]
                self._set_slide_dimensions(width, height)

        # Find and store the blank slide layout
        if blank_layout_index is not None:
            # Use the specified index
            if 0 <= blank_layout_index < len(self.pptx_presentation.slide_layouts):
                self.blank_layout = self.pptx_presentation.slide_layouts[blank_layout_index]
            else:
                # If index is out of range, fall back to a safe default
                self.blank_layout = self._find_blank_layout() or self.pptx_presentation.slide_layouts[0]
        else:
            # Auto-detect the blank layout (typically index 6, but can vary)
            self.blank_layout = self._find_blank_layout() or self.pptx_presentation.slide_layouts[6]

    def _find_blank_layout(self) -> Any:
        """Find the blank layout in the presentation.

        Attempts to find the layout with the fewest placeholders, which is typically the blank layout.

        Returns:
            The slide layout that appears to be blank, or None if no suitable layout is found
        """
        # First, check if there's a layout named "Blank" or similar
        for layout in self.pptx_presentation.slide_layouts:
            if hasattr(layout, "name") and "blank" in layout.name.lower():
                return layout

        # Find the layout with the fewest placeholders (likely to be blank)
        blank_layout = None
        min_placeholders = float("inf")

        for layout in self.pptx_presentation.slide_layouts:
            placeholder_count = len(layout.placeholders)
            if placeholder_count < min_placeholders:
                min_placeholders = placeholder_count
                blank_layout = layout

        return blank_layout if min_placeholders < 3 else None

    def _set_slide_dimensions(self, width_inches: float, height_inches: float) -> None:
        """Set the slide dimensions.

        Args:
            width_inches: Width in inches
            height_inches: Height in inches
        """
        self.pptx_presentation.slide_width = int(Inches(width_inches))
        self.pptx_presentation.slide_height = int(Inches(height_inches))

    @classmethod
    def open(cls, file_path: str | Path, blank_layout_index: int | None = None) -> Presentation:
        """Open an existing PowerPoint presentation.

        Args:
            file_path: Path to the PowerPoint file to open
            blank_layout_index: Index of blank layout in the slide_layouts (default: None, auto-detected)

        Returns:
            A new Presentation object with the loaded presentation

        Raises:
            FileNotFoundError: If the specified file doesn't exist
            ValueError: If the file is not a valid PowerPoint file
        """
        file_path_obj = Path(file_path)

        if not file_path_obj.exists():
            raise FileNotFoundError(f"Presentation file not found: {file_path}")

        try:
            pptx_presentation = PPTXPresentation(str(file_path_obj))
        except Exception as e:
            raise ValueError(f"Invalid PowerPoint file: {e}") from e

        presentation = cls(width_inches=None, height_inches=None, blank_layout_index=blank_layout_index)
        presentation.pptx_presentation = pptx_presentation
        presentation._loaded_reference = str(file_path_obj)

        # Find and store the blank slide layout
        if blank_layout_index is not None:
            # Use the specified index
            if 0 <= blank_layout_index < len(presentation.pptx_presentation.slide_layouts):
                presentation.blank_layout = presentation.pptx_presentation.slide_layouts[blank_layout_index]
            else:
                # If index is out of range, fall back to a safe default
                presentation.blank_layout = (
                    presentation._find_blank_layout() or presentation.pptx_presentation.slide_layouts[0]
                )
        else:
            # Auto-detect the blank layout (typically index 6, but can vary)
            presentation.blank_layout = (
                presentation._find_blank_layout() or presentation.pptx_presentation.slide_layouts[6]
            )

        return presentation

    def add_slide(self, layout_index: int | None = None, bg_color: str | tuple[int, int, int] | None = None) -> Slide:
        """Add a new slide to the presentation.

        Args:
            layout_index: Index of the slide layout to use (default: None uses blank layout)
            bg_color: Background color for this slide, overrides default (default: None)

        Returns:
            A new Slide object
        """
        # Use blank layout by default, or specified layout if provided
        slide_layout = self.blank_layout if layout_index is None else self.pptx_presentation.slide_layouts[layout_index]

        pptx_slide = self.pptx_presentation.slides.add_slide(slide_layout)
        slide = Slide(pptx_slide)

        # Apply background color if specified for this slide or as default
        color_to_use = bg_color if bg_color is not None else self.default_bg_color
        if color_to_use is not None:
            slide.set_background_color(color_to_use)

        return slide

    @property
    def slides(self) -> list[Slide]:
        """Get a list of all slides in the presentation.

        Returns:
            List of Slide objects
        """
        return [Slide(slide) for slide in self.pptx_presentation.slides]

    def save(self, file_path: str | Path) -> None:
        """Save the presentation to a file.

        Args:
            file_path: Path where the presentation should be saved
        """
        self.pptx_presentation.save(file_path)

    def add_slide_from_template(self, template_data: str | dict) -> Slide:
        """Add a slide using a template preset.

        Args:
            template_data: Template preset name or custom template dictionary

        Returns:
            A new Slide object configured according to the template
        """
        # Get the template data - either from a preset or use the provided dictionary
        template_name = None
        if isinstance(template_data, str):
            template_name = template_data
            try:
                # First try to get from TemplateManager (includes built-in and registered templates)
                preset = self.template_manager.get(template_data)
            except ValueError:
                # Fall back to the legacy approach if not found
                preset = self.template.get_preset(template_data)
        else:
            preset = template_data

        # Check if there's a reference PPTX file specified for this template
        reference_pptx = None
        blank_layout_index = None
        if template_name is not None:
            # First check in TemplateManager for loaded templates
            reference_pptx = self.template_manager.get_reference_pptx(template_name)
            blank_layout_index = self.template_manager.get_blank_layout_index(template_name)

            # If not found, check built-in presets
            if reference_pptx is None:
                reference_pptx = self.template.get_reference_pptx(template_name)
                blank_layout_index = self.template.get_blank_layout_index(template_name)

            # If a reference PPTX is specified and we haven't loaded it yet, create a new Presentation with it
            if reference_pptx is not None and self._loaded_reference != reference_pptx:
                # Save current properties we want to preserve
                current_width = self.pptx_presentation.slide_width
                current_height = self.pptx_presentation.slide_height

                # Load the reference PPTX
                try:
                    self.pptx_presentation = PPTXPresentation(reference_pptx)
                    self._loaded_reference = reference_pptx

                    # Restore dimensions
                    self.pptx_presentation.slide_width = current_width
                    self.pptx_presentation.slide_height = current_height

                    # Update blank layout
                    if blank_layout_index is not None:
                        if 0 <= blank_layout_index < len(self.pptx_presentation.slide_layouts):
                            self.blank_layout = self.pptx_presentation.slide_layouts[blank_layout_index]
                        else:
                            self.blank_layout = self._find_blank_layout() or self.pptx_presentation.slide_layouts[0]
                    else:
                        self.blank_layout = self._find_blank_layout() or self.pptx_presentation.slide_layouts[6]
                except Exception as e:
                    print(f"Warning: Failed to load reference PPTX '{reference_pptx}': {e}")

        # Get background color if specified
        bg_color = preset.get("bg_color", None)

        # Create a new slide using the blank layout
        slide = self.add_slide(bg_color=bg_color)

        # Add title if specified in template
        if "title" in preset:
            title_data = preset["title"]
            position = title_data.get("position", {"x": "5%", "y": "5%", "width": "90%", "height": "10%"})

            # Extract font information
            font_data = title_data.get("font", {})
            font_name = font_data.get("name", self.DEFAULT_FONT)
            font_size = font_data.get("size", 32)
            font_bold = font_data.get("bold", True)

            # Extract alignment information
            align = title_data.get("align", "center")
            vertical = title_data.get("vertical", "middle")
            color = title_data.get("color", "black")

            # Add the title text
            Text.add(
                slide=slide,
                text=title_data.get("text", "Title"),
                position=position,
                font_name=font_name,
                font_size=font_size,
                font_bold=font_bold,
                align=align,
                vertical_align=vertical,
                color=color,
            )

        # Add subtitle if specified in template
        if "subtitle" in preset:
            subtitle_data = preset["subtitle"]
            position = subtitle_data.get("position", {"x": "20%", "y": "60%", "width": "60%", "height": "20%"})

            # Extract font information
            font_data = subtitle_data.get("font", {})
            font_name = font_data.get("name", self.DEFAULT_FONT)
            font_size = font_data.get("size", 24)
            font_bold = font_data.get("bold", False)

            # Extract alignment information
            align = subtitle_data.get("align", "center")
            vertical = subtitle_data.get("vertical", "middle")
            color = subtitle_data.get("color", "black")

            # Add the subtitle text
            Text.add(
                slide=slide,
                text=subtitle_data.get("text", "Subtitle"),
                position=position,
                font_name=font_name,
                font_size=font_size,
                font_bold=font_bold,
                align=align,
                vertical_align=vertical,
                color=color,
            )

        # Add decorative bar if specified
        if "bar" in preset:
            bar_data = preset["bar"]
            position = bar_data.get("position", {"x": "0%", "y": "10%", "width": "100%", "height": "2%"})

            # Create rectangle shape for the bar
            shape = slide.add_shape(
                shape_type=MSO_SHAPE.RECTANGLE,
                x=position.get("x", "0%"),
                y=position.get("y", "10%"),
                width=position.get("width", "100%"),
                height=position.get("height", "2%"),
            )

            # Apply gradient if specified
            if "gradient" in bar_data:
                gradient = bar_data["gradient"]
                start_color = gradient.get("start_color")
                end_color = gradient.get("end_color")
                angle = gradient.get("angle", 0)

                fill = shape.fill
                fill.gradient()
                fill.gradient_stops[0].color.rgb = start_color
                fill.gradient_stops[1].color.rgb = end_color
                fill.gradient_angle = angle

        # Add content or other text elements as specified in template
        # (These will be filled in by the specific methods like add_content_slide, add_image_slide, etc.)

        # Store styling information in the slide's user data
        # This will be used by other methods like add_image_slide, add_table_slide, etc.
        slide.user_data = {
            "template_preset": preset,
            "image_style": self.template.get_image_style(preset) if "image_style" in preset else None,
            "table_style": self.template.get_table_style(preset) if "table_style" in preset else None,
            "chart_style": self.template.get_chart_style(preset) if "chart_style" in preset else None,
        }

        return slide

    def add_title_slide(self, title: str, subtitle: str | None = None) -> Slide:
        """Add a title slide with title and optional subtitle.

        Args:
            title: Text for the title
            subtitle: Text for the subtitle (default: None)

        Returns:
            A new Slide object configured as a title slide
        """
        # Get the title slide preset
        preset = self.template.get_preset("title_slide")

        # Create a new slide using the preset
        slide = self.add_slide_from_template("title_slide")

        # Update the title text
        title_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
        if title_shapes and len(title_shapes) > 0:
            title_shape = title_shapes[0]
            title_shape.text_frame.text = title

        # Add or update the subtitle text if provided
        if subtitle and len(title_shapes) > 1:
            subtitle_shape = title_shapes[1]
            subtitle_shape.text_frame.text = subtitle
        elif subtitle:
            # Get subtitle data from preset
            subtitle_data = preset.get("subtitle", {})
            position = subtitle_data.get("position", {"x": "20%", "y": "60%", "width": "60%", "height": "20%"})

            # Extract font information
            font_data = subtitle_data.get("font", {})
            font_name = font_data.get("name", self.DEFAULT_FONT)
            font_size = font_data.get("size", 24)
            font_bold = font_data.get("bold", False)

            # Extract alignment information
            align = subtitle_data.get("align", "center")
            vertical = subtitle_data.get("vertical", "middle")
            color = subtitle_data.get("color", "black")

            # Add the subtitle text
            Text.add(
                slide=slide,
                text=subtitle,
                position=position,
                font_name=font_name,
                font_size=font_size,
                font_bold=font_bold,
                align=align,
                vertical_align=vertical,
                color=color,
            )

        return slide

    def add_content_slide(self, title: str, use_bar: bool = True) -> Slide:
        """Add a content slide with title and optional horizontal bar.

        Args:
            title: Text for the title
            use_bar: Whether to include decorative bar (default: True)

        Returns:
            A new Slide object configured as a content slide
        """
        # Create a new slide using the preset
        slide = self.add_slide_from_template("content_slide")

        # Update the title text
        title_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
        if title_shapes and len(title_shapes) > 0:
            title_shape = title_shapes[0]
            title_shape.text_frame.text = title

        # If bar is not wanted, remove it
        if not use_bar:
            for shape in list(slide.shapes):
                if shape.shape_type == MSO_SHAPE.RECTANGLE and not shape.has_text_frame:
                    shape._element.getparent().remove(shape._element)

        return slide

    def add_section_slide(self, title: str, bg_color: str = "blue") -> Slide:
        """Add a section slide with a full-screen title on a colored background.

        Args:
            title: Text for the section title
            bg_color: Background color (default: "blue")

        Returns:
            A new Slide object configured as a section slide
        """
        # Create custom preset for section slide with specified background color
        preset = self.template.get_preset("section_slide")
        preset["bg_color"] = bg_color

        # Create a new slide using the preset
        slide = self.add_slide_from_template(preset)

        # Update the title text
        title_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
        if title_shapes and len(title_shapes) > 0:
            title_shape = title_shapes[0]
            title_shape.text_frame.text = title

        return slide

    def add_image_gen_slide(
        self,
        image_path: str,
        title: str | None = None,
        subtitle: str | None = None,
        label: str | None = None,
        x: float | str = "10%",
        y: float | str = "20%",
        width: float | str = "80%",
        height: float | str = "70%",
        title_height: float | str = "10%",
        subtitle_height: float | str = "5%",
        bg_color: str | tuple[int, int, int] | None = None,
        title_font_size: int = 24,
        subtitle_font_size: int = 18,
        label_font_size: int = 14,
        border: bool = False,
        border_color: str = "black",
        shadow: bool = False,
        maintain_aspect_ratio: bool = True,
    ) -> tuple[Slide, PPTXShape]:
        """Add a slide with an image and optional title, subtitle, and label.

        This method provides a more flexible alternative to add_image_slide with
        similar parameters to add_grid_slide and add_pyplot_slide.

        Args:
            image_path: Path to the image file
            title: Optional title for the slide (default: None)
            subtitle: Optional subtitle for the slide (default: None)
            label: Optional caption for the image, displayed below (default: None)
            x: X position of the image as percentage or absolute value (default: "10%")
            y: Y position of the image as percentage or absolute value (default: "20%")
            width: Width of the image as percentage or absolute value (default: "80%")
            height: Height of the image as percentage or absolute value (default: "70%")
            title_height: Height of the title area (default: "10%")
            subtitle_height: Height of the subtitle area (default: "5%")
            bg_color: Background color for the slide (default: None)
            title_font_size: Font size for the title (default: 24)
            subtitle_font_size: Font size for the subtitle (default: 18)
            label_font_size: Font size for the caption (default: 14)
            border: Whether to add a border around the image (default: False)
            border_color: Color for the border (default: "black")
            shadow: Whether to add a shadow effect to the image (default: False)
            maintain_aspect_ratio: Whether to maintain the image's aspect ratio (default: True)

        Returns:
            A tuple containing (Slide, PPTXShape) where PPTXShape is the image shape

        Example:
            ```python
            # Add an image with title and subtitle
            slide, image = pres.add_image_gen_slide(
                image_path="path/to/image.jpg",
                title="Product Showcase",
                subtitle="Latest Design",
                label="Figure 1: Product Prototype",
                maintain_aspect_ratio=True
            )
            ```
        """
        # Create a new slide
        slide = self.add_slide(bg_color=bg_color)

        # Calculate positions and dimensions
        adjusted_y = y
        adjusted_height = height

        # Add title if provided
        if title:
            # Add the title to the slide
            slide.add_text(
                text=title,
                x="0%",  # Center the title across the slide
                y="0%",
                width="100%",
                height=title_height,
                font_size=title_font_size,
                font_bold=True,
                align="center",
                vertical="middle",
            )

            # Adjust y position for what comes next
            if isinstance(y, str) and y.endswith("%"):
                y_percent = float(y.strip("%"))
                title_height_percent = float(str(title_height).strip("%"))
                adjusted_y = f"{(y_percent + title_height_percent):.2f}%"

                # Adjust height to account for title
                if isinstance(height, str) and height.endswith("%"):
                    height_percent = float(height.strip("%"))
                    adjusted_height = f"{(height_percent - title_height_percent):.2f}%"

        # Add subtitle if provided
        if subtitle:
            # Add the subtitle to the slide
            slide.add_text(
                text=subtitle,
                x="0%",  # Center the subtitle across the slide
                y=adjusted_y,
                width="100%",
                height=subtitle_height,
                font_size=subtitle_font_size,
                align="center",
                vertical="middle",
            )

            # Adjust y position for the image
            if isinstance(adjusted_y, str) and adjusted_y.endswith("%"):
                y_percent = float(adjusted_y.strip("%"))
                subtitle_height_percent = float(str(subtitle_height).strip("%"))
                adjusted_y = f"{(y_percent + subtitle_height_percent):.2f}%"

                # Adjust height to account for subtitle
                if isinstance(adjusted_height, str) and adjusted_height.endswith("%"):
                    height_percent = float(adjusted_height.strip("%"))
                    adjusted_height = f"{(height_percent - subtitle_height_percent):.2f}%"

        # Add the image to the slide
        img = Image(slide)
        image_shape = img.add(
            image_path=image_path,
            x=x,
            y=adjusted_y,
            width=width,
            height=adjusted_height,
            maintain_aspect_ratio=maintain_aspect_ratio,
        )

        # Apply styling to the image
        if border:
            image_shape.line.color.rgb = self.COLORS.get(border_color, self.COLORS["black"])
            image_shape.line.width = 1  # 1 point width for border

        # Apply shadow if specified
        if shadow:
            image_shape.shadow.inherit = False
            image_shape.shadow.visible = True
            image_shape.shadow.blur_radius = 5
            image_shape.shadow.distance = 3
            image_shape.shadow.angle = 45

        # Add label if specified
        if label:
            # Calculate the position below the image
            if (
                isinstance(adjusted_y, str)
                and adjusted_y.endswith("%")
                and isinstance(adjusted_height, str)
                and adjusted_height.endswith("%")
            ):
                image_y = float(adjusted_y.strip("%"))
                image_height = float(adjusted_height.strip("%"))
                label_y = f"{(image_y + image_height + 1):.2f}%"  # Add a small gap
            else:
                # Fall back to a reasonable default if we can't calculate exactly
                label_y = "95%"

            # Add the label text
            slide.add_text(
                text=label,
                x="0%",  # Center the label
                y=label_y,
                width="100%",
                height="5%",
                font_size=label_font_size,
                align="center",
                vertical="top",
            )

        return slide, image_shape

    def add_image_slide(
        self, title: str, image_path: str, label: str | None = None, custom_style: dict | None = None
    ) -> Slide:
        """Add a slide with a title and a centered image.

        Args:
            title: Text for the title
            image_path: Path to the image file
            label: Optional caption for the image (default: None)
            custom_style: Optional custom styling for the image (default: None)

        Returns:
            A new Slide object configured as an image slide
        """
        # Get the image slide preset
        preset = self.template.get_preset("image_slide")

        # Create a new slide using the preset
        slide = self.add_slide_from_template("image_slide")

        # Update the title text
        title_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
        if title_shapes and len(title_shapes) > 0:
            title_shape = title_shapes[0]
            title_shape.text_frame.text = title

        # Add the image
        image_area = preset.get("image_area", {}).get(
            "position", {"x": "10%", "y": "20%", "width": "80%", "height": "70%"}
        )

        # Get styling from preset or custom style
        image_style = self.template.get_image_style(preset)
        if custom_style:
            image_style.update(custom_style)

        # Create an Image object and call its add method
        img = Image(slide)
        image_shape = img.add(
            image_path=image_path,
            x=image_area.get("x", "10%"),
            y=image_area.get("y", "20%"),
            width=image_area.get("width", "80%"),
            height=image_area.get("height", "70%"),
            maintain_aspect_ratio=image_style.get("maintain_aspect_ratio", True),
        )

        # Apply styling to the image
        if image_style.get("border", False):
            image_shape.line.color.rgb = self.COLORS.get(image_style.get("border_color", "black"), self.COLORS["black"])
            image_shape.line.width = image_style.get("border_width", 1)

        # Apply shadow if specified
        if image_style.get("shadow", False):
            image_shape.shadow.inherit = False
            image_shape.shadow.visible = True
            image_shape.shadow.blur_radius = 5
            image_shape.shadow.distance = 3
            image_shape.shadow.angle = 45

        # Apply brightness/contrast if specified
        if "brightness" in image_style or "contrast" in image_style:
            brightness = image_style.get("brightness", 0)
            contrast = image_style.get("contrast", 0)
            if hasattr(image_shape, "brightness_contrast"):
                image_shape.brightness_contrast.brightness = brightness
                image_shape.brightness_contrast.contrast = contrast

        # Add label if specified
        if label:
            # Position the label below the image
            label_position = {
                "x": image_area["x"],
                "y": str(float(image_area["y"].replace("%", "")) + float(image_area["height"].replace("%", "")) + 2)
                + "%",
                "width": image_area["width"],
                "height": "5%",
            }

            Text.add(
                slide=slide,
                text=label,
                position=label_position,
                font_name=self.DEFAULT_FONT,
                font_size=14,
                font_bold=False,
                align="center",
                vertical_align="top",
                color="black",
            )

        return slide

    def add_comparison_slide(self, title: str, content_texts: list[str]) -> Slide:
        """Add a slide with title and two or more content areas for comparison.

        Args:
            title: Text for the title
            content_texts: List of texts for comparison areas (typically 2)

        Returns:
            A new Slide object configured as a comparison slide
        """
        # Create a new slide using the preset
        slide = self.add_slide_from_template("comparison_slide")

        # Update the title text
        title_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
        if title_shapes and len(title_shapes) > 0:
            title_shape = title_shapes[0]
            title_shape.text_frame.text = title

        # Get the preset for positioning information
        preset = self.template.get_preset("comparison_slide")

        # Add content texts
        if len(content_texts) >= 1:
            left_position = preset.get("left_content", {}).get(
                "position", {"x": "5%", "y": "20%", "width": "42%", "height": "75%"}
            )
            Text.add(
                slide=slide,
                text=content_texts[0],
                position=left_position,
                font_name=self.DEFAULT_FONT,
                font_size=16,
                font_bold=False,
                align="left",
                vertical_align="top",
                color="black",
            )

        if len(content_texts) >= 2:
            right_position = preset.get("right_content", {}).get(
                "position", {"x": "53%", "y": "20%", "width": "42%", "height": "75%"}
            )
            Text.add(
                slide=slide,
                text=content_texts[1],
                position=right_position,
                font_name=self.DEFAULT_FONT,
                font_size=16,
                font_bold=False,
                align="left",
                vertical_align="top",
                color="black",
            )

        return slide

    def add_table_slide(
        self,
        title: str,
        data: list[list[Any]] | pd.DataFrame,
        has_header: bool = True,
        custom_style: dict | None = None,
    ) -> Slide:
        """Add a slide with a title and a table.

        Args:
            title: Text for the title
            data: Table data as a list of lists or pandas DataFrame
            has_header: Whether the first row is a header (default: True)
            custom_style: Dictionary of style options for the table (default: None)

        Returns:
            A new Slide object configured as a table slide
        """
        # Create a new slide using the preset
        slide = self.add_slide_from_template("table_slide")

        # Update the title text
        title_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
        if title_shapes and len(title_shapes) > 0:
            title_shape = title_shapes[0]
            title_shape.text_frame.text = title

        # Get the preset for positioning information
        preset = self.template.get_preset("table_slide")
        table_position = preset.get("table_area", {}).get(
            "position", {"x": "10%", "y": "20%", "width": "80%", "height": "70%"}
        )

        # Get styling from preset or custom style
        table_style = self.template.get_table_style(preset)
        if custom_style:
            # Handle nested dictionaries
            if "first_row" in custom_style and "first_row" in table_style:
                table_style["first_row"].update(custom_style.get("first_row", {}))

            # Update all other keys
            for key, value in custom_style.items():
                if key != "first_row":
                    table_style[key] = value

        # Add the table with styling
        table = Table(slide)

        # Extract position values
        x = table_position.get("x", "10%")
        y = table_position.get("y", "20%")
        width = table_position.get("width", "80%")
        height = table_position.get("height", "60%")

        # Convert pandas DataFrame to list if needed
        table_data: list[list[Any]] = []
        if hasattr(data, "to_dict") and callable(getattr(data, "to_dict", None)) and isinstance(data, pd.DataFrame):
            # It's a pandas DataFrame
            table_data = [list(data.columns), *data.values.tolist()]
        else:
            # It's already a list of lists
            table_data = data  # type: ignore[assignment]

        # Add the table
        table.add(
            data=table_data,
            x=x,
            y=y,
            width=width,
            height=height,
            first_row_header=has_header,
            style=1 if table_style else None,
        )

        return slide

    def add_chart_slide(
        self,
        title: str,
        data: list[list[Any]] | pd.DataFrame,
        chart_type: str | None = None,
        category_column: str | None = None,
        value_columns: str | list[str] | None = None,
        custom_style: dict | None = None,
    ) -> Slide:
        """Add a slide with a title and a chart.

        Args:
            title: Text for the title
            data: Chart data as a list of lists or pandas DataFrame
            chart_type: Type of chart (default: None uses preset's chart_type)
            category_column: Name or index of the column to use as categories (default: None)
            value_columns: Names or indices of columns to use as values (default: None)
            custom_style: Dictionary of style options for the chart (default: None)

        Returns:
            A new Slide object configured as a chart slide
        """
        # Create a new slide using the preset
        slide = self.add_slide_from_template("chart_slide")

        # Update the title text
        title_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
        if title_shapes and len(title_shapes) > 0:
            title_shape = title_shapes[0]
            title_shape.text_frame.text = title

        # Get the preset for positioning information
        preset = self.template.get_preset("chart_slide")
        chart_position = preset.get("chart_area", {}).get(
            "position", {"x": "10%", "y": "20%", "width": "80%", "height": "70%"}
        )

        # Get styling from preset or custom style
        chart_style = self.template.get_chart_style(preset)
        if custom_style:
            chart_style.update(custom_style)

        # Use chart_type from parameters or from style
        if chart_type:
            chart_style["chart_type"] = chart_type

        # Import Chart class here to avoid circular import
        from easypptx.chart import Chart

        # Create a Chart object and add the chart
        chart_obj = Chart(slide)

        # Extract chart data from data parameter
        categories: list[Any] = []
        values: list[Any] = []

        # Simple extraction of categories and values for the example
        if hasattr(data, "to_dict") and callable(getattr(data, "to_dict", None)) and isinstance(data, pd.DataFrame):
            # It's a pandas DataFrame
            df_data: pd.DataFrame = data  # Create a properly typed reference
            if category_column is not None and category_column in df_data.columns:
                categories = df_data[category_column].tolist()
            else:
                categories = df_data.iloc[:, 0].tolist()

            if value_columns is not None:
                if isinstance(value_columns, str) and value_columns in df_data.columns:
                    values = df_data[value_columns].tolist()
                elif isinstance(value_columns, list) and len(value_columns) > 0 and value_columns[0] in df_data.columns:
                    values = df_data[value_columns[0]].tolist()
                else:
                    values = df_data.iloc[:, 1].tolist()
            else:
                values = df_data.iloc[:, 1].tolist()
        else:
            # It's a list of lists
            if data and len(data) > 1:
                categories = [row[0] for row in data[1:]]
                values = [row[1] for row in data[1:]]

        # Extract position values
        x = chart_position.get("x", "10%")
        y = chart_position.get("y", "20%")
        width = chart_position.get("width", "80%")
        height = chart_position.get("height", "70%")

        chart = chart_obj.add(
            chart_type=chart_style.get("chart_type", "column"),
            categories=categories,
            values=values,
            x=x,
            y=y,
            width=width,
            height=height,
            title=title,
            has_legend=chart_style.get("has_legend", True),
        )

        # Apply additional styling if applicable
        if hasattr(chart, "format") and chart_style.get("has_border", True):
            chart.format.line.color.rgb = self.COLORS.get(
                chart_style.get("border_color", "black"), self.COLORS["black"]
            )

        # Apply custom palette if specified
        if "palette" in chart_style and hasattr(chart, "series"):
            for i, series in enumerate(chart.series):
                if i < len(chart_style["palette"]):
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = chart_style["palette"][i]

        return slide

    def add_matplotlib_slide(
        self,
        title: str,
        figure,
        label: str | None = None,
        dpi: int = 300,
        file_format: str = "png",
        custom_style: dict | None = None,
    ) -> Slide:
        """Add a slide with a title and a matplotlib figure.

        Args:
            title: Text for the title
            figure: Matplotlib figure object (plt.figure() or plt.gcf())
            label: Optional caption for the figure (default: None)
            dpi: Resolution for the figure (default: 300)
            file_format: Image format ("png" or "jpg") (default: "png")
            custom_style: Optional custom styling for the image (default: None)

        Returns:
            A new Slide object with the matplotlib figure

        Example:
            ```python
            import matplotlib.pyplot as plt

            # Create a matplotlib figure
            plt.figure(figsize=(10, 6))
            plt.plot([1, 2, 3, 4], [1, 4, 9, 16])
            plt.title('Sample Plot')

            # Add it to a presentation
            slide = pres.add_matplotlib_slide(
                title="Matplotlib Example",
                figure=plt.gcf(),
                label="Figure 1: Sample Plot"
            )
            ```
        """
        # Create a new slide using the image slide preset
        slide = self.add_slide_from_template("image_slide")

        # Update the title text
        title_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
        if title_shapes and len(title_shapes) > 0:
            title_shape = title_shapes[0]
            title_shape.text_frame.text = title

        # Get the preset for positioning information
        preset = self.template.get_preset("image_slide")
        image_area = preset.get("image_area", {}).get(
            "position", {"x": "10%", "y": "20%", "width": "80%", "height": "70%"}
        )

        # Get styling from preset or custom style
        image_style = self.template.get_image_style(preset)
        if custom_style:
            image_style.update(custom_style)

        # Add the matplotlib figure
        Pyplot.add(slide=slide, figure=figure, position=image_area, dpi=dpi, file_format=file_format, style=image_style)

        # Add label if specified
        if label:
            # Position the label below the image
            label_position = {
                "x": image_area["x"],
                "y": str(float(image_area["y"].replace("%", "")) + float(image_area["height"].replace("%", "")) + 2)
                + "%",
                "width": image_area["width"],
                "height": "5%",
            }

            Text.add(
                slide=slide,
                text=label,
                position=label_position,
                font_name=self.DEFAULT_FONT,
                font_size=14,
                font_bold=False,
                align="center",
                vertical_align="top",
                color="black",
            )

        return slide

    def add_seaborn_slide(
        self,
        title: str,
        seaborn_plot,
        label: str | None = None,
        dpi: int = 300,
        file_format: str = "png",
        custom_style: dict | None = None,
    ) -> Slide:
        """Add a slide with a title and a seaborn plot.

        Args:
            title: Text for the title
            seaborn_plot: Seaborn plot object (sns.barplot, sns.heatmap, etc.)
            label: Optional caption for the figure (default: None)
            dpi: Resolution for the figure (default: 300)
            file_format: Image format ("png" or "jpg") (default: "png")
            custom_style: Optional custom styling for the image (default: None)

        Returns:
            A new Slide object with the seaborn plot

        Example:
            ```python
            import seaborn as sns

            # Create a seaborn plot
            tips = sns.load_dataset("tips")
            sns_plot = sns.barplot(x="day", y="total_bill", data=tips)

            # Add it to a presentation
            slide = pres.add_seaborn_slide(
                title="Seaborn Example",
                seaborn_plot=sns_plot,
                label="Figure 1: Tips by Day"
            )
            ```
        """
        # Extract the figure from the seaborn plot
        if hasattr(seaborn_plot, "figure"):
            figure = seaborn_plot.figure
        elif hasattr(seaborn_plot, "fig"):
            figure = seaborn_plot.fig
        else:
            import matplotlib.pyplot as plt

            figure = plt.gcf()

        # Use the matplotlib slide method
        return self.add_matplotlib_slide(
            title=title, figure=figure, label=label, dpi=dpi, file_format=file_format, custom_style=custom_style
        )

    def add_plot(
        self, title: str, plot=None, data=None, plot_type: str = "matplotlib", label: str | None = None, **kwargs
    ) -> Slide:
        """Universal method to add various types of plots to a slide.

        This is a convenience method that supports both matplotlib/seaborn plots
        and native PowerPoint charts, providing a unified interface.

        Args:
            title: Text for the slide title
            plot: Plot object (matplotlib.figure, seaborn plot, etc.) for plot_type="matplotlib" or "seaborn"
            data: Data for PowerPoint charts for plot_type="pptx_chart"
            plot_type: Type of plot ("matplotlib", "seaborn", "pptx_chart") (default: "matplotlib")
            label: Optional caption for the plot
            **kwargs: Additional arguments specific to the plot type

        Returns:
            A new Slide object with the plot

        Example:
            ```python
            # Add a matplotlib plot
            import matplotlib.pyplot as plt
            plt.plot([1, 2, 3, 4], [1, 4, 9, 16])
            slide = pres.add_plot(
                title="Matplotlib Plot",
                plot=plt.gcf(),
                plot_type="matplotlib"
            )

            # Add a seaborn plot
            import seaborn as sns
            tips = sns.load_dataset("tips")
            sns_plot = sns.barplot(x="day", y="total_bill", data=tips)
            slide = pres.add_plot(
                title="Seaborn Plot",
                plot=sns_plot,
                plot_type="seaborn"
            )

            # Add a PowerPoint chart
            import pandas as pd
            data = pd.DataFrame({"Category": ["A", "B", "C"], "Value": [1, 4, 2]})
            slide = pres.add_plot(
                title="PowerPoint Chart",
                data=data,
                plot_type="pptx_chart",
                chart_type="column",
                category_column="Category",
                value_columns="Value"
            )
            ```
        """
        if plot_type == "matplotlib":
            if plot is None:
                raise ValueError("'plot' argument is required for matplotlib plots")

            return self.add_matplotlib_slide(
                title=title,
                figure=plot,
                label=label,
                dpi=kwargs.get("dpi", 300),
                file_format=kwargs.get("file_format", "png"),
                custom_style=kwargs.get("custom_style"),
            )

        elif plot_type == "seaborn":
            if plot is None:
                raise ValueError("'plot' argument is required for seaborn plots")

            return self.add_seaborn_slide(
                title=title,
                seaborn_plot=plot,
                label=label,
                dpi=kwargs.get("dpi", 300),
                file_format=kwargs.get("file_format", "png"),
                custom_style=kwargs.get("custom_style"),
            )

        elif plot_type == "pptx_chart":
            if data is None:
                raise ValueError("'data' argument is required for PowerPoint charts")

            return self.add_chart_slide(
                title=title,
                data=data,
                chart_type=kwargs.get("chart_type"),
                category_column=kwargs.get("category_column"),
                value_columns=kwargs.get("value_columns"),
                custom_style=kwargs.get("custom_style"),
            )

        else:
            raise ValueError(f"Unsupported plot_type: {plot_type}. Use 'matplotlib', 'seaborn', or 'pptx_chart'.")

    # Direct object API methods for adding content to slides

    def add_text(
        self,
        slide: Slide,
        text: str,
        x: float | str = 1.0,
        y: float | str = 1.0,
        width: float | str = 8.0,
        height: float | str = 1.0,
        font_size: int = 18,
        font_bold: bool = False,
        font_italic: bool = False,
        font_name: str | None = None,
        align: str = "left",
        vertical: str = "top",
        color: str | tuple[int, int, int] | None = "black",
    ) -> PPTXShape:
        """Add text directly to a slide.

        Args:
            slide: The slide to add text to
            text: The text content
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Width in inches or percentage (default: 8.0)
            height: Height in inches or percentage (default: 1.0)
            font_size: Font size in points (default: 18)
            font_bold: Whether text should be bold (default: False)
            font_italic: Whether text should be italic (default: False)
            font_name: Font name (default: None uses DEFAULT_FONT)
            align: Text alignment, one of "left", "center", "right" (default: "left")
            vertical: Vertical alignment, one of "top", "middle", "bottom" (default: "top")
            color: Text color as string name from COLORS dict or RGB tuple (default: "black")
        Returns:
            The created text shape

        Example:
            ```python
            slide = pres.add_slide()
            pres.add_text(slide, "Hello World", x="10%", y="20%", font_size=24)

            # For centering text
            pres.add_text(slide, "Centered Title", x="50%", y="5%", align="center")
            ```
        """
        if font_name is None:
            font_name = self.DEFAULT_FONT

        # Forward all parameters to add_text
        return slide.add_text(
            text=text,
            x=x,
            y=y,
            width=width,
            height=height,
            font_size=font_size,
            font_bold=font_bold,
            font_italic=font_italic,
            font_name=font_name,
            align=align,
            vertical=vertical,
            color=color,
        )

    def add_image(
        self,
        slide: Slide,
        image_path: str,
        x: float | str = 1.0,
        y: float | str = 1.0,
        width: float | str | None = None,
        height: float | str | None = None,
        crop: bool = False,
        maintain_aspect_ratio: bool = True,
        center: bool = True,
        border: bool = False,
        border_color: str = "black",
        border_width: int = 1,
        shadow: bool = False,
    ) -> PPTXShape:
        """Add an image directly to a slide.

        Args:
            slide: The slide to add the image to
            image_path: Path to the image file
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Width in inches or percentage (default: None, auto-sized)
            height: Height in inches or percentage (default: None, auto-sized)
            crop: Whether to crop the image (default: False)
            maintain_aspect_ratio: Whether to maintain aspect ratio (default: True)
            center: Whether to center the image in the specified position (default: True)
            border: Whether to add a border to the image (default: False)
            border_color: Border color (default: "black")
            border_width: Border width in points (default: 1)
            shadow: Whether to add a shadow to the image (default: False)
        Returns:
            The created image shape

        Example:
            ```python
            slide = pres.add_slide()
            pres.add_image(slide, "path/to/image.jpg", x="10%", y="20%", width="60%")

            # For centering an image
            pres.add_image(slide, "path/to/image.jpg", x="50%", y="30%", width="80%", center=True)
            ```
        """
        # Use the slide's native add_image method
        image_shape = slide.add_image(image_path=image_path, x=x, y=y, width=width, height=height)

        # Apply styling
        if border:
            image_shape.line.color.rgb = self.COLORS.get(border_color, self.COLORS["black"])
            image_shape.line.width = border_width

        if shadow:
            image_shape.shadow.inherit = False
            image_shape.shadow.visible = True
            image_shape.shadow.blur_radius = 5
            image_shape.shadow.distance = 3
            image_shape.shadow.angle = 45

        return image_shape

    def add_shape(
        self,
        slide: Slide,
        shape_type: int = MSO_SHAPE.RECTANGLE,
        x: float | str = 1.0,
        y: float | str = 1.0,
        width: float | str = 2.0,
        height: float | str = 1.0,
        fill_color: str | tuple[int, int, int] | None = None,
        line_color: str | tuple[int, int, int] | None = None,
        line_width: float = 1.0,
        text: str | None = None,
        font_size: int = 14,
        font_name: str | None = None,
        font_bold: bool = False,
        font_color: str | tuple[int, int, int] | None = "black",
        text_align: str = "center",
        text_vertical: str = "middle",
    ) -> PPTXShape:
        """Add a shape directly to a slide.

        Args:
            slide: The slide to add the shape to
            shape_type: The type of shape (default: MSO_SHAPE.RECTANGLE)
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Width in inches or percentage (default: 2.0)
            height: Height in inches or percentage (default: 1.0)
            fill_color: Fill color (default: None, no fill)
            line_color: Line color (default: None, no line)
            line_width: Line width in points (default: 1.0)
            text: Text to add to the shape (default: None)
            font_size: Font size in points (default: 14)
            font_name: Font name (default: None uses DEFAULT_FONT)
            font_bold: Whether text should be bold (default: False)
            font_color: Text color (default: "black")
            text_align: Text alignment (default: "center")
            text_vertical: Vertical text alignment (default: "middle")
        Returns:
            The created shape

        Example:
            ```python
            slide = pres.add_slide()
            pres.add_shape(
                slide,
                shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                x="20%",
                y="30%",
                width="40%",
                height="20%",
                fill_color="blue",
                text="Button",
                font_color="white"
            )

            # For centered shape
            pres.add_shape(
                slide,
                x="50%",
                y="40%",
                width="80%",
                height="30%",
                fill_color="blue",
                text_align="center"
            )
            ```
        """
        # Use the slide's native add_shape method
        shape = slide.add_shape(shape_type=shape_type, x=x, y=y, width=width, height=height, fill_color=fill_color)

        # Apply line color if specified
        if line_color is not None:
            if isinstance(line_color, str) and line_color in self.COLORS:
                shape.line.color.rgb = self.COLORS[line_color]
            elif isinstance(line_color, tuple) and len(line_color) == 3:
                shape.line.color.rgb = RGBColor(*line_color)
            shape.line.width = Pt(line_width)
        else:
            shape.line.fill.background()

        # Add text if specified
        if text is not None:
            shape.text = text
            if font_name is None:
                font_name = self.DEFAULT_FONT

            # Format text
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = font_name
                paragraph.font.bold = font_bold

                # Set text color
                if isinstance(font_color, str) and font_color in self.COLORS:
                    paragraph.font.color.rgb = self.COLORS[font_color]
                elif isinstance(font_color, tuple) and len(font_color) == 3:
                    paragraph.font.color.rgb = RGBColor(*font_color)

                # Set text alignment
                if text_align in self.ALIGN:
                    paragraph.alignment = self.ALIGN[text_align]

            # Set vertical alignment
            if text_vertical in self.VERTICAL:
                shape.text_frame.vertical_anchor = self.VERTICAL[text_vertical]

        return shape

    def add_table(
        self,
        slide: Slide,
        data: list[list[Any]] | pd.DataFrame,
        x: float | str = 1.0,
        y: float | str = 1.0,
        width: float | str = 8.0,
        height: float | str = 4.0,
        has_header: bool = True,
        style: dict | None = None,
    ) -> PPTXShape:
        """Add a table directly to a slide.

        Args:
            slide: The slide to add the table to
            data: Table data as a list of lists or pandas DataFrame
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Width in inches or percentage (default: 8.0)
            height: Height in inches or percentage (default: 4.0)
            has_header: Whether the first row is a header (default: True)
            style: Dictionary of style options for the table (default: None)

        Returns:
            The created table shape

        Example:
            ```python
            slide = pres.add_slide()
            data = [["Name", "Value"], ["Item 1", 100], ["Item 2", 200]]
            pres.add_table(slide, data, x="10%", y="20%", width="80%", height="60%")
            ```
        """
        from easypptx.table import Table

        table_obj = Table(slide)

        # Get default styling if not provided
        if style is None:
            style = self.template.default_table_style.copy()

        # Convert the style parameter to int for table.add compatibility
        table_style_id = None
        if isinstance(style, dict) and "style_id" in style:
            table_style_id = style.get("style_id")

        # Convert pandas DataFrame to list if needed
        table_data: list[list[Any]] = []
        if hasattr(data, "to_dict") and callable(getattr(data, "to_dict", None)) and isinstance(data, pd.DataFrame):
            # It's a pandas DataFrame
            table_data = [list(data.columns), *data.values.tolist()]
        else:
            # It's already a list of lists
            table_data = data  # type: ignore[assignment]

        return table_obj.add(
            data=table_data, x=x, y=y, width=width, height=height, first_row_header=has_header, style=table_style_id
        )

    def add_chart(
        self,
        slide: Slide,
        data: list[list[Any]] | pd.DataFrame,
        chart_type: str = "column",
        x: float | str = 1.0,
        y: float | str = 1.0,
        width: float | str = 8.0,
        height: float | str = 4.0,
        has_legend: bool = True,
        legend_position: str = "bottom",
        category_column: str | int | None = None,
        value_columns: str | list[str] | int | list[int] | None = None,
        has_title: bool = True,
        chart_title: str | None = None,
        has_data_labels: bool = False,
        gridlines: bool = True,
    ) -> PPTXChart:
        """Add a chart directly to a slide.

        Args:
            slide: The slide to add the chart to
            data: Chart data as a list of lists or pandas DataFrame
            chart_type: Type of chart (default: "column")
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Width in inches or percentage (default: 8.0)
            height: Height in inches or percentage (default: 4.0)
            has_legend: Whether to include a legend (default: True)
            legend_position: Legend position (default: "bottom")
            category_column: Name or index of the column to use as categories (default: None)
            value_columns: Names or indices of columns to use as values (default: None)
            has_title: Whether to include a title (default: True)
            chart_title: Chart title (default: None)
            has_data_labels: Whether to include data labels (default: False)
            gridlines: Whether to include gridlines (default: True)

        Returns:
            The created chart object

        Example:
            ```python
            slide = pres.add_slide()
            data = [["Category", "Value"], ["A", 10], ["B", 20], ["C", 30]]
            pres.add_chart(
                slide,
                data,
                chart_type="pie",
                x="10%",
                y="20%",
                width="80%",
                height="60%",
                chart_title="Sales by Region"
            )
            ```
        """
        from easypptx.chart import Chart

        chart_obj = Chart(slide)

        # Check if data is a pandas DataFrame
        if hasattr(data, "to_dict") and callable(getattr(data, "to_dict", None)):
            # Use from_dataframe method
            # Handle multi-column case for pandas DataFrames
            # Get the first column for the chart - we'll need to implement multi-series charts in the future
            first_value_column = value_columns[0] if isinstance(value_columns, list) else value_columns

            if category_column is None and hasattr(data, "columns"):
                # Use the first column as the category column
                # Only for pandas DataFrames
                category_column = str(data.columns[0])

            if first_value_column is None and hasattr(data, "columns"):
                # Use the second column as the value column
                # Only for pandas DataFrames
                if len(data.columns) > 1:
                    first_value_column = str(data.columns[1])
                else:
                    raise ValueError("DataFrame must have at least two columns for automatic value extraction")

            # Convert category_column and value_column to strings for type compatibility
            if category_column is not None:
                category_column = str(category_column)
            if first_value_column is not None:
                first_value_column = str(first_value_column)

            # Make sure we're working with a proper pandas DataFrame
            # Import pandas locally to avoid circular import
            import pandas as pd

            # Convert data to DataFrame if it's not already using ternary operator
            df = pd.DataFrame(data) if not isinstance(data, pd.DataFrame) else data

            return chart_obj.from_dataframe(
                df=df,
                chart_type=chart_type,
                category_column=category_column,
                value_column=first_value_column,
                x=x,
                y=y,
                width=width,
                height=height,
                title=chart_title,
                has_legend=has_legend,
                legend_position=legend_position,
            )
        else:
            # If it's a list of lists, convert to categories and values
            categories = []
            values = []

            if data and len(data) > 1:
                if category_column is not None:
                    # Extract column data using the provided column index or name
                    if isinstance(category_column, int):
                        categories = [row[category_column] for row in data[1:]]
                    else:
                        # Find the column index by name
                        header = data[0]
                        try:
                            col_idx = header.index(category_column)
                            categories = [row[col_idx] for row in data[1:]]
                        except ValueError:
                            raise ValueError(f"Category column '{category_column}' not found in header") from None
                else:
                    # Default to first column
                    categories = [row[0] for row in data[1:]]

                if value_columns is not None:
                    # Extract values from the specified column(s)
                    if isinstance(value_columns, int | str):
                        # Single column
                        if isinstance(value_columns, int):
                            values = [row[value_columns] for row in data[1:]]
                        else:
                            # Find column by name
                            header = data[0]
                            try:
                                col_idx = header.index(value_columns)
                                values = [row[col_idx] for row in data[1:]]
                            except ValueError:
                                raise ValueError(f"Value column '{value_columns}' not found in header") from None
                    else:
                        # Multiple columns not supported in simple list format
                        raise ValueError(
                            "Multiple value columns not supported for list data. Use pandas DataFrame instead."
                        )
                else:
                    # Default to second column
                    if len(data[0]) > 1:
                        values = [row[1] for row in data[1:]]
                    else:
                        raise ValueError("Data must have at least two columns for automatic value extraction")

            return chart_obj.add(
                chart_type=chart_type,
                categories=categories,
                values=values,
                x=x,
                y=y,
                width=width,
                height=height,
                title=chart_title,
                has_legend=has_legend,
            )

    def add_grid(
        self,
        slide: Slide,
        x: float | str = "0%",
        y: float | str = "0%",
        width: float | str = "100%",
        height: float | str = "100%",
        rows: int = 1,
        cols: int = 1,
        padding: float = 5.0,
    ) -> Grid:
        """Add a grid layout to a slide.

        Args:
            slide: The slide to add the grid to
            x: X position in inches or percentage (default: "0%")
            y: Y position in inches or percentage (default: "0%")
            width: Width in inches or percentage (default: "100%")
            height: Height in inches or percentage (default: "100%")
            rows: Number of rows in the grid (default: 1)
            cols: Number of columns in the grid (default: 1)
            padding: Padding between cells as percentage of cell size (default: 5.0)

        Returns:
            The created Grid object

        Example:
            ```python
            slide = pres.add_slide()
            grid = pres.add_grid(slide, rows=2, cols=2)

            grid.add_to_cell(
                row=0,
                col=0,
                content_func=slide.add_text,
                text="Top Left",
                font_size=24,
                align="center",
                vertical="middle",
            )
            ```
        """
        from easypptx.grid import Grid

        # Create the grid
        grid = Grid(
            parent=slide,
            x=x,
            y=y,
            width=width,
            height=height,
            rows=rows,
            cols=cols,
            padding=padding,
        )

        return grid

    def add_simple_grid_slide(
        self,
        rows: int = 1,
        cols: int = 1,
        title: str | None = None,
        title_height: float | str = "10%",
        padding: float = 5.0,
        bg_color: str | tuple[int, int, int] | None = None,
    ) -> tuple[Slide, Grid]:
        """Add a slide with a simple grid layout.

        The grid will fill the entire slide, except for the title area if a title is provided.
        This is a simplified version of add_grid_slide for backward compatibility.

        Args:
            rows: Number of rows in the grid (default: 1)
            cols: Number of columns in the grid (default: 1)
            title: Optional title for the slide (default: None)
            title_height: Height of the title area (default: "10%")
            padding: Padding between cells as percentage of cell size (default: 5.0)
            bg_color: Background color for the slide (default: None)

        Returns:
            A tuple containing (Slide, Grid)

        Example:
            ```python
            slide, grid = pres.add_simple_grid_slide(rows=2, cols=3, title="My Grid Slide")

            grid.add_to_cell(
                row=0,
                col=0,
                content_func=slide.add_text,
                text="Top Left",
                font_size=24,
                align="center",
                vertical="middle",
            )
            ```
        """
        # Create a new slide
        slide = self.add_slide(bg_color=bg_color)

        # Add title if specified
        if title:
            slide.add_text(
                text=title,
                x="0%",
                y="0%",
                width="100%",
                height=title_height,
                font_size=24,
                font_bold=True,
                align="center",
                vertical="middle",
            )

            # Adjust grid position and height to account for title
            grid_y = title_height

            # Calculate grid height by subtracting title height
            if isinstance(title_height, str) and title_height.endswith("%"):
                title_height_value = float(title_height.strip("%"))
                grid_height = f"{100 - title_height_value}%"
            else:
                grid_height = f"{100 - float(title_height)}%"

            # Create the grid
            grid = self.add_grid(
                slide=slide,
                x="0%",
                y=grid_y,
                width="100%",
                height=grid_height,
                rows=rows,
                cols=cols,
                padding=padding,
            )
        else:
            # Create the grid with full slide dimensions
            grid = self.add_grid(
                slide=slide,
                x="0%",
                y="0%",
                width="100%",
                height="100%",
                rows=rows,
                cols=cols,
                padding=padding,
            )

        return slide, grid

    def add_autogrid(
        self,
        slide: Slide,
        content_funcs: list | None = None,
        rows: int | None = None,
        cols: int | None = None,
        x: float | str = "0%",
        y: float | str = "0%",
        width: float | str = "100%",
        height: float | str = "100%",
        padding: float = 5.0,
        title: str | None = None,
        title_height: float | str = "10%",
        title_align: str = "center",
        column_major: bool = True,  # Use column-major order by default
    ) -> Grid:
        """Add an autogrid layout to a slide.

        This method automatically places the provided content functions into a grid.
        If content_funcs is None, an empty grid is created that can be populated later.

        Args:
            slide: The slide to add the autogrid to
            content_funcs: List of functions that add content to the slide, or None for an empty grid
            rows: Number of rows (default: None, calculated automatically)
            cols: Number of columns (default: None, calculated automatically)
            x: X position in inches or percentage (default: "0%")
            y: Y position in inches or percentage (default: "0%")
            width: Width in inches or percentage (default: "100%")
            height: Height in inches or percentage (default: "100%")
            padding: Padding between cells as percentage of cell size (default: 5.0)
            title: Optional title for the grid (default: None)
            title_height: Height of the title area (default: "10%")
            title_align: Text alignment for the title, one of "left", "center", "right" (default: "center")
            column_major: Whether to fill cells in column-major order (default: True)
                         When True, fills cells down columns first, resulting in a visual layout
                         that matches the specified rows and columns when content is added sequentially.
                         When False, fills cells across rows first.

        Returns:
            The created Grid object

        Example:
            ```python
            slide = pres.add_slide()

            # Option 1: With content functions
            def create_text1():
                return slide.add_text("Text 1")

            def create_text2():
                return slide.add_text("Text 2")

            content_funcs = [create_text1, create_text2]
            pres.add_autogrid(slide, content_funcs, title="Auto Grid Example")

            # Option 2: Empty grid that can be populated later
            grid = pres.add_autogrid(slide, None, rows=2, cols=3)
            grid.add_to_cell(0, 0, slide.add_text, text="Cell 0,0")
            grid.add_to_cell(1, 2, slide.add_text, text="Cell 1,2")
            ```
        """
        from easypptx.grid import Grid

        # If content_funcs is None and rows/cols are provided, create an empty grid
        if content_funcs is None:
            # Make sure rows and cols are specified for empty grid
            if rows is None or cols is None:
                rows = rows or 1
                cols = cols or 1

            # Adjust grid position and dimensions if a title is provided
            adjusted_y = y
            adjusted_height = height

            if title:
                if isinstance(y, str) and y.endswith("%"):
                    y_percent = float(y.strip("%"))
                    title_height_percent = float(str(title_height).strip("%"))
                    # Adjust y position for the grid
                    adjusted_y = f"{(y_percent + title_height_percent):.2f}%"

                    # Adjust height to account for title
                    if isinstance(height, str) and height.endswith("%"):
                        height_percent = float(height.strip("%"))
                        adjusted_height = f"{(height_percent - title_height_percent):.2f}%"

                # Add the title to the slide
                slide.add_text(
                    text=title,
                    x=x,
                    y=y,
                    width=width,
                    height=title_height,
                    font_size=24,
                    font_bold=True,
                    align=title_align,
                    vertical="middle",
                )

            # Create empty grid with specified dimensions
            grid = Grid(
                parent=slide,
                x=x,
                y=adjusted_y,
                width=width,
                height=adjusted_height,
                rows=rows,
                cols=cols,
                padding=padding,
            )
        else:
            # Use the Grid.autogrid method for content_funcs
            grid = Grid.autogrid(
                parent=slide,
                content_funcs=content_funcs,
                rows=rows,
                cols=cols,
                x=x,
                y=y,
                width=width,
                height=height,
                padding=padding,
                title=title,
                title_height=title_height,
                column_major=column_major,
            )

        return grid

    def add_grid_slide(
        self,
        rows: int,
        cols: int,
        title: str | None = None,
        subtitle: str | None = None,
        title_height: float | str = "10%",
        subtitle_height: float | str = "5%",
        x: float | str = "0%",
        y: float | str = "0%",
        width: float | str = "100%",
        height: float | str = "100%",
        padding: float = 5.0,
        bg_color: str | tuple[int, int, int] | None = None,
        title_font_size: int = 24,
        subtitle_font_size: int = 18,
        title_align: str = "center",
        subtitle_align: str = "center",
    ) -> tuple[Slide, Grid]:
        """Add a slide with a grid layout.

        This method creates a new slide with an empty grid that can be populated later.
        It provides flexible options for positioning and sizing the grid, as well as
        adding a title and subtitle to the slide.

        Args:
            rows: Number of rows in the grid
            cols: Number of columns in the grid
            title: Optional title for the slide (default: None)
            subtitle: Optional subtitle for the slide (default: None)
            title_height: Height of the title area (default: "10%")
            subtitle_height: Height of the subtitle area (default: "5%")
            x: X position of the grid as percentage or absolute value (default: "0%")
            y: Y position of the grid as percentage or absolute value (default: "0%")
            width: Width of the grid as percentage or absolute value (default: "100%")
            height: Height of the grid as percentage or absolute value (default: "100%")
            padding: Padding between cells as percentage of cell size (default: 5.0)
            bg_color: Background color for the slide (default: None)
            title_font_size: Font size for the title (default: 24)
            subtitle_font_size: Font size for the subtitle (default: 18)
            title_align: Text alignment for the title, one of "left", "center", "right" (default: "center")
            subtitle_align: Text alignment for the subtitle, one of "left", "center", "right" (default: "center")

        Returns:
            A tuple containing (Slide, Grid)

        Example:
            ```python
            # Create a slide with a 3x2 grid and a title
            slide, grid = pres.add_grid_slide(
                rows=3,
                cols=2,
                title="Features Overview",
                subtitle="Product Capabilities",
                padding=5.0,
                title_align="left"
            )

            # Add content to specific cells
            grid[0, 0].add_text("Feature 1", font_bold=True)
            grid[0, 1].add_image("image1.png")

            # Add content to rows sequentially
            grid[1].add_text("Feature 2", font_bold=True)
            grid[1].add_text("Description of Feature 2")
            ```
        """
        # Create a new slide
        slide = self.add_slide(bg_color=bg_color)

        # Calculate positions and dimensions
        adjusted_y = y
        adjusted_height = height

        # Add title if provided
        if title:
            # Add the title to the slide
            slide.add_text(
                text=title,
                x=x,
                y=y,
                width=width,
                height=title_height,
                font_size=title_font_size,
                font_bold=True,
                align=title_align,
                vertical="middle",
            )

            # Adjust y position for what comes next
            if isinstance(y, str) and y.endswith("%"):
                y_percent = float(y.strip("%"))
                title_height_percent = float(str(title_height).strip("%"))
                adjusted_y = f"{(y_percent + title_height_percent):.2f}%"

                # Adjust height to account for title
                if isinstance(height, str) and height.endswith("%"):
                    height_percent = float(height.strip("%"))
                    adjusted_height = f"{(height_percent - title_height_percent):.2f}%"

        # Add subtitle if provided
        if subtitle:
            # Add the subtitle to the slide
            slide.add_text(
                text=subtitle,
                x=x,
                y=adjusted_y,
                width=width,
                height=subtitle_height,
                font_size=subtitle_font_size,
                align=subtitle_align,
                vertical="middle",
            )

            # Adjust y position for the grid
            if isinstance(adjusted_y, str) and adjusted_y.endswith("%"):
                y_percent = float(adjusted_y.strip("%"))
                subtitle_height_percent = float(str(subtitle_height).strip("%"))
                adjusted_y = f"{(y_percent + subtitle_height_percent):.2f}%"

                # Adjust height to account for subtitle
                if isinstance(adjusted_height, str) and adjusted_height.endswith("%"):
                    height_percent = float(adjusted_height.strip("%"))
                    adjusted_height = f"{(height_percent - subtitle_height_percent):.2f}%"

        # Create the grid
        grid = Grid(
            parent=slide,
            x=x,
            y=adjusted_y,
            width=width,
            height=adjusted_height,
            rows=rows,
            cols=cols,
            padding=padding,
        )

        return slide, grid

    def add_autogrid_slide(
        self,
        content_funcs: list | None = None,
        rows: int | None = None,
        cols: int | None = None,
        title: str | None = None,
        title_height: float | str = "10%",
        padding: float = 5.0,
        bg_color: str | tuple[int, int, int] | None = None,
        title_align: str = "center",
        column_major: bool = True,  # Use column-major order by default
    ) -> tuple[Slide, Grid]:
        """Add a slide with an autogrid layout.

        This method creates a new slide and automatically places the provided
        content functions into a grid. If content_funcs is None, it creates an
        empty grid with the specified rows and columns that can be populated later.

        Args:
            content_funcs: List of functions that add content to the slide, or None for an empty grid
            rows: Number of rows (default: None, calculated automatically when content_funcs provided)
            cols: Number of columns (default: None, calculated automatically when content_funcs provided)
            title: Optional title for the slide (default: None)
            title_height: Height of the title area (default: "10%")
            padding: Padding between cells as percentage of cell size (default: 5.0)
            bg_color: Background color for the slide (default: None)
            title_align: Text alignment for the title, one of "left", "center", "right" (default: "center")
            column_major: Whether to fill cells in column-major order (default: True)
                         When True, fills cells down columns first, resulting in a visual layout
                         that matches the specified rows and columns when content is added sequentially.
                         When False, fills cells across rows first.

        Returns:
            A tuple containing (Slide, Grid)

        Example:
            ```python
            # With content functions
            def create_text1():
                return slide.add_text("Text 1")

            def create_text2():
                return slide.add_text("Text 2")

            content_funcs = [create_text1, create_text2]
            slide, grid = pres.add_autogrid_slide(content_funcs, title="Auto Grid Slide", title_align="left")

            # With empty grid
            slide, grid = pres.add_autogrid_slide(None, rows=4, cols=2, title="Features")

            # Add content directly to rows
            grid[0].add_text("Feature 1", font_bold=True)
            grid[0].add_text("Description 1")
            grid[1].add_text("Feature 2", font_bold=True)
            grid[1].add_text("Description 2")
            ```
        """
        # Create a new slide
        slide = self.add_slide(bg_color=bg_color)

        # If content_funcs is None and rows/cols are provided, ensure they have values
        if content_funcs is None and (rows is None or cols is None):
            rows = rows or 1
            cols = cols or 1

        # Create the autogrid (without title, we'll add it separately to the slide)
        if title:
            # Add the title to the slide
            slide.add_text(
                text=title,
                x="0%",
                y="0%",
                width="100%",
                height=title_height,
                font_size=24,
                font_bold=True,
                align=title_align,
                vertical="middle",
            )

            # Calculate grid dimensions - preserve the original title_height format
            grid_y = title_height

            # Calculate grid height by subtracting title height
            if isinstance(title_height, str) and title_height.endswith("%"):
                title_height_value = float(title_height.strip("%"))
                grid_height = f"{100 - title_height_value}%"
            else:
                grid_height = f"{100 - float(title_height)}%"

            # Create the autogrid without its own title (we already added it)
            grid = self.add_autogrid(
                slide=slide,
                content_funcs=content_funcs,
                rows=rows,
                cols=cols,
                x="0%",
                y=grid_y,
                width="100%",
                height=grid_height,
                padding=padding,
                title=None,  # No separate title for the grid
                title_align=title_align,
                column_major=column_major,
            )
        else:
            # Create the autogrid with full slide dimensions
            grid = self.add_autogrid(
                slide=slide,
                content_funcs=content_funcs,
                rows=rows,
                cols=cols,
                x="0%",
                y="0%",
                width="100%",
                height="100%",
                padding=padding,
                title=None,  # No title
                column_major=column_major,
            )

        return slide, grid

    def add_pyplot_slide(
        self,
        figure,
        title: str | None = None,
        subtitle: str | None = None,
        label: str | None = None,
        x: float | str = "10%",
        y: float | str = "20%",
        width: float | str = "80%",
        height: float | str = "70%",
        title_height: float | str = "10%",
        subtitle_height: float | str = "5%",
        dpi: int = 300,
        file_format: str = "png",
        bg_color: str | tuple[int, int, int] | None = None,
        title_font_size: int = 24,
        subtitle_font_size: int = 18,
        label_font_size: int = 14,
        border: bool = False,
        border_color: str = "black",
        shadow: bool = False,
        maintain_aspect_ratio: bool = True,
        title_align: str = "center",
        subtitle_align: str = "center",
        label_align: str = "center",
    ) -> tuple[Slide, PPTXShape]:
        """Add a slide with a title and a matplotlib/seaborn figure.

        This method provides more flexibility than add_matplotlib_slide and add_seaborn_slide
        by allowing control over positioning, titles, and styling, similar to add_grid_slide.

        Args:
            figure: Matplotlib or Seaborn figure object
            title: Optional title for the slide (default: None)
            subtitle: Optional subtitle for the slide (default: None)
            label: Optional caption for the figure, displayed below (default: None)
            x: X position of the figure as percentage or absolute value (default: "10%")
            y: Y position of the figure as percentage or absolute value (default: "20%")
            width: Width of the figure as percentage or absolute value (default: "80%")
            height: Height of the figure as percentage or absolute value (default: "70%")
            title_height: Height of the title area (default: "10%")
            subtitle_height: Height of the subtitle area (default: "5%")
            dpi: Resolution for the figure in dots per inch (default: 300)
            file_format: Image format ("png" or "jpg") (default: "png")
            bg_color: Background color for the slide (default: None)
            title_font_size: Font size for the title (default: 24)
            subtitle_font_size: Font size for the subtitle (default: 18)
            label_font_size: Font size for the caption (default: 14)
            border: Whether to add a border around the figure (default: False)
            border_color: Color for the border (default: "black")
            shadow: Whether to add a shadow effect to the figure (default: False)
            maintain_aspect_ratio: Whether to maintain the figure's aspect ratio (default: True)
            title_align: Text alignment for the title, one of "left", "center", "right" (default: "center")
            subtitle_align: Text alignment for the subtitle, one of "left", "center", "right" (default: "center")
            label_align: Text alignment for the caption, one of "left", "center", "right" (default: "center")

        Returns:
            A tuple containing (Slide, PPTXShape) where PPTXShape is the figure shape

        Example:
            ```python
            import matplotlib.pyplot as plt

            # Create a matplotlib figure
            plt.figure(figsize=(8, 6))
            plt.plot([1, 2, 3, 4], [1, 4, 9, 16])
            plt.title('Sample Plot')
            plt.grid(True)

            # Add it to a presentation with title and subtitle
            slide, pyplot = pres.add_pyplot_slide(
                figure=plt.gcf(),
                title="Data Visualization",
                subtitle="Matplotlib Example",
                label="Figure 1: Sample Plot",
                dpi=300,
                title_align="left"
            )
            ```
        """
        # Create a new slide
        slide = self.add_slide(bg_color=bg_color)

        # Calculate positions and dimensions
        adjusted_y = y
        adjusted_height = height

        # Add title if provided
        if title:
            # Add the title to the slide
            slide.add_text(
                text=title,
                x="0%",  # Center the title across the slide
                y="0%",
                width="100%",
                height=title_height,
                font_size=title_font_size,
                font_bold=True,
                align=title_align,
                vertical="middle",
            )

            # Adjust y position for what comes next
            if isinstance(y, str) and y.endswith("%"):
                y_percent = float(y.strip("%"))
                title_height_percent = float(str(title_height).strip("%"))
                adjusted_y = f"{(y_percent + title_height_percent):.2f}%"

                # Adjust height to account for title
                if isinstance(height, str) and height.endswith("%"):
                    height_percent = float(height.strip("%"))
                    adjusted_height = f"{(height_percent - title_height_percent):.2f}%"

        # Add subtitle if provided
        if subtitle:
            # Add the subtitle to the slide
            slide.add_text(
                text=subtitle,
                x="0%",  # Center the subtitle across the slide
                y=adjusted_y,
                width="100%",
                height=subtitle_height,
                font_size=subtitle_font_size,
                align=subtitle_align,
                vertical="middle",
            )

            # Adjust y position for the pyplot
            if isinstance(adjusted_y, str) and adjusted_y.endswith("%"):
                y_percent = float(adjusted_y.strip("%"))
                subtitle_height_percent = float(str(subtitle_height).strip("%"))
                adjusted_y = f"{(y_percent + subtitle_height_percent):.2f}%"

                # Adjust height to account for subtitle
                if isinstance(adjusted_height, str) and adjusted_height.endswith("%"):
                    height_percent = float(adjusted_height.strip("%"))
                    adjusted_height = f"{(height_percent - subtitle_height_percent):.2f}%"

        # Create a style dictionary for the pyplot
        style = {
            "border": border,
            "border_color": border_color,
            "shadow": shadow,
            "maintain_aspect_ratio": maintain_aspect_ratio,
        }

        # Add the pyplot to the slide
        pyplot = Pyplot.add(
            slide=slide,
            figure=figure,
            position={
                "x": x,
                "y": adjusted_y,
                "width": width,
                "height": adjusted_height,
            },
            dpi=dpi,
            file_format=file_format,
            style=style,
        )

        # Add label if specified
        if label:
            # Calculate the position below the figure
            if (
                isinstance(adjusted_y, str)
                and adjusted_y.endswith("%")
                and isinstance(adjusted_height, str)
                and adjusted_height.endswith("%")
            ):
                figure_y = float(adjusted_y.strip("%"))
                figure_height = float(adjusted_height.strip("%"))
                label_y = f"{(figure_y + figure_height + 1):.2f}%"  # Add a small gap
            else:
                # Fall back to a reasonable default if we can't calculate exactly
                label_y = "95%"

            # Add the label text
            slide.add_text(
                text=label,
                x="0%",  # Center the label
                y=label_y,
                width="100%",
                height="5%",
                font_size=label_font_size,
                align=label_align,
                vertical="top",
            )

        return slide, pyplot

    def add_pyplot(
        self,
        slide: Slide,
        figure,
        x: float | str = 1.0,
        y: float | str = 1.0,
        width: float | str = 8.0,
        height: float | str = 4.0,
        dpi: int = 300,
        file_format: str = "png",
        border: bool = False,
        border_color: str = "black",
        border_width: int = 1,
        shadow: bool = False,
        maintain_aspect_ratio: bool = True,
        center: bool = True,
    ) -> PPTXShape:
        """Add a matplotlib or seaborn figure directly to a slide.

        Args:
            slide: The slide to add the figure to
            figure: Matplotlib figure object (plt.figure() or plt.gcf())
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Width in inches or percentage (default: 8.0)
            height: Height in inches or percentage (default: 4.0)
            dpi: Resolution for the figure (default: 300)
            file_format: Image format ("png" or "jpg") (default: "png")
            border: Whether to add a border to the image (default: False)
            border_color: Border color (default: "black")
            border_width: Border width in points (default: 1)
            shadow: Whether to add a shadow to the image (default: False)
            maintain_aspect_ratio: Whether to maintain aspect ratio (default: True)
            center: Whether to center the image in the specified position (default: True)

        Returns:
            The created image shape

        Example:
            ```python
            import matplotlib.pyplot as plt

            # Create a matplotlib figure
            plt.figure(figsize=(10, 6))
            plt.plot([1, 2, 3, 4], [1, 4, 9, 16])
            plt.title('Sample Plot')

            # Add it to a slide
            slide = pres.add_slide()
            pres.add_pyplot(
                slide,
                plt.gcf(),
                x="10%",
                y="20%",
                width="80%",
                height="60%",
                border=True,
                shadow=True
            )
            ```
        """
        position = {"x": x, "y": y, "width": width, "height": height}
        style = {
            "border": border,
            "border_color": border_color,
            "border_width": border_width,
            "shadow": shadow,
            "maintain_aspect_ratio": maintain_aspect_ratio,
            "center": center,
        }

        return Pyplot.add(slide=slide, figure=figure, position=position, dpi=dpi, file_format=file_format, style=style)


# The Grid class is now imported properly at the top of the file
