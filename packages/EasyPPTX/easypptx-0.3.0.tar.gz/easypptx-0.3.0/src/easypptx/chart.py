"""Chart handling module for EasyPPTX."""

from typing import TYPE_CHECKING, Any, ClassVar

import pandas as pd
from pptx.chart.chart import Chart as PPTXChart
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches

if TYPE_CHECKING:
    from easypptx.slide import Slide

# Type for position parameters - accepts either percentage or absolute values
PositionType = float | str


class Chart:
    """Class for handling chart operations in PowerPoint slides.

    This class provides methods for creating and manipulating charts on slides.

    Examples:
        ```python
        # Create a chart object
        chart = Chart(slide)

        # Add a bar chart
        chart.add_bar(
            categories=["A", "B", "C"],
            values=[1, 2, 3],
            title="Sample Bar Chart"
        )

        # Add a pie chart from DataFrame
        import pandas as pd
        df = pd.DataFrame({"Category": ["A", "B", "C"], "Value": [10, 20, 30]})
        chart.from_dataframe(df, chart_type="pie", x=2, y=2)
        ```
    """

    CHART_TYPES: ClassVar = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "area": XL_CHART_TYPE.AREA,
        "scatter": XL_CHART_TYPE.XY_SCATTER,
    }

    LEGEND_POSITIONS: ClassVar = {
        "right": XL_LEGEND_POSITION.RIGHT,
        "left": XL_LEGEND_POSITION.LEFT,
        "top": XL_LEGEND_POSITION.TOP,
        "bottom": XL_LEGEND_POSITION.BOTTOM,
        "corner": XL_LEGEND_POSITION.CORNER,
    }

    def __init__(self, slide_obj: "Slide") -> None:
        """Initialize a Chart object.

        Args:
            slide_obj: The Slide object to add charts to
        """
        self.slide = slide_obj

    def add(
        self,
        chart_type: str,
        categories: list,
        values: list,
        x: PositionType = 1.0,
        y: PositionType = 1.0,
        width: PositionType = 6.0,
        height: PositionType = 4.5,
        title: str | None = None,
        has_legend: bool = True,
        **kwargs: Any,
    ) -> PPTXChart:
        """Add a chart to the slide.

        Args:
            chart_type: Type of chart ('column', 'bar', 'line', 'pie', 'area', 'scatter')
            categories: List of category labels
            values: List of data values
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Width in inches or percentage (default: 6.0)
            height: Height in inches or percentage (default: 4.5)
            title: Chart title (default: None)
            has_legend: Whether to show legend (default: True)
            **kwargs: Additional chart-specific parameters

        Returns:
            The created chart object

        Raises:
            ValueError: If chart_type is not supported or data is invalid
        """
        if chart_type not in self.CHART_TYPES:
            raise ValueError(
                f"Unsupported chart type: {chart_type}. Supported types: {', '.join(self.CHART_TYPES.keys())}"
            )

        if len(categories) != len(values):
            raise ValueError("Categories and values must have the same length")

        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series("Series 1", values)

        # Get slide dimensions for percentage conversion
        slide_width = self.slide._get_slide_width()
        slide_height = self.slide._get_slide_height()

        # Convert position values to inches
        x_inches = self.slide._convert_position(x, slide_width)
        y_inches = self.slide._convert_position(y, slide_height)
        width_inches = self.slide._convert_position(width, slide_width)
        height_inches = self.slide._convert_position(height, slide_height)

        chart_shape = self.slide.pptx_slide.shapes.add_chart(
            self.CHART_TYPES[chart_type],
            Inches(x_inches),
            Inches(y_inches),
            Inches(width_inches),
            Inches(height_inches),
            chart_data,
        )

        chart = chart_shape.chart

        # Set chart title if provided
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        else:
            chart.has_title = False

        # Set legend visibility
        chart.has_legend = has_legend
        if has_legend:
            legend_position = kwargs.get("legend_position", "right")
            if isinstance(legend_position, str) and legend_position in self.LEGEND_POSITIONS:
                chart.legend.position = self.LEGEND_POSITIONS[legend_position]

        return chart

    def add_bar(
        self,
        categories: list,
        values: list,
        x: PositionType = 1.0,
        y: PositionType = 1.0,
        width: PositionType = 6.0,
        height: PositionType = 4.5,
        title: str | None = None,
        **kwargs: Any,
    ) -> PPTXChart:
        """Add a bar chart to the slide.

        Args:
            categories: List of category labels
            values: List of data values
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Width in inches or percentage (default: 6.0)
            height: Height in inches or percentage (default: 4.5)
            title: Chart title (default: None)
            **kwargs: Additional chart-specific parameters

        Returns:
            The created chart object
        """
        return self.add(
            chart_type="bar",
            categories=categories,
            values=values,
            x=x,
            y=y,
            width=width,
            height=height,
            title=title,
            **kwargs,
        )

    def add_column(
        self,
        categories: list,
        values: list,
        x: PositionType = 1.0,
        y: PositionType = 1.0,
        width: PositionType = 6.0,
        height: PositionType = 4.5,
        title: str | None = None,
        **kwargs: Any,
    ) -> PPTXChart:
        """Add a column chart to the slide.

        Args:
            categories: List of category labels
            values: List of data values
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Width in inches or percentage (default: 6.0)
            height: Height in inches or percentage (default: 4.5)
            title: Chart title (default: None)
            **kwargs: Additional chart-specific parameters

        Returns:
            The created chart object
        """
        return self.add(
            chart_type="column",
            categories=categories,
            values=values,
            x=x,
            y=y,
            width=width,
            height=height,
            title=title,
            **kwargs,
        )

    def add_pie(
        self,
        categories: list,
        values: list,
        x: PositionType = 1.0,
        y: PositionType = 1.0,
        width: PositionType = 6.0,
        height: PositionType = 4.5,
        title: str | None = None,
        **kwargs: Any,
    ) -> PPTXChart:
        """Add a pie chart to the slide.

        Args:
            categories: List of category labels
            values: List of data values
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Width in inches or percentage (default: 6.0)
            height: Height in inches or percentage (default: 4.5)
            title: Chart title (default: None)
            **kwargs: Additional chart-specific parameters

        Returns:
            The created chart object
        """
        return self.add(
            chart_type="pie",
            categories=categories,
            values=values,
            x=x,
            y=y,
            width=width,
            height=height,
            title=title,
            **kwargs,
        )

    def from_dataframe(
        self,
        df: "pd.DataFrame",
        chart_type: str,
        category_column: str,
        value_column: str,
        x: PositionType = 1.0,
        y: PositionType = 1.0,
        width: PositionType = 6.0,
        height: PositionType = 4.5,
        title: str | None = None,
        has_legend: bool = True,
        **kwargs: Any,
    ) -> PPTXChart:
        """Create a chart from a pandas DataFrame.

        Args:
            df: Pandas DataFrame
            chart_type: Type of chart ('column', 'bar', 'line', 'pie', 'area', 'scatter')
            category_column: Column name to use for categories
            value_column: Column name to use for values
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Width in inches or percentage (default: 6.0)
            height: Height in inches or percentage (default: 4.5)
            title: Chart title (default: None)
            has_legend: Whether to show legend (default: True)
            **kwargs: Additional chart-specific parameters

        Returns:
            The created chart object

        Raises:
            ValueError: If columns don't exist in DataFrame
        """
        if category_column not in df.columns:
            raise ValueError(f"Category column '{category_column}' not found in DataFrame")
        if value_column not in df.columns:
            raise ValueError(f"Value column '{value_column}' not found in DataFrame")

        categories = df[category_column].tolist()
        values = df[value_column].tolist()

        return self.add(
            chart_type=chart_type,
            categories=categories,
            values=values,
            x=x,
            y=y,
            width=width,
            height=height,
            title=title,
            has_legend=has_legend,
            **kwargs,
        )
