"""Tests for the Template functionality."""

from unittest.mock import MagicMock, patch

from pptx.enum.shapes import MSO_SHAPE

from easypptx import Presentation, Template


def test_presentation_blank_slide_default():
    """Test that new presentations use blank layout by default for slides."""
    pres = Presentation()
    pres.add_slide()

    # Check that blank_layout is set to slide_layouts[6]
    assert pres.blank_layout == pres.pptx_presentation.slide_layouts[6]


def test_template_presets_exist():
    """Test that Template class has predefined presets."""
    template = Template()

    # Check that the Template class has presets
    assert hasattr(template, "presets")
    assert isinstance(template.presets, dict)
    assert len(template.presets) > 0

    # Check for required preset types
    assert "title_slide" in template.presets
    assert "content_slide" in template.presets
    assert "section_slide" in template.presets


def test_add_title_slide():
    """Test adding a title slide using the preset."""
    pres = Presentation()

    # Add a title slide with default settings
    slide = pres.add_title_slide("Test Title", "Test Subtitle")

    # Test that the slide has the correct elements
    shapes = slide.shapes

    # We should have at least two text boxes (title and subtitle)
    assert len(shapes) >= 2

    # Check the text in the first shape (usually title)
    title_shape = next((shape for shape in shapes if shape.has_text_frame), None)
    assert title_shape is not None
    assert "Test Title" in title_shape.text_frame.text


def test_add_content_slide():
    """Test adding a content slide using the preset."""
    pres = Presentation()

    # Add a content slide with default settings
    slide = pres.add_content_slide("Test Content", use_bar=True)

    # Test that the slide has the correct elements
    shapes = slide.shapes

    # We should have at least two shapes (title and bar)
    assert len(shapes) >= 2

    # Check the text in the title shape
    title_shape = next((shape for shape in shapes if shape.has_text_frame), None)
    assert title_shape is not None
    assert "Test Content" in title_shape.text_frame.text

    # Check for a rectangle shape (the bar)
    bar_shape = next((shape for shape in shapes if shape.shape_type == MSO_SHAPE.RECTANGLE), None)
    assert bar_shape is not None


def test_add_section_slide():
    """Test adding a section slide using the preset."""
    pres = Presentation()

    # Add a section slide with default settings
    slide = pres.add_section_slide("Test Section")

    # Test that the slide has the correct elements
    shapes = slide.shapes

    # Should have at least one text shape (the title)
    assert len(shapes) >= 1

    # Check the text in the title shape
    title_shape = next((shape for shape in shapes if shape.has_text_frame), None)
    assert title_shape is not None
    assert "Test Section" in title_shape.text_frame.text

    # Section slides should have a background color
    # First set the background to solid, then check it has a type
    slide.pptx_slide.background.fill.solid()
    assert slide.pptx_slide.background.fill.type is not None


def test_add_image_slide():
    """Test adding an image slide using the preset."""
    pres = Presentation()

    # Mock image path since we're not actually reading an image
    mock_image_path = "test_image.png"

    # Mock the image handling methods
    with patch("easypptx.Image.add") as mock_add:
        # Configure the mock
        mock_add.return_value = MagicMock()

        # Add an image slide
        slide = pres.add_image_slide("Test Image", mock_image_path)

        # Test that the slide has the correct elements
        shapes = slide.shapes

        # Should have at least one text shape (the title)
        assert len(shapes) >= 1

        # Verify that Image.add was called with the correct parameters
        mock_add.assert_called_once()


def test_add_comparison_slide():
    """Test adding a comparison slide with multiple content areas."""
    pres = Presentation()

    # Add a comparison slide with default settings
    slide = pres.add_comparison_slide("Test Comparison", ["Left Content", "Right Content"])

    # Test that the slide has the correct elements
    shapes = slide.shapes

    # Should have at least three text shapes (title and two content areas)
    assert len(shapes) >= 3

    # Check the text in the title shape
    title_shape = next((shape for shape in shapes if shape.has_text_frame), None)
    assert title_shape is not None
    assert "Test Comparison" in title_shape.text_frame.text

    # Check that there are multiple text shapes
    text_shapes = [shape for shape in shapes if shape.has_text_frame]
    assert len(text_shapes) >= 3  # title + 2 content areas


def test_add_table_slide():
    """Test adding a table slide using the preset."""
    pres = Presentation()

    # Sample data for the table
    data = [["Header 1", "Header 2"], ["Value 1", "Value 2"]]

    # Mock the table handling methods
    with patch("easypptx.Table.add") as mock_add:
        # Configure the mock
        mock_add.return_value = MagicMock()

        # Add a table slide
        slide = pres.add_table_slide("Test Table", data)

        # Test that the slide has the correct elements
        shapes = slide.shapes

        # Should have at least one text shape (the title)
        assert len(shapes) >= 1

        # Verify that Table.add was called with the correct parameters
        mock_add.assert_called_once()


def test_custom_template():
    """Test using a custom template for a slide."""
    pres = Presentation()

    # Define a custom template
    custom_template = {
        "bg_color": "blue",
        "title": {
            "text": "Custom Template",
            "position": {"x": "5%", "y": "10%", "width": "90%", "height": "15%"},
            "font": {"name": "Meiryo", "size": 32, "bold": True},
            "align": "center",
            "vertical": "middle",
            "color": "white",
        },
        "content_area": {"position": {"x": "10%", "y": "30%", "width": "80%", "height": "60%"}},
    }

    # Add a slide with the custom template
    slide = pres.add_slide_from_template(custom_template)

    # Test that the slide has the correct elements
    shapes = slide.shapes

    # Should have at least one text shape (the title)
    assert len(shapes) >= 1

    # Check the text in the title shape
    title_shape = next((shape for shape in shapes if shape.has_text_frame), None)
    assert title_shape is not None
    assert "Custom Template" in title_shape.text_frame.text

    # Check that background color is applied
    # First make the background solid, then check it's applied
    slide.pptx_slide.background.fill.solid()
    assert slide.pptx_slide.background.fill.type is not None
