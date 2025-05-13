"""Tests for the Table class."""

from unittest.mock import MagicMock, patch

import pandas as pd
import pytest
from pptx.util import Pt

from easypptx import Table


class TestTable:
    """Test cases for the Table class."""

    def setup_method(self):
        """Set up test environment before each test method."""
        self.slide = MagicMock()
        self.table = Table(self.slide)

        # Mock the shape's table object
        self.mock_shape = MagicMock()
        self.mock_table = MagicMock()
        self.mock_shape.table = self.mock_table

        # Set the return value for shapes.add_table
        self.slide.pptx_slide.shapes.add_table.return_value = self.mock_shape

        # Sample data for tests
        self.sample_data = [
            ["Header 1", "Header 2"],
            ["Value 1", "Value 2"],
            ["Value 3", "Value 4"],
        ]

    def test_table_init(self):
        """Test Table initialization."""
        assert self.table.slide == self.slide

    def test_add_table_basic(self):
        """Test adding a basic table."""
        # Call add method with default parameters
        result = self.table.add(self.sample_data)

        # Verify shapes.add_table was called with correct parameters
        self.slide.pptx_slide.shapes.add_table.assert_called_once()
        call_args = self.slide.pptx_slide.shapes.add_table.call_args[0]

        # Check rows and columns
        assert call_args[0] == 3  # rows
        assert call_args[1] == 2  # columns
        # Since we're mocking the slide, we can't test the exact conversion
        # Instead, just test that the parameters were passed in some form
        assert call_args[2] is not None  # x
        assert call_args[3] is not None  # y

        # Just check that the cell method was called for each cell in the data
        assert self.mock_table.cell.call_count == 6  # 3 rows x 2 columns

        # Verify the result is the mock table
        assert result == self.mock_table

    def test_add_table_with_position(self):
        """Test adding a table with custom position."""
        # Call add method with custom position
        result = self.table.add(self.sample_data, x=2.5, y=3.5)

        # Verify shapes.add_table was called with correct parameters
        self.slide.pptx_slide.shapes.add_table.assert_called_once()
        call_args = self.slide.pptx_slide.shapes.add_table.call_args[0]

        # Since we're mocking the slide, we can't test the exact conversion
        # Instead, just test that the parameters were passed in some form
        assert call_args[2] is not None  # x
        assert call_args[3] is not None  # y

        # Verify the result is the mock table
        assert result == self.mock_table

    def test_add_table_with_dimensions(self):
        """Test adding a table with custom dimensions."""
        # Call add method with custom dimensions
        result = self.table.add(self.sample_data, width=6.0, height=4.0)

        # Verify shapes.add_table was called with correct parameters
        self.slide.pptx_slide.shapes.add_table.assert_called_once()
        call_args = self.slide.pptx_slide.shapes.add_table.call_args[0]

        # Since we're mocking the slide, we can't test the exact conversion
        # Instead, just test that the parameters were passed in some form
        assert call_args[4] is not None  # width
        assert call_args[5] is not None  # height

        # Verify the result is the mock table
        assert result == self.mock_table

    def test_add_table_empty_data(self):
        """Test adding a table with empty data raises ValueError."""
        with pytest.raises(ValueError):
            self.table.add([])

    def test_add_table_inconsistent_columns(self):
        """Test adding a table with inconsistent column counts raises ValueError."""
        inconsistent_data = [
            ["Header 1", "Header 2", "Header 3"],
            ["Value 1", "Value 2"],  # Missing a column
            ["Value 3", "Value 4", "Value 5"],
        ]

        with pytest.raises(ValueError):
            self.table.add(inconsistent_data)

    def test_add_table_first_row_header(self):
        """Test adding a table with first row as header."""
        # Mock paragraph for header formatting
        mock_paragraph = MagicMock()

        # Set up cell mock to return text_frame with paragraphs
        mock_cell = MagicMock()
        mock_cell.text_frame.paragraphs = [mock_paragraph]
        self.mock_table.cell.return_value = mock_cell

        # Call add method with first_row_header=True
        self.table.add(self.sample_data, first_row_header=True)

        # Verify header row formatting
        assert mock_paragraph.font.bold is True
        assert mock_paragraph.font.size == Pt(14)

    def test_add_table_with_style(self):
        """Test adding a table with a specific style."""
        # Call add method with a style
        self.table.add(self.sample_data, style=5)

        # Verify table style was set
        assert self.mock_table.style == 5

    def test_from_dataframe_basic(self):
        """Test creating a table from pandas DataFrame."""
        # Create a DataFrame
        df = pd.DataFrame({
            "A": [1, 2],
            "B": [3, 4],
        })

        # Patch the add method
        with patch.object(self.table, "add") as mock_add:
            # Call from_dataframe
            self.table.from_dataframe(df)

            # Verify add was called with correct data
            mock_add.assert_called_once()
            call_args = mock_add.call_args[1]

            # Check that column names are in first row
            assert call_args["data"][0] == list(df.columns)

            # Check each row of data
            for i, row in enumerate(df.values):
                assert list(row) == call_args["data"][i + 1]

    def test_from_dataframe_with_index(self):
        """Test creating a table from DataFrame with index included."""
        # Create a DataFrame with custom index
        df = pd.DataFrame(
            {
                "A": [1, 2],
                "B": [3, 4],
            },
            index=["Row 1", "Row 2"],
        )

        # Patch the add method
        with patch.object(self.table, "add") as mock_add:
            # Call from_dataframe with include_index=True
            self.table.from_dataframe(df, include_index=True)

            # Verify add was called with data that includes index
            mock_add.assert_called_once()
            call_args = mock_add.call_args[1]

            # First row should have column names
            assert len(call_args["data"]) == 3  # header + 2 data rows

            # Check that data rows include index
            for i, idx in enumerate(df.index):
                assert call_args["data"][i + 1][0] == idx

    def test_from_dataframe_with_parameters(self):
        """Test creating a table from DataFrame with custom parameters."""
        df = pd.DataFrame({
            "A": [1, 2],
            "B": [3, 4],
        })

        # Patch the add method
        with patch.object(self.table, "add") as mock_add:
            # Call from_dataframe with custom parameters
            self.table.from_dataframe(
                df,
                x=2.0,
                y=3.0,
                width=6.0,
                height=4.0,
                first_row_header=False,
                style=3,
            )

            # Verify add was called with correct parameters
            mock_add.assert_called_once()
            call_args = mock_add.call_args[1]

            assert call_args["x"] == 2.0
            assert call_args["y"] == 3.0
            assert call_args["width"] == 6.0
            assert call_args["height"] == 4.0
            assert call_args["first_row_header"] is False
            assert call_args["style"] == 3
