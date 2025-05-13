"""
TOML Templates Demo

This example demonstrates how to use the TOML templates created for EasyPPTX.
It shows how to:
1. Load templates from TOML files
2. Create slides based on the loaded templates
3. Customize slide content while maintaining template styling
"""

from pathlib import Path

import matplotlib.pyplot as plt
from pptx.enum.shapes import MSO_SHAPE

from easypptx import Presentation
from easypptx.template import TemplateManager

# Create output directory for presentations
output_dir = Path("output")
output_dir.mkdir(exist_ok=True)

# Path to templates directory
templates_dir = Path("../templates")

# Initialize template manager and load templates
tm = TemplateManager(template_dir=str(templates_dir))

# Load templates from TOML files
business_title_name = tm.load(str(templates_dir / "business_title.toml"))
business_content_name = tm.load(str(templates_dir / "business_content.toml"))
dashboard_name = tm.load(str(templates_dir / "dashboard.toml"))
marketing_name = tm.load(str(templates_dir / "marketing.toml"))
tech_dark_name = tm.load(str(templates_dir / "tech_dark.toml"))

# Print available templates
print(f"Loaded templates: {', '.join(tm.list_templates())}")

# Create a new presentation
pres = Presentation()

# Create a title slide explaining the demo
slide = pres.add_slide(title="TOML Templates Demo")
slide.add_text(
    text="This presentation demonstrates how to use TOML-based templates with EasyPPTX.",
    x="10%",
    y="30%",
    width="80%",
    height="20%",
    font_size=24,
    align="center",
)
slide.add_text(
    text="The following templates are included:\n• Business Title\n• Business Content\n• Dashboard\n• Marketing\n• Tech Dark Theme",
    x="10%",
    y="50%",
    width="80%",
    height="30%",
    font_size=20,
    align="center",
)

# 1. Business Title Template Example
title_slide = pres.add_slide_from_template(business_title_name)
title_slide.add_text(
    text="Quarterly Business Review",
    x="10%",
    y="30%",
    width="80%",
    height="20%",
    font_size=44,
    font_bold=True,
    align="center",
    color="white",
)
title_slide.add_text(
    text="Q2 2025 Financial Results",
    x="10%",
    y="55%",
    width="80%",
    height="10%",
    font_size=24,
    align="center",
    color="#66ccff",
)
title_slide.add_text(
    text="ACME Corporation | Confidential",
    x="0%",
    y="90%",
    width="100%",
    height="5%",
    font_size=14,
    font_bold=True,
    align="center",
    color="#cceeff",
)

# 2. Business Content Template Example
content_slide = pres.add_slide_from_template(business_content_name)
content_slide.add_text(
    text="Financial Highlights",
    x="5%",
    y="5%",
    width="90%",
    height="10%",
    font_size=32,
    font_bold=True,
    align="left",
    color="#003366",
)
# Add title bar (a shape)
content_slide.add_shape(
    shape_type="RECTANGLE",
    x="5%",
    y="15%",
    width="90%",
    height="1%",
    fill_color="#003366",
)
# Add bullet points
content_slide.add_text(
    text=(
        "• Revenue increased by 15% year-over-year to $24.5M\n"
        "• Gross margin improved to 68% (up from 62% in Q1)\n"
        "• Operating expenses decreased by 3% due to efficiency initiatives\n"
        "• Net profit increased by 22% to $6.8M\n"
        "• Cash flow from operations up 18% to $8.2M\n\n"
        "All KPIs exceeded quarterly targets, setting us on track for our annual goals."
    ),
    x="7%",
    y="22%",
    width="86%",
    height="65%",
    font_size=20,
    align="left",
    color="black",
)
content_slide.add_text(
    text="ACME Corporation | Confidential",
    x="0%",
    y="90%",
    width="100%",
    height="5%",
    font_size=14,
    font_bold=True,
    align="center",
    color="#003366",
)
content_slide.add_text(
    text="2",
    x="90%",
    y="90%",
    width="5%",
    height="5%",
    font_size=12,
    align="right",
    color="#003366",
)

# 3. Dashboard Template Example
dashboard_slide = pres.add_slide_from_template(dashboard_name)
dashboard_slide.add_text(
    text="Q2 2025 Performance Dashboard",
    x="0%",
    y="2%",
    width="100%",
    height="8%",
    font_size=28,
    font_bold=True,
    align="center",
    color="#333333",
)

# Create a simple bar chart
fig1, ax1 = plt.subplots(figsize=(4, 3))
categories = ["Q1", "Q2", "Q3", "Q4"]
values = [18.5, 24.5, 22.0, 26.5]
ax1.bar(categories, values, color="#4472C4")
ax1.set_title("Quarterly Revenue ($M)")
ax1.grid(axis="y", linestyle="--", alpha=0.7)
for i, v in enumerate(values):
    ax1.text(i, v + 0.5, f"${v}M", ha="center")

# Create a simple line chart
fig2, ax2 = plt.subplots(figsize=(4, 3))
months = ["Jan", "Feb", "Mar", "Apr", "May", "Jun"]
data = [62, 64, 65, 66, 67, 68]
ax2.plot(months, data, marker="o", color="#4472C4", linewidth=2)
ax2.set_title("Gross Margin (%)")
ax2.set_ylim(60, 70)
ax2.grid(True, linestyle="--", alpha=0.7)

# KPI displays
dashboard_slide.add_text(
    text="Revenue",
    x="7%",
    y="14%",
    width="40%",
    height="5%",
    font_size=20,
    font_bold=True,
    color="#333333",
)
dashboard_slide.add_text(
    text="$24.5M\n+15% YoY",
    x="7%",
    y="20%",
    width="40%",
    height="30%",
    font_size=36,
    font_bold=True,
    align="center",
    vertical="middle",
    color="#4472C4",
)

dashboard_slide.add_text(
    text="Gross Margin",
    x="53%",
    y="14%",
    width="40%",
    height="5%",
    font_size=20,
    font_bold=True,
    color="#333333",
)
dashboard_slide.add_text(
    text="68%\n+6% YoY",
    x="53%",
    y="20%",
    width="40%",
    height="30%",
    font_size=36,
    font_bold=True,
    align="center",
    vertical="middle",
    color="#4472C4",
)

# Add charts to dashboard
dashboard_slide.add_text(
    text="Quarterly Revenue",
    x="7%",
    y="52%",
    width="40%",
    height="5%",
    font_size=20,
    font_bold=True,
    color="#333333",
)
dashboard_slide.add_pyplot(
    figure=fig1,
    x="7%",
    y="58%",
    width="40%",
    height="30%",
    dpi=150,
)

dashboard_slide.add_text(
    text="Gross Margin Trend",
    x="53%",
    y="52%",
    width="40%",
    height="5%",
    font_size=20,
    font_bold=True,
    color="#333333",
)
dashboard_slide.add_pyplot(
    figure=fig2,
    x="53%",
    y="58%",
    width="40%",
    height="30%",
    dpi=150,
)

dashboard_slide.add_text(
    text="Data updated as of: June 30, 2025",
    x="0%",
    y="95%",
    width="100%",
    height="3%",
    font_size=12,
    align="center",
    color="#666666",
)

# 4. Marketing Template Example
marketing_slide = pres.add_slide_from_template(marketing_name)
marketing_slide.add_text(
    text="Summer Product Launch",
    x="5%",
    y="5%",
    width="90%",
    height="15%",
    font_size=36,
    font_bold=True,
    align="left",
    color="#ff6600",
)
# Add accent bar
marketing_slide.add_shape(
    shape_type="RECTANGLE",
    x="5%",
    y="21%",
    width="30%",
    height="2%",
    fill_color="#ff6600",
)
# Image area placeholder
marketing_slide.add_shape(
    shape_type="RECTANGLE",
    x="5%",
    y="25%",
    width="45%",
    height="60%",
    fill_color="#e0e0e0",
    text="Product Image",
    font_size=24,
    text_align="center",
)
# Content area
marketing_slide.add_text(
    text="Introducing Quantum X Pro",
    x="55%",
    y="25%",
    width="40%",
    height="10%",
    font_size=28,
    font_bold=True,
    align="left",
    color="#333333",
)
marketing_slide.add_text(
    text=(
        "The next generation of our flagship product is here. "
        "Quantum X Pro offers:\n\n"
        "• 2x performance improvement\n"
        "• Enhanced user interface\n"
        "• New AI-powered features\n"
        "• Industry-leading battery life\n\n"
        "Available in stores starting July 15th."
    ),
    x="55%",
    y="37%",
    width="40%",
    height="30%",
    font_size=18,
    align="left",
    color="#333333",
)
# Call to action
marketing_slide.add_shape(
    shape_type="ROUNDED_RECTANGLE",
    x="55%",
    y="70%",
    width="40%",
    height="10%",
    fill_color="#ff6600",
    text="PRE-ORDER NOW",
    font_size=20,
    font_bold=True,
    font_color="white",
    text_align="center",
)
marketing_slide.add_text(
    text="ACME Corporation | Summer Campaign 2025",
    x="0%",
    y="90%",
    width="100%",
    height="5%",
    font_size=14,
    align="center",
    color="#666666",
)

# 5. Tech Dark Theme Example
tech_slide = pres.add_slide_from_template(tech_dark_name)
tech_slide.add_text(
    text="Implementing the New API",
    x="5%",
    y="5%",
    width="90%",
    height="15%",
    font_size=36,
    font_bold=True,
    align="left",
    color="#4fc3f7",
)
# Add accent bar
tech_slide.add_shape(
    shape_type="RECTANGLE",
    x="5%",
    y="21%",
    width="30%",
    height="2%",
    fill_color="#4fc3f7",
)
# Code example
code_sample = """# Example API request
import requests

def get_user_data(user_id):
    url = f"https://api.example.com/users/{user_id}"
    headers = {"Authorization": f"Bearer {API_TOKEN}"}

    response = requests.get(url, headers=headers)
    response.raise_for_status()

    return response.json()

# Usage
user_data = get_user_data("user123")
print(f"User: {user_data['name']}")"""

tech_slide.add_text(
    text=code_sample,
    x="5%",
    y="25%",
    width="90%",
    height="30%",
    font_size=16,
    font_name="Consolas",
    align="left",
    color="#e0e0e0",
    bg_color="#1e1e1e",
)
# Highlight area
tech_slide.add_shape(
    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
    x="70%",
    y="25%",
    width="25%",
    height="25%",
    fill_color="#4fc3f7",
    fill_opacity=0.2,
    line_color="#4fc3f7",
    line_width=2,
)
tech_slide.add_text(
    text=(
        "• API uses JWT authentication\n"
        "• Responses are in JSON format\n"
        "• Rate limited to 100 requests/minute\n"
        "• Async implementation recommended\n"
        "• Error responses include detailed info\n"
    ),
    x="7%",
    y="62%",
    width="86%",
    height="25%",
    font_size=20,
    align="left",
    color="#e0e0e0",
)
tech_slide.add_text(
    text="Technical Documentation | Confidential",
    x="0%",
    y="95%",
    width="95%",
    height="5%",
    font_size=14,
    align="right",
    color="#e0e0e0",
)

# Add a final slide summarizing the TOML templates
slide = pres.add_slide(title="TOML Templates Benefits")
slide.add_text(
    text=(
        "Benefits of using TOML templates:\n\n"
        "• Consistent branding across presentations\n"
        "• Reusable design elements\n"
        "• Easy to share and version control\n"
        "• Human-readable format\n"
        "• Can be customized for different contexts\n"
        "• Supports complex layouts and styling\n"
    ),
    x="10%",
    y="20%",
    width="80%",
    height="40%",
    font_size=24,
)

slide.add_text(
    text=(
        "To use TOML templates:\n\n"
        "1. Create template files with .toml extension\n"
        "2. Load them with the TemplateManager\n"
        "3. Use add_slide_from_template() to create slides\n"
        "4. Customize content while maintaining consistent styling\n"
    ),
    x="10%",
    y="60%",
    width="80%",
    height="30%",
    font_size=20,
)

# Save the presentation
pres.save(output_dir / "toml_templates_demo.pptx")
print(f"Presentation saved to {output_dir / 'toml_templates_demo.pptx'}")
