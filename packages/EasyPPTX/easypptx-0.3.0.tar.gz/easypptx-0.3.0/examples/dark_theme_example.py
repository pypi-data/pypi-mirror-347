"""
Dark theme example for EasyPPTX.

This example demonstrates creating presentations with dark backgrounds
and vibrant text colors for a modern, high-contrast look.
"""

from pathlib import Path

import pandas as pd
from pptx.enum.shapes import MSO_SHAPE

from easypptx import Chart, Presentation, Table, Text

# Create output directory
output_dir = Path("output")
output_dir.mkdir(exist_ok=True)


def create_dark_presentation():
    """Create a presentation with a dark theme."""
    print("Creating dark theme presentation...")

    # Create a presentation with black background as default
    pres = Presentation(default_bg_color="black")

    # Title slide
    slide = pres.add_slide()
    text = Text(slide)
    text.add_title("Dark Theme Presentation", font_size=54, color="cyan", align="center")
    text.add_paragraph(
        "Modern and Stylish Design", y="30%", font_size=32, font_bold=False, color="white", align="center"
    )
    text.add_paragraph("Created with EasyPPTX", y="80%", font_size=20, font_italic=True, color="gray", align="center")

    # Content slide with multiple elements
    slide = pres.add_slide()
    text = Text(slide)
    text.add_title("Key Features", color="yellow", align="center")

    # Create feature blocks with auto-alignment
    features = [
        {"type": "text", "text": "Dark Backgrounds", "color": "white", "font_bold": True, "align": "center"},
        {"type": "text", "text": "Vibrant Colors", "color": "cyan", "font_bold": True, "align": "center"},
        {"type": "text", "text": "Modern Layout", "color": "magenta", "font_bold": True, "align": "center"},
        {"type": "text", "text": "Auto Alignment", "color": "green", "font_bold": True, "align": "center"},
    ]

    slide.add_multiple_objects(objects_data=features, layout="grid", padding_percent=10.0, start_y="30%", height="30%")

    # Add decorative elements
    slide.add_shape(
        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, x="5%", y="70%", width="90%", height="20%", fill_color="darkgray"
    )

    text.add_paragraph(
        "High contrast design for better readability and visual impact",
        x="10%",
        y="75%",
        width="80%",
        height="10%",
        font_size=20,
        color="white",
        align="center",
        vertical="middle",
    )

    # Data visualization slide
    slide = pres.add_slide()
    text = Text(slide)
    text.add_title("Data Visualization", color="orange", align="center")

    # Sample data
    data = {"Category": ["A", "B", "C", "D", "E"], "Value": [75, 45, 90, 35, 60]}
    df = pd.DataFrame(data)

    # Add chart
    chart = Chart(slide)
    chart.from_dataframe(
        df,
        chart_type="column",
        category_column="Category",
        value_column="Value",
        x="10%",
        y="25%",
        width="45%",
        height="60%",
        title="Performance Metrics",
        has_legend=True,
    )

    # Add table
    tbl = Table(slide)
    tbl.from_dataframe(df, x="60%", y="30%", width="30%", first_row_header=True, style="Medium Style 2 - Accent 1")

    text.add_paragraph(
        "Charts and tables maintain readability on dark backgrounds",
        x="60%",
        y="60%",
        width="30%",
        font_size=16,
        color="lightgray",
        align="center",
    )

    # Save the presentation
    pres.save(output_dir / "dark_theme_example.pptx")
    print(f"Saved to {output_dir / 'dark_theme_example.pptx'}")


def create_gradient_example():
    """Create a presentation with color gradient effects."""
    print("Creating gradient effect example...")

    # Create a presentation with deep blue background
    pres = Presentation(default_bg_color=(0, 20, 40))

    # Title slide
    slide = pres.add_slide(bg_color=(0, 15, 30))  # Slightly different shade for variety
    text = Text(slide)
    text.add_title("Gradient Effect Demo", font_size=54, color="cyan", align="center")

    # Add gradient shapes to create visual effect
    # Top right corner accent
    slide.add_shape(
        shape_type=MSO_SHAPE.OVAL,
        x="70%",
        y="-10%",
        width="60%",
        height="50%",
        fill_color=(0, 80, 120),  # Lighter blue
    )

    # Bottom left corner accent
    slide.add_shape(
        shape_type=MSO_SHAPE.OVAL,
        x="-20%",
        y="70%",
        width="60%",
        height="50%",
        fill_color=(60, 0, 80),  # Purple
    )

    text.add_paragraph("Creating depth with shapes and colors", y="40%", font_size=28, color="white", align="center")

    # Demo slide with layered shapes
    slide = pres.add_slide(bg_color=(0, 20, 40))
    text = Text(slide)
    text.add_title("Layered Design Elements", color="yellow", align="center")

    # Background shapes for visual interest
    for i, color_tuple in enumerate([
        (0, 60, 120),  # Medium blue
        (0, 80, 100),  # Teal blue
        (0, 100, 80),  # Turquoise
        (40, 100, 60),  # Blue-green
    ]):
        slide.add_shape(
            shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
            x=f"{15 + i * 5}%",
            y=f"{30 + i * 10}%",
            width="70%",
            height="15%",
            fill_color=color_tuple,
        )

    # Add text on top
    for i, (content, color) in enumerate([
        ("Modern Design", "white"),
        ("Professional Look", "cyan"),
        ("High Contrast", "yellow"),
        ("Visual Hierarchy", "lightgray"),
    ]):
        text.add_paragraph(
            content,
            x="25%",
            y=f"{32 + i * 10}%",
            width="60%",
            height="10%",
            font_size=24,
            color=color,
            align="center",
            vertical="middle",
        )

    # Save the presentation
    pres.save(output_dir / "gradient_effect_example.pptx")
    print(f"Saved to {output_dir / 'gradient_effect_example.pptx'}")


# Execute examples
if __name__ == "__main__":
    print("Creating dark theme presentations...\n")

    create_dark_presentation()
    print()

    create_gradient_example()
    print()

    print("All dark theme examples created successfully!")
