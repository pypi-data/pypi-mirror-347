"""Tests for the Presentation class."""

import tempfile
from pathlib import Path

import pytest
from pptx.util import Inches

from easypptx import Presentation


def test_presentation_init():
    """Test presentation initialization."""
    pres = Presentation()
    assert pres is not None
    assert pres.pptx_presentation is not None


def test_presentation_add_slide():
    """Test adding slides to a presentation."""
    pres = Presentation()

    # Add a slide
    slide = pres.add_slide()
    assert slide is not None

    # Check slides property
    slides = pres.slides
    assert len(slides) == 1


def test_presentation_save():
    """Test saving a presentation to a file."""
    pres = Presentation()

    # Add a slide
    pres.add_slide()

    # Save to a temporary file
    with tempfile.NamedTemporaryFile(suffix=".pptx") as tmp:
        temp_path = Path(tmp.name)
        pres.save(temp_path)
        assert temp_path.exists()
        assert temp_path.stat().st_size > 0


def test_presentation_open():
    """Test opening an existing presentation."""
    # Create and save a presentation
    pres = Presentation()
    pres.add_slide()

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        temp_path = Path(tmp.name)

    try:
        pres.save(temp_path)

        # Open the saved presentation
        opened_pres = Presentation.open(temp_path)
        assert opened_pres is not None
        assert len(opened_pres.slides) == 1
    finally:
        # Clean up the temporary file
        if temp_path.exists():
            temp_path.unlink()


def test_presentation_open_nonexistent_file():
    """Test opening a nonexistent file raises FileNotFoundError."""
    with pytest.raises(FileNotFoundError):
        Presentation.open("nonexistent_file.pptx")


def test_presentation_default_aspect_ratio():
    """Test that new presentations have 16:9 aspect ratio by default."""
    pres = Presentation()

    # Get the slide width and height
    slide_width = pres.pptx_presentation.slide_width
    slide_height = pres.pptx_presentation.slide_height

    # Calculate the aspect ratio (width / height)
    aspect_ratio = slide_width / slide_height

    # Check that it's approximately 16:9 (1.77778...)
    assert round(aspect_ratio, 2) == 1.78


def test_presentation_with_aspect_ratio():
    """Test creating presentations with different aspect ratios."""
    # Test with 16:9 aspect ratio
    pres_16_9 = Presentation(aspect_ratio="16:9")
    aspect_ratio_16_9 = pres_16_9.pptx_presentation.slide_width / pres_16_9.pptx_presentation.slide_height
    assert round(aspect_ratio_16_9, 2) == 1.78

    # Test with 4:3 aspect ratio
    pres_4_3 = Presentation(aspect_ratio="4:3")
    aspect_ratio_4_3 = pres_4_3.pptx_presentation.slide_width / pres_4_3.pptx_presentation.slide_height
    assert round(aspect_ratio_4_3, 2) == 1.33


def test_presentation_custom_dimensions():
    """Test creating presentations with custom dimensions."""
    # Set custom width and height in inches
    width_inches = 13.33
    height_inches = 7.5

    pres = Presentation(width_inches=width_inches, height_inches=height_inches)

    # Convert to EMU (English Metric Units)
    expected_width = int(Inches(width_inches))
    expected_height = int(Inches(height_inches))

    assert pres.pptx_presentation.slide_width == expected_width
    assert pres.pptx_presentation.slide_height == expected_height
