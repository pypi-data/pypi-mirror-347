"""Slide module for EasyPPTX."""

from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.shapes.autoshape import Shape as PPTXShape
from pptx.slide import Slide as PPTXSlide
from pptx.util import Inches, Pt

# Type for position parameters - accepts either percentage or absolute values
PositionType = float | str


class Slide:
    """Class representing a slide in a PowerPoint presentation.

    This class provides methods for adding and manipulating content on a slide,
    such as text, images, tables, and shapes.

    Attributes:
        pptx_slide: The underlying python-pptx Slide object

    Examples:
        ```python
        # Add text to a slide using inches
        slide.add_text("Hello World", x=2, y=2)

        # Add text to a slide using percentages
        slide.add_text("Hello World", x="20%", y="30%")

        # Add an image to a slide
        slide.add_image("image.png", x=1, y=1, width=5, height=3)

        # Add an image using percentages
        slide.add_image("image.png", x="10%", y="20%", width="50%", height="30%")
        ```
    """

    def __init__(self, pptx_slide: PPTXSlide) -> None:
        """Initialize a Slide object.

        Args:
            pptx_slide: The python-pptx Slide object
        """
        self.pptx_slide = pptx_slide
        self.user_data: dict[str, Any] = {}

        # Store template defaults when applied from a template
        self.template_defaults: dict[str, dict[str, Any]] = {
            "text": {},
            "image": {},
            "shape": {},
            "table": {},
            "chart": {},
            "global": {},
        }

        # Cache slide dimensions to avoid recalculating them
        self._slide_width = self._get_slide_width()
        self._slide_height = self._get_slide_height()

    def apply_template_defaults(self, template_data: dict[str, Any]) -> None:
        """Apply template defaults to this slide.

        This method extracts default method arguments from a template and stores them
        for later use by the add_xxx methods.

        Args:
            template_data: Template data dictionary
        """
        # Extract defaults for different element types
        if "defaults" in template_data:
            defaults = template_data["defaults"]

            # Apply defaults for each element type
            for element_type in ["text", "image", "shape", "table", "chart", "global"]:
                if element_type in defaults:
                    self.template_defaults[element_type] = defaults[element_type]

    def merge_with_defaults(self, method_type: str, kwargs: dict[str, Any]) -> dict[str, Any]:
        """Merge provided arguments with template defaults.

        Args:
            method_type: The type of method ("text", "image", etc.)
            kwargs: Keyword arguments provided to the method

        Returns:
            Dictionary with merged arguments, where provided args override defaults
        """
        # Start with a copy of the method-specific defaults
        method_defaults: dict[str, Any] = self.template_defaults.get(method_type, {}).copy()

        # Get global defaults if available
        global_defaults: dict[str, Any] = self.template_defaults.get("global", {})

        # Create merged defaults: global defaults overridden by method-specific defaults
        merged_defaults: dict[str, Any] = global_defaults.copy()
        for key, value in method_defaults.items():
            merged_defaults[key] = value

        # Override with provided kwargs
        result: dict[str, Any] = merged_defaults.copy()
        for key, value in kwargs.items():
            if value is not None:  # Only override if the value is not None
                result[key] = value

        return result

    def _convert_position(self, value: PositionType, slide_dimension: int, is_width: bool = True) -> float:
        """Convert a position value to inches.

        This method supports two types of positioning:
        1. Absolute positioning in inches (e.g., 1.0, 2.5)
        2. Percentage-based positioning (e.g., "10%", "50%") relative to slide dimensions

        Percentage-based positioning makes layouts responsive to different slide sizes
        and aspect ratios, similar to the approach used in CSS layouts.

        Args:
            value: Position value (percentage string like "20%" or absolute inches)
            slide_dimension: The total slide dimension (width or height) in EMUs
            is_width: Whether this value is for width calculation (for aspect ratio adaptation)

        Returns:
            Position value in inches

        Examples:
            # Convert 50% of slide width to inches
            x_inches = slide._convert_position("50%", slide_width)

            # Use absolute positioning (returns the same value)
            x_inches = slide._convert_position(2.5, slide_width)  # Returns 2.5
        """
        # Convert EMUs to inches for calculations
        dimension_inches = slide_dimension / 914400

        if isinstance(value, str) and value.endswith("%"):
            # Convert percentage to inches - proper conversion based on slide dimensions
            percent = float(value.strip("%"))

            # Ensure accurate percentage calculation by limiting to valid range (0-100%)
            percent = max(0, min(100, percent))

            # Calculate the position in inches
            position_inches = (percent / 100) * dimension_inches
            return position_inches
        else:
            # Return absolute position in inches, ensuring it's a float
            return float(value)

    def _get_slide_width(self) -> int:
        """Get the slide width in EMUs.

        This method safely retrieves the slide width, trying different access paths
        to accommodate different contexts (real slides, test mocks, etc.)

        Returns:
            The slide width in English Metric Units (EMUs)
        """
        # Method 1: Try accessing through the presentation pack directly
        try:
            if hasattr(self.pptx_slide, "part") and hasattr(self.pptx_slide.part, "package"):
                package = self.pptx_slide.part.package
                if hasattr(package, "presentation_part") and hasattr(package.presentation_part, "presentation"):
                    return package.presentation_part.presentation.slide_width
        except (AttributeError, TypeError):
            pass

        # Method 2: Try accessing through the presentation directly
        try:
            if hasattr(self.pptx_slide, "part") and hasattr(self.pptx_slide.part, "package"):
                presentation = self.pptx_slide.part.package.presentation
                if hasattr(presentation, "slide_width"):
                    return presentation.slide_width
        except (AttributeError, TypeError):
            pass

        # Method 3: Try through parent-child relationship of Parts
        try:
            parent = self.pptx_slide.part.parent
            if parent and hasattr(parent, "presentation"):
                return parent.presentation.slide_width
        except (AttributeError, TypeError):
            pass

        # Default value - standard 16:9 presentation (13.33 x 7.5 inches)
        return 12192768  # 13.33 inches = 914400 EMUs per inch * 13.33

    def _get_slide_height(self) -> int:
        """Get the slide height in EMUs.

        This method safely retrieves the slide height, trying different access paths
        to accommodate different contexts (real slides, test mocks, etc.)

        Returns:
            The slide height in English Metric Units (EMUs)
        """
        # Method 1: Try accessing through the presentation pack directly
        try:
            if hasattr(self.pptx_slide, "part") and hasattr(self.pptx_slide.part, "package"):
                package = self.pptx_slide.part.package
                if hasattr(package, "presentation_part") and hasattr(package.presentation_part, "presentation"):
                    return package.presentation_part.presentation.slide_height
        except (AttributeError, TypeError):
            pass

        # Method 2: Try accessing through the presentation directly
        try:
            if hasattr(self.pptx_slide, "part") and hasattr(self.pptx_slide.part, "package"):
                presentation = self.pptx_slide.part.package.presentation
                if hasattr(presentation, "slide_height"):
                    return presentation.slide_height
        except (AttributeError, TypeError):
            pass

        # Method 3: Try through parent-child relationship of Parts
        try:
            parent = self.pptx_slide.part.parent
            if parent and hasattr(parent, "presentation"):
                return parent.presentation.slide_height
        except (AttributeError, TypeError):
            pass

        # Default value - standard 16:9 presentation (13.33 x 7.5 inches)
        return 6858000  # 7.5 inches = 914400 EMUs per inch

    def add_text(
        self,
        text: str,
        x: PositionType | None = None,
        y: PositionType | None = None,
        width: PositionType | None = None,
        height: PositionType | None = None,
        font_size: int | None = None,
        font_bold: bool | None = None,
        font_italic: bool | None = None,
        font_name: str | None = None,
        align: str | None = None,
        vertical: str | None = None,
        color: str | tuple[int, int, int] | None = None,
        **kwargs,  # Accept any additional parameters and ignore them
    ) -> PPTXShape:
        """Add a text box to the slide.

        Args:
            text: The text content
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Width in inches or percentage (default: 8.0)
            height: Height in inches or percentage (default: 1.0)
            font_size: Font size in points (default: 18)
            font_bold: Whether text should be bold (default: False)
            font_italic: Whether text should be italic (default: False)
            font_name: Font name (default: "Meiryo")
            align: Text alignment, one of "left", "center", "right" (default: "left")
            vertical: Vertical alignment, one of "top", "middle", "bottom" (default: "top")
            color: Text color as string name from COLORS dict or RGB tuple (default: None)
            **kwargs: Additional parameters (ignored)

        Returns:
            The created shape object
        """
        # Get defaults from template (method-specific and global)
        text_defaults = self.template_defaults.get("text", {})
        global_defaults = self.template_defaults.get("global", {})

        # Apply defaults for None values with cascading fallbacks:
        # 1. Use parameter if provided (not None)
        # 2. Otherwise use text-specific default from template
        # 3. Otherwise use global default from template
        # 4. Otherwise use fallback hardcoded default

        # Position and size
        x_val = x if x is not None else text_defaults.get("x", global_defaults.get("x", 1.0))
        y_val = y if y is not None else text_defaults.get("y", global_defaults.get("y", 1.0))
        width_val = width if width is not None else text_defaults.get("width", global_defaults.get("width", 8.0))
        height_val = height if height is not None else text_defaults.get("height", global_defaults.get("height", 1.0))

        # Font properties
        font_size_val = (
            font_size if font_size is not None else text_defaults.get("font_size", global_defaults.get("font_size", 18))
        )
        font_bold_val = (
            font_bold
            if font_bold is not None
            else text_defaults.get("font_bold", global_defaults.get("font_bold", False))
        )
        font_italic_val = (
            font_italic
            if font_italic is not None
            else text_defaults.get("font_italic", global_defaults.get("font_italic", False))
        )
        font_name_val = (
            font_name
            if font_name is not None
            else text_defaults.get("font_name", global_defaults.get("font_name", "Meiryo"))
        )

        # Alignment
        align_val = align if align is not None else text_defaults.get("align", global_defaults.get("align", "left"))
        vertical_val = (
            vertical if vertical is not None else text_defaults.get("vertical", global_defaults.get("vertical", "top"))
        )

        # Color - can be None as a valid value
        color_val = color
        if color_val is None:
            # Try to get color from text defaults
            color_val = text_defaults.get("color")
            # If still None, try to get from global defaults
            if color_val is None:
                color_val = global_defaults.get("color")

        # Convert list colors to tuples if needed
        if isinstance(color_val, list) and len(color_val) == 3:
            from typing import cast

            # Cast the color_val to a specific type for mypy
            color_list = cast(list[int], color_val)
            color_val = (color_list[0], color_list[1], color_list[2])

        # Use cached slide dimensions
        slide_width = self._slide_width
        slide_height = self._slide_height

        # Convert position values to inches
        x_inches = self._convert_position(x_val, slide_width, is_width=True)
        y_inches = self._convert_position(y_val, slide_height, is_width=False)
        width_inches = self._convert_position(width_val, slide_width, is_width=True)
        height_inches = self._convert_position(height_val, slide_height, is_width=False)

        # Create the textbox
        text_box = self.pptx_slide.shapes.add_textbox(
            Inches(x_inches), Inches(y_inches), Inches(width_inches), Inches(height_inches)
        )
        text_frame = text_box.text_frame
        text_frame.text = text
        text_frame.word_wrap = True  # Enable word wrap for better text display

        # Set auto-size based on available properties in python-pptx
        if hasattr(text_frame, "auto_size"):
            # In newer versions of python-pptx
            try:
                from pptx.enum.text import MSO_AUTO_SIZE

                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except (ImportError, AttributeError):
                # Fallback for older versions
                text_frame.auto_size = True

        # Set vertical alignment
        if vertical_val in ["top", "middle", "bottom"]:
            from easypptx.presentation import Presentation

            text_frame.vertical_anchor = Presentation.VERTICAL[vertical_val]

        # Apply text formatting
        p = text_frame.paragraphs[0]
        p.font.size = Pt(font_size_val)
        p.font.bold = font_bold_val
        p.font.italic = font_italic_val
        p.font.name = font_name_val

        # Set horizontal alignment
        if align_val in ["left", "center", "right"]:
            from easypptx.presentation import Presentation

            p.alignment = Presentation.ALIGN[align_val]

        # Set text color
        if color_val:
            from easypptx.presentation import Presentation

            if isinstance(color_val, str) and color_val in Presentation.COLORS:
                p.font.color.rgb = Presentation.COLORS[color_val]
            elif isinstance(color_val, tuple) and len(color_val) == 3:
                p.font.color.rgb = RGBColor(*color_val)

        return text_box

    def add_image(
        self,
        image_path: str,
        x: PositionType = 1.0,
        y: PositionType = 1.0,
        width: PositionType | None = None,
        height: PositionType | None = None,
        **kwargs,  # Accept any additional parameters and ignore them
    ) -> PPTXShape:
        """Add an image to the slide.

        Args:
            image_path: Path to the image file
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Width in inches or percentage (default: None, maintains aspect ratio)
            height: Height in inches or percentage (default: None, maintains aspect ratio)
            **kwargs: Additional parameters (ignored)

        Returns:
            The created picture shape

        Raises:
            FileNotFoundError: If the image file doesn't exist
        """
        # Use cached slide dimensions
        slide_width = self._slide_width
        slide_height = self._slide_height

        # Convert position values to inches
        x_inches = self._convert_position(x, slide_width, is_width=True)
        y_inches = self._convert_position(y, slide_height, is_width=False)

        # Convert size values to inches if provided
        width_inches = None
        height_inches = None

        if width is not None:
            width_inches = Inches(self._convert_position(width, slide_width, is_width=True))
        if height is not None:
            height_inches = Inches(self._convert_position(height, slide_height, is_width=False))

        return self.pptx_slide.shapes.add_picture(
            image_path, Inches(x_inches), Inches(y_inches), width_inches, height_inches
        )

    @property
    def shapes(self) -> list[PPTXShape]:
        """Get all shapes on the slide.

        Returns:
            List of shape objects
        """
        return list(self.pptx_slide.shapes)

    def clear(self) -> None:
        """Remove all shapes from the slide."""
        for shape in self.pptx_slide.shapes:
            self.pptx_slide.shapes._spTree.remove(shape._element)

    def add_shape(
        self,
        shape_type: MSO_SHAPE = MSO_SHAPE.RECTANGLE,
        x: PositionType = 1.0,
        y: PositionType = 1.0,
        width: PositionType = 5.0,
        height: PositionType = 1.0,
        fill_color: str | tuple[int, int, int] | None = None,
        **kwargs,  # Accept any additional parameters and ignore them
    ) -> PPTXShape:
        """Add a shape to the slide.

        Args:
            shape_type: The shape type (default: MSO_SHAPE.RECTANGLE)
            x: X position in inches or percentage (default: 1.0)
            y: Y position in inches or percentage (default: 1.0)
            width: Width in inches or percentage (default: 5.0)
            height: Height in inches or percentage (default: 1.0)
            fill_color: Fill color as string name from COLORS dict or RGB tuple (default: None)
            **kwargs: Additional parameters (ignored)

        Returns:
            The created shape object
        """
        # Use cached slide dimensions
        slide_width = self._slide_width
        slide_height = self._slide_height

        # Convert position values to inches
        x_inches = self._convert_position(x, slide_width, is_width=True)
        y_inches = self._convert_position(y, slide_height, is_width=False)
        width_inches = self._convert_position(width, slide_width, is_width=True)
        height_inches = self._convert_position(height, slide_height, is_width=False)

        # Create the shape
        shape = self.pptx_slide.shapes.add_shape(
            shape_type, Inches(x_inches), Inches(y_inches), Inches(width_inches), Inches(height_inches)
        )

        # Set fill color if specified
        if fill_color:
            from easypptx.presentation import Presentation

            if isinstance(fill_color, str) and fill_color in Presentation.COLORS:
                shape.fill.solid()
                shape.fill.fore_color.rgb = Presentation.COLORS[fill_color]
            elif isinstance(fill_color, tuple) and len(fill_color) == 3:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*fill_color)

        return shape

    def add_multiple_objects(
        self,
        objects_data: list[dict[str, Any]],
        layout: str = "grid",
        padding_percent: float = 5.0,
        start_x: PositionType = "5%",
        start_y: PositionType = "5%",
        width: PositionType = "90%",
        height: PositionType = "90%",
        **kwargs,  # Accept any additional parameters and ignore them
    ) -> list[PPTXShape]:
        """Add multiple objects to the slide with automatic alignment.

        Args:
            objects_data: List of dictionaries containing object data
                Each dict should have 'type' ('text', 'image', or 'shape') and type-specific parameters
            layout: Layout type ('grid', 'horizontal', 'vertical')
            padding_percent: Padding between objects as percentage of container
            start_x: Starting X position of container in inches or percentage
            start_y: Starting Y position of container in inches or percentage
            width: Width of container in inches or percentage
            height: Height of container in inches or percentage
            **kwargs: Additional parameters (ignored)

        Returns:
            List of created shape objects
        """
        # Use cached slide dimensions
        slide_width = self._slide_width
        slide_height = self._slide_height

        # Convert container position and size to inches
        container_x = self._convert_position(start_x, slide_width)
        container_y = self._convert_position(start_y, slide_height)
        container_width = self._convert_position(width, slide_width)
        container_height = self._convert_position(height, slide_height)

        # Calculate padding
        padding = padding_percent / 100.0

        # Determine number of rows and columns for grid layout
        num_objects = len(objects_data)

        if layout == "horizontal":
            cols = num_objects
            rows = 1
        elif layout == "vertical":
            cols = 1
            rows = num_objects
        else:  # Default to grid
            import math

            cols = math.ceil(math.sqrt(num_objects))
            rows = math.ceil(num_objects / cols)

        # Calculate cell dimensions including padding
        cell_width = container_width / cols
        cell_height = container_height / rows

        # Adjust for padding
        obj_width = cell_width * (1 - padding)
        obj_height = cell_height * (1 - padding)

        created_objects = []

        for i, obj_data in enumerate(objects_data):
            # Calculate position for this object
            col = i % cols
            row = i // cols

            obj_x = container_x + (col * cell_width) + (cell_width * padding / 2)
            obj_y = container_y + (row * cell_height) + (cell_height * padding / 2)

            # Add object based on type
            obj_type = obj_data.get("type", "text")

            if obj_type == "text":
                # Extract text-specific parameters with defaults
                text = obj_data.get("text", "")
                font_size = obj_data.get("font_size", 18)
                font_bold = obj_data.get("font_bold", False)
                font_italic = obj_data.get("font_italic", False)
                font_name = obj_data.get("font_name", "Meiryo")
                align = obj_data.get("align", "center")
                vertical = obj_data.get("vertical", "middle")
                color = obj_data.get("color", "black")

                obj = self.add_text(
                    text=text,
                    x=obj_x,
                    y=obj_y,
                    width=obj_width,
                    height=obj_height,
                    font_size=font_size,
                    font_bold=font_bold,
                    font_italic=font_italic,
                    font_name=font_name,
                    align=align,
                    vertical=vertical,
                    color=color,
                )

            elif obj_type == "image":
                # Extract image-specific parameters
                image_path = obj_data.get("image_path", "")

                obj = self.add_image(image_path=image_path, x=obj_x, y=obj_y, width=obj_width, height=obj_height)

            elif obj_type == "shape":
                # Extract shape-specific parameters
                shape_type = obj_data.get("shape_type", MSO_SHAPE.RECTANGLE)
                fill_color = obj_data.get("fill_color", None)

                obj = self.add_shape(
                    shape_type=shape_type, x=obj_x, y=obj_y, width=obj_width, height=obj_height, fill_color=fill_color
                )

            created_objects.append(obj)

        return created_objects

    @property
    def title(self) -> str | None:
        """Get the slide title.

        Returns:
            The slide title if it exists, None otherwise
        """
        if self.pptx_slide.shapes.title:
            return self.pptx_slide.shapes.title.text
        return None

    @title.setter
    def title(self, value: str) -> None:
        """Set the slide title.

        Args:
            value: The title text
        """
        if self.pptx_slide.shapes.title:
            self.pptx_slide.shapes.title.text = value

    def set_background_color(self, color: str | tuple[int, int, int], **kwargs) -> None:
        """Set the background color of the slide.

        Args:
            color: Background color as string name from COLORS dict or RGB tuple
            **kwargs: Additional parameters (ignored)
        """
        # Get the background fill object
        background = self.pptx_slide.background

        # First ensure background is not None or transparent
        background_p = background._element
        background_p.attrib["bwMode"] = "auto"

        # Ensure we have a solid fill
        fill = background.fill
        solid_fill = fill.solid()

        # Apply the color
        if isinstance(color, str):
            from easypptx.presentation import Presentation

            if color in Presentation.COLORS:
                fill.fore_color.rgb = Presentation.COLORS[color]
            else:
                raise ValueError(f"Color '{color}' not found in predefined colors")
        elif isinstance(color, tuple) and len(color) == 3:
            fill.fore_color.rgb = RGBColor(*color)
        else:
            raise ValueError("Color must be a string name or RGB tuple")

        return solid_fill
