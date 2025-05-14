"""
Example showing the object-related APIs in EasyPPTX.

This example demonstrates how to use the direct object manipulation methods
in the Presentation class to add various elements to slides.
"""

from pathlib import Path

import pandas as pd
from pptx.enum.shapes import MSO_SHAPE

from easypptx import Presentation

# Create a folder for outputs if it doesn't exist
output_dir = Path("output")
output_dir.mkdir(exist_ok=True)

# Create a new presentation
pres = Presentation()

# Create a title slide
title_slide = pres.add_title_slide(title="Object API Examples", subtitle="Direct manipulation of slide elements")

# Slide 1: Text
text_slide = pres.add_content_slide(title="Text Examples")

# Add text directly to the slide
pres.add_text(slide=text_slide, text="This is regular text", x="10%", y="20%", width="80%", height="8%", font_size=18)

pres.add_text(
    slide=text_slide,
    text="This is bold, italic, blue text",
    x="10%",
    y="30%",
    width="80%",
    height="8%",
    font_size=18,
    font_bold=True,
    font_italic=True,
    color="blue",
)

pres.add_text(
    slide=text_slide, text="Right-aligned text", x="10%", y="40%", width="80%", height="8%", font_size=18, align="right"
)

pres.add_text(
    slide=text_slide,
    text="Middle-aligned text with custom color",
    x="10%",
    y="50%",
    width="80%",
    height="8%",
    font_size=18,
    align="center",
    vertical="middle",
    color=(150, 75, 0),
)

# Slide 2: Shapes
shape_slide = pres.add_content_slide(title="Shape Examples")

# Add various shapes
pres.add_shape(
    slide=shape_slide,
    shape_type=MSO_SHAPE.RECTANGLE,
    x="10%",
    y="20%",
    width="30%",
    height="15%",
    fill_color="blue",
    text="Rectangle",
    font_color="white",
)

pres.add_shape(
    slide=shape_slide,
    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
    x="60%",
    y="20%",
    width="30%",
    height="15%",
    fill_color="green",
    line_color="darkgray",
    line_width=2.0,
    text="Rounded Rectangle",
    font_color="white",
)

pres.add_shape(
    slide=shape_slide,
    shape_type=MSO_SHAPE.OVAL,
    x="10%",
    y="45%",
    width="30%",
    height="15%",
    fill_color="red",
    text="Oval",
    font_color="white",
)

pres.add_shape(
    slide=shape_slide,
    shape_type=MSO_SHAPE.PENTAGON,
    x="60%",
    y="45%",
    width="30%",
    height="15%",
    fill_color="orange",
    text="Pentagon",
    font_color="white",
)

pres.add_shape(
    slide=shape_slide,
    shape_type=MSO_SHAPE.CHEVRON,
    x="10%",
    y="70%",
    width="30%",
    height="15%",
    fill_color="cyan",
    text="Chevron",
    font_color="black",
)

pres.add_shape(
    slide=shape_slide,
    shape_type=MSO_SHAPE.STAR_5_POINT,
    x="60%",
    y="70%",
    width="20%",
    height="15%",
    fill_color="yellow",
    text="Star",
    font_color="black",
)

# Slide 3: Table
table_slide = pres.add_content_slide(title="Table Example")

# Create table data
table_data = [
    ["Product", "Q1", "Q2", "Q3", "Q4"],
    ["Product A", 120, 140, 135, 150],
    ["Product B", 85, 90, 95, 110],
    ["Product C", 45, 55, 65, 70],
    ["Product D", 30, 25, 40, 35],
]

# Add table to slide
pres.add_table(
    slide=table_slide,
    data=table_data,
    x="10%",
    y="20%",
    width="80%",
    height="50%",
    has_header=True,
    style={
        "first_row": {"bold": True, "bg_color": "blue", "text_color": "white"},
        "banded_rows": True,
        "band_color": "lightgray",
    },
)

# Slide 4: Chart
chart_slide = pres.add_content_slide(title="Chart Example")

# Create chart data
chart_data = pd.DataFrame({
    "Quarter": ["Q1", "Q2", "Q3", "Q4"],
    "Sales": [120, 145, 160, 180],
    "Expenses": [95, 110, 115, 130],
    "Profit": [25, 35, 45, 50],
})

# Add chart to slide
pres.add_chart(
    slide=chart_slide,
    data=chart_data,
    chart_type="column",
    x="10%",
    y="20%",
    width="80%",
    height="60%",
    has_legend=True,
    legend_position="bottom",
    category_column="Quarter",
    value_columns=["Sales", "Expenses", "Profit"],
    has_title=True,
    chart_title="Quarterly Performance",
    has_data_labels=True,
    gridlines=True,
)

# Slide 5: Image
image_slide = pres.add_content_slide(title="Image Example")

# Create a sample image path
# Check if the output/images directory exists, if not use a placeholder text
image_dir = output_dir / "images"
if image_dir.exists() and any(image_dir.glob("*.png")):
    # Use the first image found
    images = list(image_dir.glob("*.png"))
    if images:
        sample_image = str(images[0])
    else:
        # Create a dummy image
        sample_image = None
        pres.add_text(
            slide=image_slide,
            text="Sample image not found in output/images directory",
            x="10%",
            y="40%",
            width="80%",
            height="10%",
            font_size=16,
            align="center",
        )
else:
    # Create a dummy image
    sample_image = None
    pres.add_text(
        slide=image_slide,
        text="Sample image not found in output/images directory",
        x="10%",
        y="40%",
        width="80%",
        height="10%",
        font_size=16,
        align="center",
    )

# If we have an image, add it to the slide
if sample_image:
    pres.add_image(
        slide=image_slide,
        image_path=sample_image,
        x="20%",
        y="20%",
        width="60%",
        height="60%",
        maintain_aspect_ratio=True,
        border=True,
        border_color="blue",
        shadow=True,
    )

# Slide 6: Combination of elements
combined_slide = pres.add_content_slide(title="Combined Elements")

# Add a shape as background panel
pres.add_shape(
    slide=combined_slide,
    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
    x="5%",
    y="15%",
    width="90%",
    height="75%",
    fill_color=(240, 240, 240),
    line_color="gray",
)

# Add title text
pres.add_text(
    slide=combined_slide,
    text="Dashboard Title",
    x="10%",
    y="20%",
    width="80%",
    height="10%",
    font_size=24,
    font_bold=True,
    align="center",
)

# Add a small table
table_data_small = [
    ["Metric", "Value"],
    ["Total Revenue", "$1,245,000"],
    ["Growth Rate", "+15.2%"],
    ["Conversion", "4.8%"],
]

pres.add_table(
    slide=combined_slide, data=table_data_small, x="10%", y="35%", width="35%", height="20%", has_header=True
)

# Add small chart
chart_data_small = pd.DataFrame({"Month": ["Jan", "Feb", "Mar", "Apr"], "Value": [42, 35, 65, 58]})

pres.add_chart(
    slide=combined_slide,
    data=chart_data_small,
    chart_type="line",
    x="55%",
    y="35%",
    width="35%",
    height="20%",
    category_column="Month",
    value_columns="Value",
    has_title=True,
    chart_title="Monthly Trend",
    has_legend=False,
)

# Add text boxes at the bottom
pres.add_shape(
    slide=combined_slide,
    shape_type=MSO_SHAPE.RECTANGLE,
    x="10%",
    y="65%",
    width="20%",
    height="15%",
    fill_color="blue",
    text="$1.24M\nTotal Revenue",
    font_color="white",
    font_size=14,
)

pres.add_shape(
    slide=combined_slide,
    shape_type=MSO_SHAPE.RECTANGLE,
    x="40%",
    y="65%",
    width="20%",
    height="15%",
    fill_color="green",
    text="+15.2%\nGrowth",
    font_color="white",
    font_size=14,
)

pres.add_shape(
    slide=combined_slide,
    shape_type=MSO_SHAPE.RECTANGLE,
    x="70%",
    y="65%",
    width="20%",
    height="15%",
    fill_color="orange",
    text="4.8%\nConversion",
    font_color="white",
    font_size=14,
)

# Add a closing slide
thank_you_slide = pres.add_section_slide(title="Thank You!")

# Save the presentation
output_path = output_dir / "object_api_example.pptx"
pres.save(output_path)
print(f"Presentation saved to {output_path}")
