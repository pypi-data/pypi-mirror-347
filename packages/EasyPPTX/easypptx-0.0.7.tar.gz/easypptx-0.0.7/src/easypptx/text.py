"""Text handling module for EasyPPTX."""

from typing import TYPE_CHECKING

from pptx.dml.color import RGBColor
from pptx.shapes.autoshape import Shape as PPTXShape
from pptx.text.text import TextFrame
from pptx.util import Pt

if TYPE_CHECKING:
    from easypptx.slide import Slide

# Type for position parameters - accepts either percentage or absolute values
PositionType = float | str
RGBColorTuple = tuple[int, int, int]


class Text:
    """Class for handling text operations in PowerPoint slides.

    This class provides methods for creating and formatting text elements.

    Examples:
        ```python
        # Create a text object
        text = Text(slide)

        # Add a title
        text.add_title("Presentation Title")

        # Add a paragraph
        text.add_paragraph("This is a paragraph", font_size=24)

        # Using the static method
        Text.add(slide, "Static text", position={"x": "10%", "y": "20%", "width": "80%", "height": "5%"})
        ```
    """

    def __init__(self, slide_obj: "Slide") -> None:
        """Initialize a Text object.

        Args:
            slide_obj: The Slide object to add text to
        """
        self.slide = slide_obj

    def add_title(
        self,
        text: str,
        font_size: int = 44,
        font_name: str = "Meiryo",
        color: str | RGBColorTuple | None = "black",
        align: str = "center",
        x: PositionType = "10%",
        y: PositionType = "5%",
        width: PositionType = "80%",
        height: PositionType = "15%",
    ) -> PPTXShape:
        """Add a title to the slide.

        Args:
            text: The title text
            font_size: Font size in points (default: 44)
            font_name: Font name (default: "Meiryo")
            color: Text color as string name from COLORS dict or RGB tuple (default: "black")
            align: Text alignment, one of "left", "center", "right" (default: "center")
            x: X position in inches or percentage (default: "10%")
            y: Y position in inches or percentage (default: "5%")
            width: Width in inches or percentage (default: "80%")
            height: Height in inches or percentage (default: "15%")

        Returns:
            The created shape object
        """
        shape = self.slide.add_text(
            text=text,
            x=x,
            y=y,
            width=width,
            height=height,
            font_size=font_size,
            font_bold=True,
            font_name=font_name,
            color=color,
            align=align,
            vertical="middle",
        )
        return shape

    def add_paragraph(
        self,
        text: str,
        x: PositionType = 1.0,
        y: PositionType = 2.0,
        width: PositionType = 8.0,
        height: PositionType = 1.0,
        font_size: int = 18,
        font_bold: bool = False,
        font_italic: bool = False,
        font_name: str = "Meiryo",
        align: str = "left",
        vertical: str = "top",
        color: str | RGBColorTuple | None = "black",
    ) -> PPTXShape:
        """Add a paragraph of text to the slide.

        Args:
            text: The paragraph text
            x: X position in inches or percentage (default: "10%")
            y: Y position in inches or percentage (default: "25%")
            width: Width in inches or percentage (default: "80%")
            height: Height in inches or percentage (default: "10%")
            font_size: Font size in points (default: 18)
            font_bold: Whether text should be bold (default: False)
            font_italic: Whether text should be italic (default: False)
            font_name: Font name (default: "Meiryo")
            align: Text alignment, one of "left", "center", "right" (default: "left")
            vertical: Vertical alignment, one of "top", "middle", "bottom" (default: "top")
            color: Text color as string name from COLORS dict or RGB tuple (default: "black")

        Returns:
            The created shape object
        """
        shape = self.slide.add_text(
            text=text,
            x=x,
            y=y,
            width=width,
            height=height,
            font_size=font_size,
            font_bold=font_bold,
            font_italic=font_italic,
            font_name=font_name,
            align=align,
            vertical=vertical,
            color=color,
        )

        return shape

    @staticmethod
    def add(
        slide,
        text: str,
        position: dict[str, PositionType],
        font_name: str = "Meiryo",
        font_size: int = 18,
        font_bold: bool = False,
        font_italic: bool = False,
        align: str = "left",
        vertical_align: str = "top",
        color: str | RGBColorTuple | None = "black",
    ) -> PPTXShape:
        """Static method to add text to a slide.

        Args:
            slide: Slide object to add text to
            text: Text content
            position: Dictionary with x, y, width, height as percentages or inches
            font_name: Font name (default: "Meiryo")
            font_size: Font size in points (default: 18)
            font_bold: Whether text should be bold (default: False)
            font_italic: Whether text should be italic (default: False)
            align: Text alignment, one of "left", "center", "right" (default: "left")
            vertical_align: Vertical alignment, one of "top", "middle", "bottom" (default: "top")
            color: Text color as string name from COLORS dict or RGB tuple (default: "black")

        Returns:
            The created shape object
        """
        # Create a text object for this slide
        text_obj = Text(slide)

        # Add the text with the specified parameters
        return text_obj.add_paragraph(
            text=text,
            x=position.get("x", "10%"),
            y=position.get("y", "10%"),
            width=position.get("width", "80%"),
            height=position.get("height", "10%"),
            font_name=font_name,
            font_size=font_size,
            font_bold=font_bold,
            font_italic=font_italic,
            align=align,
            vertical=vertical_align,
            color=color,
        )

    @staticmethod
    def format_text_frame(
        text_frame: TextFrame,
        font_size: int | None = None,
        font_bold: bool | None = None,
        font_italic: bool | None = None,
        font_name: str | None = None,
        color: str | RGBColorTuple | None = None,
        align: str | None = None,
        vertical: str | None = None,
    ) -> None:
        """Format an existing text frame.

        Args:
            text_frame: The text frame to format
            font_size: Font size in points (default: None)
            font_bold: Whether text should be bold (default: None)
            font_italic: Whether text should be italic (default: None)
            font_name: Font name (default: None)
            color: Text color as string name from COLORS dict or RGB tuple (default: None)
            align: Text alignment, one of "left", "center", "right" (default: None)
            vertical: Vertical alignment, one of "top", "middle", "bottom" (default: None)
        """
        # Set vertical alignment for the text frame
        if vertical:
            from easypptx.presentation import Presentation

            if vertical in Presentation.VERTICAL:
                text_frame.vertical_anchor = Presentation.VERTICAL[vertical]

        # Format all paragraphs
        for paragraph in text_frame.paragraphs:
            if font_size is not None:
                paragraph.font.size = Pt(font_size)
            if font_bold is not None:
                paragraph.font.bold = font_bold
            if font_italic is not None:
                paragraph.font.italic = font_italic
            if font_name is not None:
                paragraph.font.name = font_name

            # Set text alignment
            if align:
                from easypptx.presentation import Presentation

                if align in Presentation.ALIGN:
                    paragraph.alignment = Presentation.ALIGN[align]

            # Set text color
            if color:
                from easypptx.presentation import Presentation

                if isinstance(color, str) and color in Presentation.COLORS:
                    paragraph.font.color.rgb = Presentation.COLORS[color]
                elif isinstance(color, tuple) and len(color) == 3:
                    paragraph.font.color.rgb = RGBColor(*color)
