"""
Example showing how to use templates with TOML export/import in EasyPPTX.

This example demonstrates how to:
1. Create a custom template
2. Export a template to TOML or JSON
3. Import a template from TOML or JSON
4. Use the template to create slides

Templates can be shared, reused, and customized for consistent presentations.
"""

from pathlib import Path

from pptx.enum.shapes import MSO_SHAPE

from easypptx import Presentation
from easypptx.template import TemplateManager

# Create a folder for outputs if it doesn't exist
output_dir = Path("output")
output_dir.mkdir(exist_ok=True)

# Create a folder for templates if it doesn't exist
template_dir = output_dir / "templates"
template_dir.mkdir(exist_ok=True)

# Create a new presentation
pres = Presentation()

# Create a title slide
title_slide = pres.add_title_slide(
    title="Template Management with TOML", subtitle="Creating, Exporting, and Importing Templates"
)

# 1. Working with Built-in Templates
# =================================
section_slide = pres.add_section_slide(title="Built-in Templates")

# Create a slide to describe built-in templates
builtin_slide = pres.add_content_slide(title="Working with Built-in Templates")

# Add text explaining built-in templates
pres.add_text(
    slide=builtin_slide,
    text=(
        "EasyPPTX comes with several built-in templates:\n\n"
        "• title_slide: For presentation title slides\n"
        "• content_slide: For general content slides\n"
        "• section_slide: For section dividers\n"
        "• comparison_slide: For side-by-side comparisons\n"
        "• image_slide: For slides with a dominant image\n"
        "• table_slide: For data tables\n"
        "• chart_slide: For data visualizations\n"
        "• And many more..."
    ),
    x="10%",
    y="20%",
    width="80%",
    height="70%",
    font_size=18,
)

# Create an example of each built-in template
example_title_slide = pres.add_title_slide(title="Title Slide Template", subtitle="Used for introductory slides")

example_content_slide = pres.add_content_slide(title="Content Slide Template")
pres.add_text(
    slide=example_content_slide,
    text="Regular content slides have a title, optional bar, and content area.",
    x="10%",
    y="20%",
    width="80%",
    height="10%",
    font_size=18,
)

example_section_slide = pres.add_section_slide(title="Section Slide Template")

# 2. Create a Custom Template
# ==========================
section_slide = pres.add_section_slide(title="Custom Templates")

# Create a slide to explain custom templates
custom_template_slide = pres.add_content_slide(title="Creating Custom Templates")
pres.add_text(
    slide=custom_template_slide,
    text=(
        "You can create custom templates by:\n\n"
        "1. Building a template dictionary manually\n"
        "2. Modifying an existing template\n"
        "3. Creating a template manager and registering templates\n\n"
        "Custom templates can include:\n"
        "• Background colors\n"
        "• Title and subtitle styling\n"
        "• Positioning information\n"
        "• Default styling for elements (text, shapes, etc.)\n"
    ),
    x="10%",
    y="20%",
    width="80%",
    height="70%",
    font_size=18,
)

# Initialize a template manager
tm = TemplateManager(template_dir=str(template_dir))

# Create a custom template for a product slide
product_template = {
    "bg_color": "white",
    "title": {
        "text": "Product Template",
        "position": {"x": "5%", "y": "5%", "width": "90%", "height": "10%"},
        "font": {"name": "Meiryo", "size": 32, "bold": True},
        "align": "left",
        "vertical": "middle",
        "color": "blue",
    },
    "image_area": {"position": {"x": "5%", "y": "20%", "width": "40%", "height": "70%"}},
    "detail_area": {"position": {"x": "50%", "y": "20%", "width": "45%", "height": "70%"}},
    "image_style": {
        "border": True,
        "border_color": "blue",
        "border_width": 2,
        "shadow": True,
        "maintain_aspect_ratio": True,
    },
    "key_feature_title": {
        "position": {"x": "50%", "y": "20%", "width": "45%", "height": "10%"},
        "font": {"name": "Meiryo", "size": 22, "bold": True},
        "align": "left",
        "color": "blue",
    },
    "bullet_points": {
        "position": {"x": "50%", "y": "32%", "width": "45%", "height": "30%"},
        "font": {"name": "Meiryo", "size": 16, "bold": False},
        "align": "left",
    },
    "price_box": {
        "position": {"x": "50%", "y": "65%", "width": "45%", "height": "15%"},
        "shape_type": MSO_SHAPE.ROUNDED_RECTANGLE,
        "fill_color": "blue",
        "text_color": "white",
        "font": {"name": "Meiryo", "size": 24, "bold": True},
        "align": "center",
    },
}

# Create a custom template for a comparison slide
comparison_template = {
    "bg_color": (245, 245, 245),  # Light gray background
    "title": {
        "text": "Comparison Template",
        "position": {"x": "5%", "y": "5%", "width": "90%", "height": "10%"},
        "font": {"name": "Meiryo", "size": 32, "bold": True},
        "align": "center",
        "vertical": "middle",
        "color": "black",
    },
    "vs_text": {
        "position": {"x": "45%", "y": "40%", "width": "10%", "height": "20%"},
        "font": {"name": "Meiryo", "size": 24, "bold": True},
        "align": "center",
        "vertical": "middle",
        "color": "gray",
    },
    "left_title": {
        "position": {"x": "5%", "y": "20%", "width": "40%", "height": "10%"},
        "font": {"name": "Meiryo", "size": 24, "bold": True},
        "align": "center",
        "vertical": "middle",
        "color": "blue",
    },
    "right_title": {
        "position": {"x": "55%", "y": "20%", "width": "40%", "height": "10%"},
        "font": {"name": "Meiryo", "size": 24, "bold": True},
        "align": "center",
        "vertical": "middle",
        "color": "green",
    },
    "left_content": {"position": {"x": "5%", "y": "32%", "width": "40%", "height": "60%"}},
    "right_content": {"position": {"x": "55%", "y": "32%", "width": "40%", "height": "60%"}},
}

# Register the custom templates
tm.register("product_slide", product_template)
tm.register("comparison_slide", comparison_template)

# Create slides to demonstrate the custom templates
custom_slide1 = pres.add_slide()
pres.add_text(
    slide=custom_slide1,
    text="Product Template (Registered)",
    x="5%",
    y="5%",
    width="90%",
    height="10%",
    font_size=32,
    font_bold=True,
    align="left",
    color="blue",
)

# Use elements from the product template
pres.add_shape(
    slide=custom_slide1,
    shape_type=MSO_SHAPE.RECTANGLE,
    x="5%",
    y="20%",
    width="40%",
    height="70%",
    fill_color="lightgray",
    text="Product Image\nArea",
    font_size=18,
    text_align="center",
)

pres.add_text(
    slide=custom_slide1,
    text="Key Features",
    x="50%",
    y="20%",
    width="45%",
    height="10%",
    font_size=22,
    font_bold=True,
    align="left",
    color="blue",
)

pres.add_text(
    slide=custom_slide1,
    text="• Feature 1\n• Feature 2\n• Feature 3\n• Feature 4",
    x="50%",
    y="32%",
    width="45%",
    height="30%",
    font_size=16,
    align="left",
)

pres.add_shape(
    slide=custom_slide1,
    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
    x="50%",
    y="65%",
    width="45%",
    height="15%",
    fill_color="blue",
    text="$99.99",
    font_size=24,
    font_bold=True,
    font_color="white",
    text_align="center",
)

# Create a slide for the comparison template
custom_slide2 = pres.add_slide()
pres.add_text(
    slide=custom_slide2,
    text="Comparison Template (Registered)",
    x="5%",
    y="5%",
    width="90%",
    height="10%",
    font_size=32,
    font_bold=True,
    align="center",
)

pres.add_text(
    slide=custom_slide2,
    text="VS",
    x="45%",
    y="40%",
    width="10%",
    height="20%",
    font_size=24,
    font_bold=True,
    align="center",
    color="gray",
)

pres.add_text(
    slide=custom_slide2,
    text="Option A",
    x="5%",
    y="20%",
    width="40%",
    height="10%",
    font_size=24,
    font_bold=True,
    align="center",
    color="blue",
)

pres.add_text(
    slide=custom_slide2,
    text="Option B",
    x="55%",
    y="20%",
    width="40%",
    height="10%",
    font_size=24,
    font_bold=True,
    align="center",
    color="green",
)

pres.add_shape(
    slide=custom_slide2,
    shape_type=MSO_SHAPE.RECTANGLE,
    x="5%",
    y="32%",
    width="40%",
    height="60%",
    fill_color="lightgray",
    text="Left Content Area",
    font_size=18,
    text_align="center",
)

pres.add_shape(
    slide=custom_slide2,
    shape_type=MSO_SHAPE.RECTANGLE,
    x="55%",
    y="32%",
    width="40%",
    height="60%",
    fill_color="lightgray",
    text="Right Content Area",
    font_size=18,
    text_align="center",
)

# 3. Export Templates to TOML and JSON
# ===================================
section_slide = pres.add_section_slide(title="Template Export & Import")

# Create a slide to explain template export
export_slide = pres.add_content_slide(title="Exporting Templates")
pres.add_text(
    slide=export_slide,
    text=(
        "Templates can be exported to TOML or JSON format:\n\n"
        "• TOML is the default format (more readable)\n"
        "• JSON is also supported\n\n"
        "Built-in templates and custom templates can be exported:\n\n"
        'template_manager.save("template_name", file_path, format="toml")\n'
        'template_manager.save("custom_template", file_path, format="json")\n\n'
        "If file_path is None, templates are saved to the template directory."
    ),
    x="10%",
    y="20%",
    width="80%",
    height="70%",
    font_size=18,
)

# Save the custom templates to TOML and JSON
product_toml_path = tm.save("product_slide", str(template_dir / "product_slide.toml"), format="toml")
comparison_json_path = tm.save("comparison_slide", str(template_dir / "comparison_slide.json"), format="json")

# Also export a built-in template
title_toml_path = tm.save("title_slide", str(template_dir / "title_slide.toml"), format="toml")

# Create a slide to explain template import
import_slide = pres.add_content_slide(title="Importing Templates")
pres.add_text(
    slide=import_slide,
    text=(
        "Templates can be imported from TOML or JSON files:\n\n"
        "• File format is detected by extension (.toml or .json)\n"
        "• Templates are registered with the TemplateManager\n\n"
        "Import methods are simple:\n\n"
        'template_manager.load("path/to/template.toml")\n'
        'template_manager.load("path/to/template.json", "custom_name")\n\n'
        "If no template_name is provided, the filename is used."
    ),
    x="10%",
    y="20%",
    width="80%",
    height="70%",
    font_size=18,
)

# Let's clear all registered templates to demonstrate loading from files
tm.registered_templates = {}

# Load templates from files
product_name = tm.load(product_toml_path)
comparison_name = tm.load(comparison_json_path)
title_name = tm.load(title_toml_path, "my_title_slide")

# Check that templates were loaded successfully
template_names = tm.list_templates()

# Create a slide showing the loaded templates
loaded_slide = pres.add_content_slide(title="Loaded Templates")
pres.add_text(
    slide=loaded_slide,
    text=(
        f"Successfully loaded templates:\n\n"
        f"• {product_name} (from TOML)\n"
        f"• {comparison_name} (from JSON)\n"
        f"• {title_name} (from TOML with custom name)\n\n"
        f"Available templates: {', '.join(template_names)}"
    ),
    x="10%",
    y="20%",
    width="80%",
    height="30%",
    font_size=18,
)

# Create a slide using a template from a file
# Use the add_slide_from_template method which works with TemplateManager
product_slide = pres.add_slide_from_template("product_slide")

# Manually add content based on template positions
pres.add_text(
    slide=product_slide,
    text="Product: Premium Headphones",
    x="5%",
    y="5%",
    width="90%",
    height="10%",
    font_size=32,
    font_bold=True,
    align="left",
    color="blue",
)

pres.add_shape(
    slide=product_slide,
    shape_type=MSO_SHAPE.RECTANGLE,
    x="5%",
    y="20%",
    width="40%",
    height="70%",
    fill_color="lightgray",
    text="Product Image",
    font_size=18,
    text_align="center",
)

pres.add_text(
    slide=product_slide,
    text="Key Features",
    x="50%",
    y="20%",
    width="45%",
    height="10%",
    font_size=22,
    font_bold=True,
    align="left",
    color="blue",
)

pres.add_text(
    slide=product_slide,
    text="• Noise cancellation\n• 30-hour battery life\n• High-quality sound\n• Comfortable fit",
    x="50%",
    y="32%",
    width="45%",
    height="30%",
    font_size=16,
    align="left",
)

pres.add_shape(
    slide=product_slide,
    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
    x="50%",
    y="65%",
    width="45%",
    height="15%",
    fill_color="blue",
    text="$299.99",
    font_size=24,
    font_bold=True,
    font_color="white",
    text_align="center",
)

# Create a Thank You slide
thank_you_slide = pres.add_section_slide(title="Thank You!")

# Save the presentation
pres.save(output_dir / "template_toml_example.pptx")
print(f"Presentation saved to {output_dir / 'template_toml_example.pptx'}")
print(f"Templates saved to {template_dir}")
print(f"Available templates: {', '.join(tm.list_templates())}")
