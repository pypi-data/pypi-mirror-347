import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from numerize import numerize 
from scipy.stats import pearsonr
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt
import os
import math
from .functions.chart_functions.line import *
from .functions.chart_functions.bar import *
from .functions.chart_functions.facetLine import *
from .functions.chart_functions.facetBar import *
from .functions.chart_functions.filledLine import *
from .functions.chart_functions.facetFilledLine import *
from .functions.chart_functions.lineWithHighlight import *
from .functions.dataManipulation import *


#Define functions for table/cell formatting
def SubElement(parent, tagname, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element

def _set_cell_border(cell, border_color="000000", border_width='12700'):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for lines in ['a:lnL','a:lnR','a:lnT','a:lnB']:
        ln = SubElement(tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = SubElement(ln, 'a:NoFill')
        srgbClr = SubElement(solidFill, 'a:srgbClr', val=border_color)
        prstDash = SubElement(ln, 'a:prstDash', val='solid')
        round_ = SubElement(ln, 'a:round')
        headEnd = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')

    return cell

#function to set the powerpoint template being used
def setTemplate(fileName):
    from pptx import Presentation
    
    #validate input
    if type(fileName) != str:
        raise Exception("You must input your filename as a string")

    #Load in template for presentation
    try:
        global prs 
        prs = Presentation(fileName)
    except:
        raise Exception("File not found")
        
#function for setting a color palette for charts
def setColors(colors):
    global colorPalette
    colorPalette = colors
    
#master function for creating slides
def createSlides(slides):
    
    #loop through each item in the array
    for z in range(len(slides)):

        #Grab our slide definition
        slideDefinition = slides[z]

        #####################
        ### Create Slide and place easy info
        #####################

        #create slide
        layout = prs.slide_layouts[slideDefinition['item-index']['slide']]
        slide = prs.slides.add_slide(layout)

        #set title and subtitle
        if 'title' in slideDefinition:
            slide.placeholders[slideDefinition['item-index']['title']].text = slideDefinition['title']

        #insert placeholder if desired, otherwise delete
        if "description" in slideDefinition:
            slide.placeholders[slideDefinition['item-index']['description']].text = slideDefinition['description']                                

        #insert subtitle if present in dictionary
        if "subtitle" in slideDefinition:
            slide.placeholders[slideDefinition['item-index']['subtitle']].text = slideDefinition['subtitle']

        
        
        ### Now place custom text items if present
        if 'custom-text-items' in slideDefinition:
            for textItem in slideDefinition['custom-text-items']:
                slide.placeholders[textItem['item-index']].text = textItem['value']

        
        #####################
        ### Create each chart and place it into the slide
        #####################

        #loop through
        for i in range(len(slideDefinition['charts'])):

            #grab the chart definition
            chartDefinition = slideDefinition['charts'][i]

            #set color palette. If pre-set, define it. If not, use default
            try:
                colorPalette
            except NameError:
                mainColors = px.colors.qualitative.Plotly
            else:
                mainColors = colorPalette
                                
            #get data defined
            temp = chartDefinition['data']

            #filter and group data
            dataOb = prepareData(temp, chartDefinition)
            temp = dataOb['dataframe']
            groupList = dataOb['groupList']

            ### Create the chart figure
            #line chart
            if chartDefinition['type'] == 'line':
    
                #Create line chart
                fig = createLineChart(temp, chartDefinition, mainColors)
    
            #bar chart
            if chartDefinition['type'] == 'bar':
    
                #Create bar chart
                fig = createBarChart(temp, chartDefinition, mainColors, groupList)                                    
    
            #facet line
            if chartDefinition['type'] == 'facetLine':
    
                #Create line chart
                fig = createFacetLineChart(temp, chartDefinition, mainColors)                        
                
            #Facet Bar Chart
            if chartDefinition['type'] == 'facetBar':
    
                #Create line chart
                fig = createFacetBarChart(temp, chartDefinition, mainColors)                     
    
            #Filled line chart
            if chartDefinition['type'] == 'filledLine':
                
                #Create Chart
                fig = createFilledLineChart(temp, chartDefinition, mainColors)
    
            #Facet Fill Line
            if chartDefinition['type'] == 'facetFilledLine':
                
                #Create Chart
                fig = createFacetFilledLineChart(temp, chartDefinition, mainColors)
            
            #Line With Highlight
            if chartDefinition['type'] == 'lineWithHighlight':

                #Create Chart
                fig = createLineWithHighlightChart(temp, chartDefinition, mainColors)

            #save figure
            if chartDefinition['type'] != 'table':
    
                #check if the folder for charts exists. If not, create it
                if not os.path.exists('charts'):
                    os.makedirs('charts')
                
                #setup params
                filename = f'charts/chart{z}_{i}.png'

                #Grab the placeholder and get the dimensions
                image_width = slide.placeholders[chartDefinition['item-index']].width
                image_height = slide.placeholders[chartDefinition['item-index']].height
                image_ratio = image_width / image_height

                #figure out the width and height
                fin_width = 500
                fin_height = 500
                if image_width > image_height:
                    fin_width = 500 * image_ratio
                else:
                    fin_height = 500 * image_ratio
                
                #save out the files
                fig.write_image(filename, scale=2, width=fin_width, height=fin_height)


            ### Now move onto placing the chart into the slide
            if chartDefinition['type'] != 'table':
                #insert image
                picture = slide.placeholders[chartDefinition['item-index']].insert_picture(filename)

            else:
                #we're going to insert a table                
                shape = slide.placeholders[chartDefinition['item-index']].insert_table(rows=len(temp)+1, cols=len(temp.columns))
                table = shape.table
                
                #iterate through every row and column and place the value that is present in the df
                #for loop for the rows
                for i in range(len(temp) + 1):
                    #for each row, get the value of the column
                    for i2 in range(len(temp.columns)):
                        cell = table.cell(i,i2)
                        #if we're dealing with the header
                        if i == 0:
                            cell.text = temp.columns[i2]
                        else:
                            text = temp.iloc[i-1, i2]
                            textFormat = chartDefinition['column_formats'][i2]
    
                            #catch Nan values for numeric based values
                            if 'float' in str(type(text)) or 'int' in str(type(text)):
                                if math.isnan(text):
                                    cell.text = ''
                                    continue
                                    
                            #catch Nan values for string based values
                            if 'str' in str(type(float)) or 'NoneType' in str(type(text)):
                                if text is None:
                                    cell.text = ''
                                    continue
                            
                            if textFormat == 'number':
                                cell.text = str(int(text))
                            elif textFormat == 'money':
                                cell.text = "$" + str(int(text))
                            elif textFormat == 'percent':
                                cell.text = str(int(text * 100)) + "%"
                            elif textFormat == 'twoDigitNum':
                                cell.text = str(round(text, 2))
                            elif textFormat == 'date':
                                cell.text = text.strftime('%m/%d/%Y')
                            else:
                                cell.text = str(text)
    
                            
                #central formatting for every cell
                for i in range(len(temp) + 1):
                    for i2 in range(len(temp.columns)):
                        #Remove the border for each cell
                        cell = table.cell(i,i2)
                        cell = _set_cell_border(cell, '000000', '0')
                        # cell = _set_cell_border(cell, '000000', '0', 'all')
                        
                        #format color
                        cell = table.cell(i,i2)
                        paragraph = cell.text_frame.paragraphs[0]
                        colorString = chartDefinition['text_color']
                        color = RGBColor.from_string(colorString.replace('#',''))
                        paragraph.font.color.rgb = color
    
                #If we need to change the header fill
                if 'header_fill_color' in chartDefinition:
                    for i in range(len(temp.columns)):
                        cell = table.cell(0,i)
                        cell.fill.solid()
                        colorString = chartDefinition['header_fill_color']
                        color = RGBColor.from_string(colorString.replace('#',''))
                        cell.fill.fore_color.rgb = color
                        
                #If we need to change the header text color
                if 'header_text_color' in chartDefinition:
                    for i in range(len(temp.columns)):
                        cell = table.cell(0,i)
                        colorString = chartDefinition['header_text_color']
                        color = RGBColor.from_string(colorString.replace('#',''))
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.font.color.rgb = color
    
                #If we need to change the font size of the header
                if 'header_font_size' in chartDefinition:
                    for i in range(len(temp.columns)):
                        cell = table.cell(0,i)
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.font.size = Pt(int(chartDefinition['header_font_size']))
    
                #If we need to change the fill color for each cell
                if 'fill_color' in chartDefinition:
                    #get the data for the fill coloring
                    fillData = chartDefinition['fill_color']
                    fillData = fillData.replace("#", '', regex=True)
                    
                    #loop through each cell
                    for i in range(len(temp) + 1):
                        for i2 in range(len(temp.columns)):
                            #skip the header
                            if i != 0:
                                cell = table.cell(i,i2)
                                cell.fill.solid()
                                color = RGBColor.from_string(fillData.iloc[i-1, i2])
                                cell.fill.fore_color.rgb = color
    
                #if we need to change the font size for the main cells
                if 'text_font_size' in chartDefinition:
                    #loop through each cell
                    for i in range(len(temp) + 1):
                        for i2 in range(len(temp.columns)):
                            #skip the header
                            if i != 0:
                                cell = table.cell(i, i2)
                                paragraph = cell.text_frame.paragraphs[0]
                                paragraph.font.size = Pt(int(chartDefinition['text_font_size']))                        
                
                ### Now center the table in the middle of the slide
                #get base variables
                slideHeight = 5143500
                heightOffset = chartDefinition['top_offset'] if 'top_offset' in chartDefinition else 0
                titleHeight = slide.placeholders[slideDefinition['item-index']['title']].height if 'title' in slideDefinition['item-index'] else 0
                tableHeight = slide.placeholders[chartDefinition['item-index']].height
    
                #calculate where the table needs to start
                middleOfSlide = int(slideHeight / 2) + int(titleHeight / 2)
                halfTableHeight = int(tableHeight / 2)
                idealTableStart = int(middleOfSlide - halfTableHeight + heightOffset)
    
                #set the top of the table
                slide.placeholders[chartDefinition['item-index']].top = idealTableStart


    #finally save out file
    prs.save("output.pptx")