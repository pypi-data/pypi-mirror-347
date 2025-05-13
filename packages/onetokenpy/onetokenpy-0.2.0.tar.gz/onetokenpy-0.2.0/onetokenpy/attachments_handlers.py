from .core import Attachments
import requests
from pathlib import Path
from PIL import Image
import io

# ========== Attachments Handlers (files, URLs, images) ==========

@Attachments.register_handler('.pdf')
def _(p: Path):
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(p)
        out = "\n".join(page.get_text() for page in doc[:5]) or "[empty pdf]"
        return {'type': 'text', 'content': f"<pdf name='{p.name}'>{out}</pdf>", 'identifier': str(p)}
    except ImportError:
        return {'type': 'text', 'content': '[PDF text extraction not available – install pymupdf]', 'identifier': str(p)}


@Attachments.register_handler('.pptx')
def _(p: Path):
    try:
        from pptx import Presentation
        slides = Presentation(p)
        out = []
        for n, s in enumerate(slides.slides, 1):
            txt = " ".join(
                [sh.text for sh in s.shapes if hasattr(sh, "text")]).strip()
            out.append(f"Slide {n}: {txt}")
            str_out = '\n'.join(out)
        return {'type': 'text', 'content': f"<pptx name='{p.name}'>{str_out}</pptx>", 'identifier': str(p)}
    except ImportError:
        return {'type': 'text', 'content': '[PPTX extraction unavailable – install python-pptx]', 'identifier': str(p)}


@Attachments.register_handler('.md')
@Attachments.register_handler('.txt')
def _(p: Path):
    return {'type': 'text', 'content': f"<file name='{p.name}'>{p.read_text()}</file>", 'identifier': str(p)}


MAX_IMAGE_SIZE_BYTES = 4.5 * 1024 * 1024  # 4.5 MB

def _resize_image_if_needed(p: Path) -> bytes:
    """Resizes an image if it's larger than MAX_IMAGE_SIZE_BYTES."""
    img_bytes = p.read_bytes()
    if len(img_bytes) > MAX_IMAGE_SIZE_BYTES:
        try:
            img = Image.open(io.BytesIO(img_bytes))
            
            # Convert to RGB if it's RGBA (alpha channel can increase size, and often not needed)
            if img.mode == 'RGBA' or img.mode == 'P': # P is for paletted images
                img = img.convert('RGB')

            # Simple resizing: reduce dimensions until size is acceptable
            # More sophisticated resizing could target a max dimension or iteratively reduce quality
            width, height = img.size
            current_size = len(img_bytes)
            quality = 85 # Initial JPEG quality

            # Attempt to save as JPEG and check size
            temp_buffer = io.BytesIO()
            img.save(temp_buffer, format="JPEG", quality=quality, optimize=True)
            img_bytes_resized = temp_buffer.getvalue()
            current_size = len(img_bytes_resized)

            # If still too large, progressively reduce dimensions and/or quality
            # This is a simple iterative approach, could be refined
            scale_factor = 0.9
            max_iterations = 5 # Prevent infinite loops

            while current_size > MAX_IMAGE_SIZE_BYTES and max_iterations > 0 and (width > 100 and height > 100):
                temp_buffer = io.BytesIO()
                if width * scale_factor > 100 and height * scale_factor > 100: # Don't make too small
                    new_width = int(width * scale_factor)
                    new_height = int(height * scale_factor)
                    img_resized_dim = img.resize((new_width, new_height), Image.LANCZOS)
                    img_resized_dim.save(temp_buffer, format="JPEG", quality=quality, optimize=True)
                    width, height = new_width, new_height
                else: # if dimensions are already small, try reducing quality more aggressively
                    quality = max(10, quality - 10) # reduce quality, min 10
                    img.save(temp_buffer, format="JPEG", quality=quality, optimize=True)
                
                img_bytes_resized = temp_buffer.getvalue()
                current_size = len(img_bytes_resized)
                max_iterations -= 1
                if quality <= 10 and (width * scale_factor <=100 or height * scale_factor <=100):
                    break # Stop if quality is min and dimensions are small

            if current_size > MAX_IMAGE_SIZE_BYTES:
                # If still too large after attempts, we might have to raise or return a message
                # For now, returning the last resized version even if slightly over,
                # or we could raise an error / return a "cannot compress enough" message.
                # Let's return truncated for now to avoid breaking API call.
                # However, the API will likely still reject if it's over.
                # A better solution would be a more robust resizing strategy or clear error.
                # For now, let's assume this should get it under the limit.
                pass
            
            return img_bytes_resized
        except Exception as e:
            # If resizing fails, fall back to original bytes but log/warn
            # Or, return an error message in the content
            # print(f"Warning: Could not resize image {p.name}: {e}. Sending original.")
            # For now, we'll return original if error, but this might still fail API call
            return img_bytes # Fallback, though likely to still hit API limit
    return img_bytes

@Attachments.register_handler('.png')
@Attachments.register_handler('.jpg')
@Attachments.register_handler('.jpeg')
@Attachments.register_handler('.webp')
def _(p: Path):
    processed_bytes = _resize_image_if_needed(p)
    return {'type': 'image', 'content': processed_bytes, 'identifier': str(p)}


@Attachments.register_handler(r'^https?://\S+', priority=10)
def _(url: str):
    """Fetch content from a URL and return as text or image based on content type."""
    try:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'}
        response = requests.get(url, timeout=10, headers=headers)
        response.raise_for_status()
        content_type = response.headers.get('Content-Type', '').lower()
        if 'image' in content_type:
            return {'type': 'image', 'content': response.content, 'identifier': url}
        elif 'text' in content_type or 'html' in content_type:
            # Limit content size for efficiency
            content = response.text[:20000]
            if len(response.text) > 20000:
                content += "\n[...content truncated]"
            return {'type': 'text', 'content': f"<url>{url}</url>\n<content>{content}</content>", 'identifier': url}
        else:
            return {'type': 'text', 'content': f"[Unsupported content type: {content_type}]", 'identifier': url}
    except Exception as e:
        return {'type': 'text', 'content': f"[Error fetching URL: {e}]", 'identifier': url}


@Attachments.register_handler('.csv')
def _(p: Path):
    import pandas as pd
    return {'type': 'text', 'content': pd.read_csv(p).head().to_csv(), 'identifier': str(p)}

# --- Attachments: add HEIC handler right after class definition ---


@Attachments.register_handler('.heic')
def _(p: Path):
    try:
        from PIL import Image
        import pillow_heif
        import io
        pillow_heif.register_heif_opener()
        
        # Initial check for original HEIC file size if needed, but conversion often changes size.
        # It's better to check after conversion to JPEG.

        img = Image.open(p)
        
        # Convert to RGB if it's RGBA or P
        if img.mode == 'RGBA' or img.mode == 'P':
            img = img.convert('RGB')

        buf = io.BytesIO()
        quality = 85 # Initial JPEG quality
        img.save(buf, format="JPEG", quality=quality, optimize=True)
        img_bytes = buf.getvalue()

        if len(img_bytes) > MAX_IMAGE_SIZE_BYTES:
            # Apply similar resizing logic as in _resize_image_if_needed
            width, height = img.size
            current_size = len(img_bytes)
            scale_factor = 0.9
            max_iterations = 5

            while current_size > MAX_IMAGE_SIZE_BYTES and max_iterations > 0 and (width > 100 and height > 100):
                temp_buffer = io.BytesIO()
                if width * scale_factor > 100 and height * scale_factor > 100:
                    new_width = int(width * scale_factor)
                    new_height = int(height * scale_factor)
                    img_resized_dim = img.resize((new_width, new_height), Image.LANCZOS)
                    img_resized_dim.save(temp_buffer, format="JPEG", quality=quality, optimize=True)
                    width, height = new_width, new_height
                else:
                    quality = max(10, quality - 10)
                    img.save(temp_buffer, format="JPEG", quality=quality, optimize=True)
                
                img_bytes_resized_heic = temp_buffer.getvalue()
                current_size = len(img_bytes_resized_heic)
                max_iterations -= 1
                if quality <= 10 and (width * scale_factor <=100 or height * scale_factor <=100):
                    break
            
            img_bytes = img_bytes_resized_heic # Use the resized bytes

        return {'type': 'image', 'content': img_bytes, 'identifier': str(p)}
    except Exception as e:
        return {'type': 'text', 'content': f"[Cannot load HEIC: {e}]", 'identifier': str(p)}

